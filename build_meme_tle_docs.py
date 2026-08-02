#!/usr/bin/env python3
"""
build_meme_tle_docs.py
======================
由 MEME vs TLE 研究結果產生「期中報告 Word（.docx）」與「簡報 PPT（.pptx）」，
乾淨學術樣式。內容與 docs/meme_tle_report/MEME_TLE_comparison_interim.md 一致，
圖表取自 docs/meme_tle_report/figs/。可在圖表更新後重跑以重建兩份文件。

用法： python build_meme_tle_docs.py
輸出：
  docs/meme_tle_report/MEME_TLE_comparison_interim.docx
  docs/meme_tle_report/MEME_TLE_comparison_interim.pptx
"""
from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor, Inches

from pptx import Presentation
from pptx.dml.color import RGBColor as PColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PInches, Pt as PPt, Emu

# ── 樣式常數 ──────────────────────────────────────────────────────────────────
BLUE   = "0072B2"    # 主色（MEME / 標題）
ORANGE = "E69F00"    # TLE
RED    = "D55E00"    # 機動 / 警示
GREEN  = "009E73"
INK    = "1A1A1A"
MUTED  = "6B6B6B"
LIGHT  = "F2F6F9"
CJK    = "Microsoft JhengHei"
LATIN  = "Calibri"

ROOT = Path("f:/GitHub/Sat_TraingDataExtension")
RPT  = ROOT / "docs" / "meme_tle_report"
FIG  = RPT / "figs"

META = [
    ("計畫案號", "TASA-S-1150268"),
    ("資料基準", "2026-05-02 ～ 2026-07-08（Starlink MEME）"),
    ("分析樣本", "50 顆代表性子集（35 顆已知機動 + 15 顆靜止衛星）"),
    ("對應程式", "study1/2/3_*.py（重用 compare_tle_vs_ephemeris.py）"),
]


# ════════════════════════════════════════════════════════════════════════════
#  WORD
# ════════════════════════════════════════════════════════════════════════════

def _run(p, text, size=11, bold=False, color=INK, italic=False):
    r = p.add_run(text)
    r.font.size = Pt(size); r.bold = bold; r.italic = italic
    r.font.color.rgb = RGBColor.from_string(color)
    r.font.name = LATIN
    rpr = r._element.get_or_add_rPr()
    rf = rpr.get_or_add_rFonts()
    rf.set(qn("w:eastAsia"), CJK)
    return r


def _heading(doc, text, level=1):
    p = doc.add_paragraph()
    p.space_before = Pt(10)
    sizes = {1: 15, 2: 12.5}
    _run(p, text, size=sizes.get(level, 12), bold=True,
         color=BLUE if level == 1 else INK)
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.space_before = Pt(12 if level == 1 else 8)
    return p


def _bullet(doc, text, sub=False):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.left_indent = Inches(0.5 if sub else 0.3)
    _run(p, text)
    p.paragraph_format.space_after = Pt(2)
    return p


def _para(doc, text, size=11, color=INK, italic=False):
    p = doc.add_paragraph()
    _run(p, text, size=size, color=color, italic=italic)
    return p


def _table(doc, headers, rows, widths=None):
    t = doc.add_table(rows=1, cols=len(headers))
    t.style = "Light Grid Accent 1"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        c = t.rows[0].cells[i]
        c.paragraphs[0].clear()
        _run(c.paragraphs[0], h, bold=True, size=10, color="FFFFFF")
        _shade(c, BLUE)
    for row in rows:
        cells = t.add_row().cells
        for i, val in enumerate(row):
            cells[i].paragraphs[0].clear()
            _run(cells[i].paragraphs[0], str(val), size=9.5)
    if widths:
        for i, w in enumerate(widths):
            for r in t.rows:
                r.cells[i].width = Inches(w)
    return t


def _shade(cell, hexcolor):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = tcPr.makeelement(qn("w:shd"), {qn("w:val"): "clear", qn("w:fill"): hexcolor})
    tcPr.append(shd)


def _figure(doc, name, caption, width=6.1):
    path = FIG / name
    if not path.exists():
        _para(doc, f"[缺圖 {name}]", color=RED); return
    doc.add_picture(str(path), width=Inches(width))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _run(cap, caption, size=9, italic=True, color=MUTED)
    cap.paragraph_format.space_after = Pt(8)


def build_docx():
    doc = Document()
    # 預設樣式字型
    normal = doc.styles["Normal"]
    normal.font.name = LATIN; normal.font.size = Pt(11)
    normal.element.get_or_add_rPr().get_or_add_rFonts().set(qn("w:eastAsia"), CJK)
    for s in doc.sections:
        s.top_margin = s.bottom_margin = Inches(0.8)
        s.left_margin = s.right_margin = Inches(0.9)

    # 標題
    tp = doc.add_paragraph(); tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _run(tp, "MEME 精密星曆 vs. 公開 TLE", size=20, bold=True, color=BLUE)
    sp = doc.add_paragraph(); sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _run(sp, "軌道預測誤差與外推特性研究　—　期中報告", size=13, color=INK)
    doc.add_paragraph()
    mt = doc.add_table(rows=0, cols=2); mt.alignment = WD_TABLE_ALIGNMENT.CENTER
    for k, v in META:
        cells = mt.add_row().cells
        _run(cells[0].paragraphs[0], k, bold=True, size=10, color=BLUE)
        _run(cells[1].paragraphs[0], v, size=10)
        cells[0].width = Inches(1.3); cells[1].width = Inches(5.0)
    doc.add_paragraph()

    # 一、背景
    _heading(doc, "一、背景與基本概念")
    _para(doc, "本研究以 SpaceX 精密星曆（MEME）作為近似地面真相，反向量化公開 TLE 的誤差與外推特性，"
               "回應契約「專案背景」對 TLE 精度限制之關切，並為 Layer 1 差分偵測的誤報來源提供實測依據。")
    _table(doc, ["資料源", "產生方", "精度", "特性"], [
        ["公開 TLE", "美國太空監視網（18 SDS）", "位置誤差公里級", "免費、全目錄、以 SGP4 傳播"],
        ["MEME 精密星曆", "SpaceX（每顆自報）", "次公里～百公尺級", "72h 涵蓋、1 分鐘間隔、每 ~8h 重發"],
    ], widths=[1.2, 1.9, 1.3, 2.0])
    _para(doc, "")
    _para(doc, "關鍵洞見：每個 MEME 檔的第一筆（外推齡 0）外推量最小、最接近定軌真值；相鄰檔每 ~8h 發布、"
               "重疊約 88%。此結構衍生兩類互補比較——MEME vs TLE（量 TLE 誤差）與 MEME vs MEME"
               "（量 MEME 自身外推誤差，無需外部傳播器）。", )
    _figure(doc, "fig1_concept.png",
            "圖 1　MEME 星曆檔結構：72h 涵蓋、1 分鐘間隔、每 ~8h 重發、重疊 ~88%；實心點為各檔第一筆近真值。")

    # 二、方法
    _heading(doc, "二、使用方法")
    _table(doc, ["程式", "研究", "方法摘要"], [
        ["study1_tle_error_distribution.py", "MEME vs TLE 誤差分布",
         "每檔第一筆為近真值；每快照選最新先行 TLE 以 SGP4 傳播，算 TLE−MEME 殘差與 TLE 齡"],
        ["study2_meme_self_prediction.py", "MEME vs MEME（0–72h）",
         "同一時刻：晚檔第一筆為真值、早檔外推為預測，殘差依外推時程分箱"],
        ["study3_tle_frozen_and_gap.py", "TLE 凍結曲線 + 斷點",
         "(A) 凍結單筆 TLE 外推 1–7 天；(B) 自動偵測斷點，外推越過斷點並過濾期間機動衛星"],
    ], widths=[2.0, 1.5, 3.0])
    _para(doc, "")
    _para(doc, "誤差一律以 RTN 座標分解（R 徑向 / T 沿軌 / N 面外），以 3D 位置誤差範數為主指標。", size=10, color=MUTED)

    # 三、MEME vs TLE
    _heading(doc, "三、結果 A：MEME vs TLE — 公開 TLE 的誤差分布")
    for b in [
        "位置誤差中位數 P50 = 2.5 km；新鮮 TLE（齡 < 0.5 天）中位數僅 1.68 km。",
        "誤差隨 TLE 年齡快速增長：< 0.5 天 1.7 km → 1–1.5 天 13.9 km → 1.5–2 天 29.5 km。",
        "沿軌（T）分量主導——符合 SGP4 對半長軸/平均運動最敏感的物理，解釋 TLE 難偵測微小面外機動。",
        "誤差尾巴（P99 ≈ 297 km、最大 4,451 km）來自觀測期內實際機動的衛星，TLE 無法即時反映。",
    ]:
        _bullet(doc, b)
    _figure(doc, "fig2_study1_tle_error.png",
            "圖 2　(a) TLE 位置誤差 ECDF；(b) 誤差隨 TLE 齡增長（中位數＋IQR 帶）。")

    # 四、MEME vs MEME
    _heading(doc, "四、結果 B：MEME vs MEME — 精密星曆的自我預測外推誤差")
    for b in [
        "8h：P50 ≈ 0.09 km；24h：≈ 0.7 km；48h：≈ 1.6 km；72h：≈ 1.8 km，48h 後趨於飽和。",
        "誤差同樣沿軌主導；P90 尾巴（~8 km）對應少數在時程內機動的衛星。",
        "全程次公里～數公里，較同時程公開 TLE 低約一個數量級——定量證實 MEME 作為近真值參考的正當性。",
    ]:
        _bullet(doc, b)
    _figure(doc, "fig3_study2_meme_self.png",
            "圖 3　MEME 自我預測位置誤差 vs 外推時程（0–72h），中位數與 P90。")

    # 五、凍結 + 斷點
    _heading(doc, "五、結果 C：TLE 凍結外推曲線與下載斷點驗證")
    _heading(doc, "5.1　TLE 凍結退化曲線（1–7 天）", level=2)
    _bullet(doc, "P50：1 天 8.6 km → 3 天 44 km → 7 天 118–157 km，隨齡近似冪次增長，沿軌主導。")
    _figure(doc, "fig4_study3a_frozen.png",
            "圖 4　TLE 凍結外推退化曲線，1–7 天（50 顆，已濾除機動）。", width=5.2)
    _heading(doc, "5.2　下載斷點 ~7 天 spot-check（含機動過濾）", level=2)
    _para(doc, "MEME 下載於 2026 年 6 月中曾中斷約 7 天。取斷點前最後一筆 TLE 外推越過斷點至斷點後"
               "第一筆 MEME 真值。必須先過濾斷點期間機動的衛星，否則機動 ΔV 主導殘差：")
    _table(doc, ["子集", "n", "7 天誤差 P50", "P90", "最大"], [
        ["純外推（未機動）", "38", "246 km", "741 km", "1,470 km"],
        ["機動衛星", "12", "3,473 km", "8,786 km", "12,671 km"],
    ], widths=[1.8, 0.7, 1.5, 1.3, 1.3])
    _para(doc, "")
    _bullet(doc, "機動衛星中位誤差是純外推的 ~14 倍；純外推組最乾淨 5 顆 7 天僅 5–15 km（TLE 純外推可信下限）。")
    _bullet(doc, "robust 機動過濾：以整檔前 ~10 個軌道之平均半長軸（消去 J2 短週期振盪）偵測 8h 階躍"
                 "（門檻 0.2 km）與跨斷點淨移動（阻力感知門檻 0.6 km）；已知靜止衛星雜訊地板 ≤ 0.18 km。")
    _figure(doc, "fig5_study3b_gap.png",
            "圖 5　斷點 spot-check：機動過濾將兩族群清楚分離（對數縱軸）。", width=5.4)

    # 六、綜合比較
    _heading(doc, "六、綜合比較（最有意義的比較項目）")
    _para(doc, "將三種預測的「中位位置誤差 vs 時程」疊於同一對數圖，呈現清楚的精度層級："
               "MEME 自我預測 ≪ TLE（實務，選最新）≪ TLE（凍結，單筆老化）。")
    _figure(doc, "fig6_hierarchy.png",
            "圖 6　預測誤差層級：同一時程下 MEME 精度最高、凍結 TLE 最差。")
    _para(doc, "其他值得延伸的比較項目（已具資料基礎）：RTN 分量隨時程的演化、TLE 誤差 vs 太陽通量 F10.7 / B*、"
               "誤差 vs 衛星高度帶（500 vs 550 km shell）。", size=10, color=MUTED)

    # 七、後續研究建議
    _heading(doc, "七、後續研究建議")
    for i, b in enumerate([
        "擴至全 285 顆，給出艦隊級統計與信賴區間。",
        "MEME 長時程外推（>3 天）需數值傳播器（J2＋阻力＋SRP＋日月）自 MEME 初始狀態外推；"
        "本報告斷點 spot-check 以 TLE 為外推載具，量到的是 TLE 而非 MEME 的外推誤差。",
        "精修小機動過濾：純外推殘餘尾巴源於 < 0.2 km 小機動偵測不到，屬 TLE/osculating 方法物理下限。",
        "回饋 Layer 1/3：以 TLE 誤差 vs 齡曲線校準自適應閾值 θ_a(age) 與 Layer 3 守門特徵。",
    ], 1):
        p = doc.add_paragraph(style="List Number")
        _run(p, b)

    out = RPT / "MEME_TLE_comparison_interim.docx"
    doc.save(str(out))
    print(f"[docx] {out}")


# ════════════════════════════════════════════════════════════════════════════
#  PPT
# ════════════════════════════════════════════════════════════════════════════

def _set_font(run, size=18, bold=False, color=INK):
    run.font.size = PPt(size); run.font.bold = bold
    run.font.color.rgb = PColor.from_string(color)
    run.font.name = CJK
    rpr = run._r.get_or_add_rPr()
    ea = rpr.makeelement(qn("a:ea"), {"typeface": CJK})
    rpr.append(ea)


SW, SH = PInches(13.333), PInches(7.5)


def _slide(prs, title, subtitle=None, page=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 頂部細色條
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, PInches(0.16))
    bar.fill.solid(); bar.fill.fore_color.rgb = PColor.from_string(BLUE); bar.line.fill.background()
    # 標題
    tb = s.shapes.add_textbox(PInches(0.6), PInches(0.35), PInches(12.1), PInches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    _set_font(tf.paragraphs[0].add_run(), 27, True, BLUE)
    tf.paragraphs[0].runs[0].text = title
    if subtitle:
        p = tf.add_paragraph(); _set_font(p.add_run(), 14, False, MUTED)
        p.runs[0].text = subtitle
    # 頁碼
    if page is not None:
        pt = s.shapes.add_textbox(PInches(12.5), PInches(7.0), PInches(0.7), PInches(0.4))
        _set_font(pt.text_frame.paragraphs[0].add_run(), 10, False, MUTED)
        pt.text_frame.paragraphs[0].runs[0].text = str(page)
    return s


def _bullets(slide, items, left, top, width, height, size=16):
    tb = slide.shapes.add_textbox(PInches(left), PInches(top), PInches(width), PInches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = PPt(8)
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        r = p.add_run(); _set_font(r, size - (2 if lvl else 0), False, INK)
        r.text = ("– " if lvl else "• ") + it
        p.level = lvl
    return tb


def _pic(slide, name, left, top, width=None, height=None):
    path = FIG / name
    if not path.exists():
        return
    kw = {}
    if width: kw["width"] = PInches(width)
    if height: kw["height"] = PInches(height)
    slide.shapes.add_picture(str(path), PInches(left), PInches(top), **kw)


def _figbox(slide, left, top, width, height, hexcolor=LIGHT):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(left), PInches(top),
                                 PInches(width), PInches(height))
    box.fill.solid(); box.fill.fore_color.rgb = PColor.from_string(hexcolor)
    box.line.color.rgb = PColor.from_string("D8DEE3"); box.line.width = PPt(0.75)
    box.shadow.inherit = False
    return box


def build_pptx():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH

    # 1. 封面
    s = prs.slides.add_slide(prs.slide_layouts[6])
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, PInches(2.4), SW, PInches(2.7))
    band.fill.solid(); band.fill.fore_color.rgb = PColor.from_string(BLUE); band.line.fill.background()
    tb = s.shapes.add_textbox(PInches(0.8), PInches(2.7), PInches(11.7), PInches(2.1))
    tf = tb.text_frame; tf.word_wrap = True
    _set_font(tf.paragraphs[0].add_run(), 38, True, "FFFFFF")
    tf.paragraphs[0].runs[0].text = "MEME 精密星曆 vs. 公開 TLE"
    p = tf.add_paragraph(); _set_font(p.add_run(), 22, False, "EAF3F9")
    p.runs[0].text = "軌道預測誤差與外推特性研究　—　期中報告"
    mb = s.shapes.add_textbox(PInches(0.8), PInches(5.5), PInches(11.7), PInches(1.6))
    mtf = mb.text_frame
    for i, (k, v) in enumerate(META):
        p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
        r = p.add_run(); _set_font(r, 13, False, MUTED); r.text = f"{k}：{v}"

    # 2. 研究背景與目標
    s = _slide(prs, "研究背景與目標", page=2)
    _bullets(s, [
        "公開 TLE：免費、全目錄，但位置誤差公里級、以 SGP4 傳播。",
        "MEME 精密星曆：SpaceX 每顆 Starlink 自報，精度高 1–2 個數量級（次公里～百公尺）。",
        "目標：以 MEME 為「近似地面真相」，反向量化公開 TLE 的誤差與外推特性。",
        "呼應契約「專案背景」對 TLE 精度限制之關切，為 Layer 1 誤報來源提供實測依據。",
    ], 0.7, 1.7, 12.0, 4.5, size=19)

    # 3. 核心概念：MEME 檔結構
    s = _slide(prs, "核心概念：MEME 檔結構與「第一筆＝近真值」", page=3)
    _figbox(s, 0.7, 1.7, 7.4, 4.6)
    _pic(s, "fig1_concept.png", 0.85, 2.0, width=7.1)
    _bullets(s, [
        "每檔涵蓋 72h、1 分鐘間隔。",
        "第一筆（外推齡 0）最接近定軌真值。",
        "相鄰檔 ~8h 發布、重疊 ~88%。",
        "衍生兩類比較：",
        ("MEME vs TLE：量 TLE 誤差", 1),
        ("MEME vs MEME：量 MEME 外推誤差", 1),
    ], 8.4, 1.9, 4.5, 4.4, size=16)

    # 4. 方法總覽
    s = _slide(prs, "方法總覽：三支分析程式", page=4)
    _bullets(s, [
        "study1 — MEME vs TLE 誤差分布：每檔第一筆為近真值，選最新先行 TLE 傳播比對。",
        "study2 — MEME vs MEME（0–72h）：晚檔真值 vs 早檔外推，無需外部傳播器。",
        "study3 — TLE 凍結曲線（1–7 天）＋ 斷點 spot-check（含 robust 機動過濾）。",
        "誤差一律以 RTN 分解（R 徑向 / T 沿軌 / N 面外），主指標為 3D 位置誤差。",
        "全部重用既有 compare_tle_vs_ephemeris.py 基礎設施，未改原檔。",
    ], 0.7, 1.7, 12.0, 4.5, size=18)

    # 5. 結果 A
    s = _slide(prs, "結果 A：MEME vs TLE — 公開 TLE 誤差分布", page=5)
    _figbox(s, 6.5, 1.7, 6.5, 4.7)
    _pic(s, "fig2_study1_tle_error.png", 6.6, 2.9, width=6.3)
    _bullets(s, [
        "P50 = 2.5 km；新鮮 TLE(<0.5天) 僅 1.68 km。",
        "隨齡快增：1.7 → 13.9 → 29.5 km（0.5→2 天）。",
        "沿軌 T 分量主導。",
        "尾巴（最大 4,451 km）來自實際機動衛星。",
    ], 0.7, 1.9, 5.6, 4.2, size=17)

    # 6. 結果 B
    s = _slide(prs, "結果 B：MEME vs MEME — 自我預測外推誤差", page=6)
    _figbox(s, 6.5, 1.7, 6.5, 4.7)
    _pic(s, "fig3_study2_meme_self.png", 6.7, 2.4, width=6.1)
    _bullets(s, [
        "8h ≈ 0.09 km；24h ≈ 0.7 km；48h ≈ 1.6 km；72h ≈ 1.8 km。",
        "全程次公里～數公里，低於同時程 TLE 約一個數量級。",
        "無需外部傳播器、無 SGP4 誤差混入。",
        "定量證實 MEME 作為近真值參考的正當性。",
    ], 0.7, 1.9, 5.6, 4.2, size=17)

    # 7. 結果 C-1 凍結曲線
    s = _slide(prs, "結果 C-1：TLE 凍結外推退化曲線（1–7 天）", page=7)
    _figbox(s, 6.5, 1.7, 6.5, 4.7)
    _pic(s, "fig4_study3a_frozen.png", 6.7, 2.4, width=6.1)
    _bullets(s, [
        "凍結單筆 TLE、任其老化外推。",
        "P50：1 天 8.6 km → 3 天 44 km → 7 天 ~120–157 km。",
        "隨齡近似冪次增長，沿軌主導。",
        "已濾除機動衛星（見下頁方法）。",
    ], 0.7, 1.9, 5.6, 4.2, size=17)

    # 8. 結果 C-2 斷點（頭條）
    s = _slide(prs, "結果 C-2：下載斷點 ~7 天 spot-check（頭條）", page=8)
    _figbox(s, 6.5, 1.7, 6.5, 4.9)
    _pic(s, "fig5_study3b_gap.png", 6.65, 2.5, width=6.2)
    _bullets(s, [
        "純外推(38顆)：7 天 P50 246 km。",
        "機動衛星(12顆)：P50 3,473 km。",
        ("→ 中位誤差差 ~14 倍", 0),
        "最乾淨 5 顆純外推僅 5–15 km。",
        "過濾法：整檔 ~10 軌道平均半長軸偵測階躍。",
    ], 0.7, 1.9, 5.6, 4.5, size=16)

    # 9. 綜合比較
    s = _slide(prs, "綜合比較：預測誤差層級", page=9)
    _figbox(s, 3.4, 1.7, 6.5, 4.9)
    _pic(s, "fig6_hierarchy.png", 3.55, 2.5, width=6.2)
    _bullets(s, [
        "MEME ≪ 實務 TLE ≪ 凍結 TLE。",
        "一圖說明「以 MEME 為真相、TLE 為待驗證」。",
        "延伸項目：",
        ("RTN 分量演化", 1),
        ("誤差 vs F10.7 / B*", 1),
        ("誤差 vs 高度帶", 1),
    ], 0.7, 1.9, 2.6, 4.5, size=15)

    # 10. 後續研究建議
    s = _slide(prs, "後續研究建議", page=10)
    _bullets(s, [
        "擴至全 285 顆，給出艦隊級統計與信賴區間。",
        "MEME 長時程外推(>3 天)需數值傳播器；斷點 spot-check 目前以 TLE 為外推載具。",
        "精修小機動過濾（< 0.2 km 小機動為物理下限）。",
        "回饋 Layer 1/3：以 TLE 誤差-齡曲線校準自適應閾值與守門特徵。",
    ], 0.7, 1.8, 12.0, 4.5, size=19)

    # 11. 結語
    s = _slide(prs, "結語", page=11)
    _bullets(s, [
        "以 MEME 精密星曆建立了公開 TLE 誤差的實測基準（P50 ~2.5 km、隨齡數倍劣化）。",
        "MEME 自我預測誤差次公里級，定量支撐其作為近真值/訓練標籤的正當性。",
        "斷點 spot-check 以機動過濾清楚分離純外推(246 km)與機動(3,473 km)兩族群。",
        "結果可直接回饋 Layer 1 閾值設計與 Layer 3 守門特徵校準。",
    ], 0.7, 1.8, 12.0, 4.5, size=19)

    out = RPT / "MEME_TLE_comparison_interim.pptx"
    prs.save(str(out))
    print(f"[pptx] {out}")


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print("\n完成。")
