#!/usr/bin/env python3
"""
build_review_ppt.py
===================
向審查委員說明的「期中報告簡報」（.pptx），乾淨學術樣式、16:9。
涵蓋全計畫：Layer 1 / Layer 3 / Layer 2 現況 / maneuver_app.py / 本期 MEME vs TLE 研究，
以及落差風險、完成度、後續交付。

輸出：docs/期中報告簡報_審查委員_20260712.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

BLUE, DK, ORANGE, RED, GREEN = "0072B2", "12405C", "E69F00", "D55E00", "009E73"
INK, MUTED, LIGHT = "1A1A1A", "6B6B6B", "F2F6F9"
CJK = "Microsoft JhengHei"

ROOT = Path("f:/GitHub/Sat_TraingDataExtension")
DOCS = ROOT / "docs"
MFIG = DOCS / "meme_tle_report" / "figs"

SW, SH = Inches(13.333), Inches(7.5)


def _font(run, size=18, bold=False, color=INK, italic=False):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.name = CJK
    rpr = run._r.get_or_add_rPr()
    rpr.append(rpr.makeelement(qn("a:ea"), {"typeface": CJK}))


def _txt(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _rect(slide, x, y, w, h, color, line=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor.from_string(color)
    if line:
        r.line.color.rgb = RGBColor.from_string(line); r.line.width = Pt(0.75)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r


def _slide(prs, title, subtitle=None, page=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.333, 0.16, BLUE)
    tf = _txt(s, 0.55, 0.32, 12.2, 1.0)
    _font(tf.paragraphs[0].add_run(), 26, True, BLUE); tf.paragraphs[0].runs[0].text = title
    if subtitle:
        p = tf.add_paragraph(); _font(p.add_run(), 13, False, MUTED); p.runs[0].text = subtitle
    if page is not None:
        pf = _txt(s, 12.5, 7.02, 0.7, 0.4)
        _font(pf.paragraphs[0].add_run(), 10, False, MUTED); pf.paragraphs[0].runs[0].text = str(page)
    return s


def _bullets(slide, items, x, y, w, h, size=17):
    tf = _txt(slide, x, y, w, h)
    for i, it in enumerate(items):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7); p.level = lvl
        r = p.add_run(); _font(r, size - (2 if lvl else 0), bool(lvl == -1), INK)
        r.text = ("– " if lvl == 1 else "• ") + it
    return tf


def _pic(slide, path: Path, x, y, w=None, h=None, frame=True):
    if not path.exists():
        tf = _txt(slide, x, y, w or 4, 0.5)
        _font(tf.paragraphs[0].add_run(), 11, False, RED); tf.paragraphs[0].runs[0].text = f"[缺圖 {path.name}]"
        return
    from PIL import Image
    iw, ih = Image.open(path).size
    if w and not h:
        h = w * ih / iw
    if frame:
        _rect(slide, x - 0.08, y - 0.08, (w or 4) + 0.16, (h or 3) + 0.16, "FFFFFF", line="D8DEE3")
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w) if w else None, Inches(h) if h else None)


def _kpi(slide, items, x, y):
    """一排 KPI 方塊：items = [(數值, 標籤, 顏色), ...]"""
    bw, gap = 2.7, 0.25
    for i, (val, lab, col) in enumerate(items):
        bx = x + i * (bw + gap)
        _rect(slide, bx, y, bw, 1.5, LIGHT, line="D8DEE3")
        tf = _txt(slide, bx + 0.1, y + 0.16, bw - 0.2, 0.8)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _font(tf.paragraphs[0].add_run(), 26, True, col); tf.paragraphs[0].runs[0].text = val
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        _font(p.add_run(), 12, False, INK); p.runs[0].text = lab


def build():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    pg = [1]

    def nextpg():
        pg[0] += 1; return pg[0]

    # 1 封面
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 2.3, 13.333, 2.9, BLUE)
    tf = _txt(s, 0.9, 2.55, 11.6, 2.4)
    _font(tf.paragraphs[0].add_run(), 34, True, "FFFFFF")
    tf.paragraphs[0].runs[0].text = "智慧化低軌通訊衛星軌道異常及太空事件偵測演算法研究"
    p = tf.add_paragraph(); _font(p.add_run(), 22, False, "EAF3F9"); p.runs[0].text = "期中進度報告 — 審查簡報"
    mf = _txt(s, 0.9, 5.55, 11.6, 1.6)
    for k, v in [("計畫案號", "TASA-S-1150268"), ("主辦單位", "國家太空中心（TASA）"),
                 ("執行單位", "社團法人中華民國國防科技學術研究學會"), ("資料基準日", "2026-07-06")]:
        pp = mf.paragraphs[0] if k == "計畫案號" else mf.add_paragraph()
        _font(pp.add_run(), 13, False, MUTED); pp.runs[0].text = f"{k}：{v}"

    # 2 大綱
    s = _slide(prs, "簡報大綱", page=nextpg())
    _bullets(s, [
        "計畫背景與契約三層偵測架構",
        "總體進度概覽",
        "Layer 1：TLE 差分機動偵測（P1–P6、54 天全量驗證）",
        "Layer 3：LightGBM 機動分類器（標籤洩漏修正、外部驗證）",
        "Layer 2：統計偵測層現況與替代方案",
        "maneuver_app.py 視覺化儀表板",
        "本期重點：MEME 精密星曆 vs 公開 TLE 誤差與外推研究",
        "已知落差、風險、完成度與後續交付規劃",
    ], 0.9, 1.7, 11.5, 5.2, size=18)

    # 3 背景與目標
    s = _slide(prs, "計畫背景與目標", page=nextpg())
    _bullets(s, [
        "低軌道 >10,000 顆活躍衛星，機動偵測是太空態勢感知（SSA）核心需求。",
        "痛點：訊號與噪音混疊（大氣阻力 vs 低推力機動）、資料稀疏不均、衛星個體差異。",
        "特別是「偵測靈敏度與誤報率難以兼顧」為國際公認且已量化之難題。",
        "目標：僅用公開 TLE，建立可大規模部署的機動/異常偵測演算法。",
        "契約三層架構：Layer 1 閾值基準層／Layer 2 統計偵測層／Layer 3 AI 偵測層。",
    ], 0.9, 1.7, 11.6, 4.8, size=18)

    # 4 總體進度
    s = _slide(prs, "契約三層架構與總體進度", page=nextpg())
    _kpi(s, [("✅ 完成", "Layer 1 閾值基準層", GREEN),
             ("⚠️ 替代達成", "Layer 2 統計偵測層", ORANGE),
             ("✅ 完成", "Layer 3 傳統 ML", GREEN),
             ("55–60%", "整體完成度（6/24 估）", BLUE)], 0.6, 1.7)
    _bullets(s, [
        "Layer 1：14,090 顆、54 天全量驗證，P1–P6 完成。",
        "Layer 2：契約具名方法（CUSUM/BOCPD/SSA）未實作，已以等效統計方法替代（詳後）。",
        "Layer 3：LightGBM/XGBoost 完成；LSTM/Transformer 依排程屬第三～五月，未逾期。",
        "maneuver_app.py 儀表板完成度顯著領先 6/24 描述；本期新增 MEME vs TLE 誤差研究。",
    ], 0.9, 3.7, 11.6, 3.0, size=16)

    # 5 Layer1 方法
    s = _slide(prs, "Layer 1｜TLE 差分機動偵測：方法", "六級誤報抑制策略 P1–P6", page=nextpg())
    _pic(s, DOCS / "fig3_flowchart.png", 7.0, 1.7, w=6.0)
    _bullets(s, [
        "以連續 TLE 之軌道根數差值（Δa、Δi、Δe、ΔRAAN）識別機動。",
        "J2 攝動修正後計算 RAAN 殘差，去除自然漂移誤判。",
        "P1 單調衰減抑制、P2 高度自適應閾值、P3 B* 輔助、",
        ("P4 多窗口補充、P5 F10.7 太陽通量、P6 星座專屬閾值。", 1),
    ], 0.9, 1.9, 5.9, 4.5, size=16)

    # 6 Layer1 結果
    s = _slide(prs, "Layer 1｜消融實驗與全量指標", page=nextpg())
    _pic(s, DOCS / "fig6_fp_waterfall.png", 7.0, 1.9, w=6.0)
    _kpi(s, [("26.9%", "Overall Recall", BLUE), ("5.4%", "FAR", GREEN),
             ("98.2%", "Precision@1000", BLUE)], 0.7, 1.8)
    _bullets(s, [
        "P1–P4 消融：假陽性 68 → 29，精確率 94.8% → 97.5%。",
        "P1 貢獻最大（−40% FP），對應純大氣阻力衰減。",
        "54 天、14,090 顆全量評估取得上列指標。",
    ], 0.9, 3.6, 5.9, 3.0, size=16)

    # 7 Layer1 驗證
    s = _slide(prs, "Layer 1｜多角度獨立驗證", page=nextpg())
    _pic(s, DOCS / "fig7_recall_at_n.png", 6.9, 1.8, w=6.1)
    _bullets(s, [
        "Recall@N 信心排名：N=1000 時 Precision 維持 95–98%。",
        "獨立 MEME Hold-out 事件驗證（99 個真實 V 形事件）：",
        ("事件級 Recall 57.6%，平均偵測前置 24.4 小時。", 1),
        "跨時間段穩定性測試一致率 89.7%。",
        "CDM 弱監督為負向結果，誠實揭露、指出改進方向。",
    ], 0.9, 1.9, 5.7, 4.6, size=16)

    # 8 Layer3 標籤洩漏
    s = _slide(prs, "Layer 3｜LightGBM：發現並修正標籤洩漏", page=nextpg())
    _pic(s, DOCS / "paper2_fig1_ml_pipeline.png", 7.1, 2.0, w=5.9)
    _bullets(s, [
        "20 個衛星級聚合特徵；衛星層級分層切分杜絕資料洩漏。",
        "發現 4 個特徵（flag_rate 等）與標籤共用規則 → tautological，AUC 虛高至 1.0。",
        "剔除後新增 da_monotonic_decay、bstar_f107_normalized。",
        "方法論教訓：SHAP 無法自動偵測標籤洩漏，需追溯計算是否獨立於標籤。",
    ], 0.9, 2.0, 6.0, 4.4, size=15.5)

    # 9 Layer3 效能 + 泛化落差
    s = _slide(prs, "Layer 3｜效能與外部驗證（誠實揭露泛化落差）", page=nextpg())
    _pic(s, DOCS / "paper2_fig5_roc_comparison.png", 7.0, 2.0, w=6.0)
    _kpi(s, [("99.5%", "Precision（測試集）", BLUE), ("97.5%", "Recall（測試集）", BLUE),
             ("39.7%", "Recall（外部 MEME）", RED)], 0.7, 1.75)
    _bullets(s, [
        "修正後三種樹模型（RF/XGB/LGBM）效能相近，AUC≈0.996。",
        "初版「LightGBM 大幅領先」係洩漏特徵所致，非演算法優勢。",
        "外部驗證（MEME GT）Recall 降至 39.7% — 泛化落差誠實揭露。",
    ], 0.9, 3.55, 5.9, 3.0, size=15.5)

    # 10 Layer2 現況
    s = _slide(prs, "Layer 2｜統計偵測層現況與澄清", page=nextpg())
    _bullets(s, [
        "契約指定之 CUSUM / BOCPD / SSA 三種具名方法尚未實作。",
        "已以三項等效統計方法達成同層功能：",
        ("P4 多窗口滑動掃描", 1),
        ("3σ MAD 同儕比較異常偵測（anomaly_detector.py，適用非合作目標）", 1),
        ("EP 漂移殘差同儕比較（ep_slope_detector.py，聯集 Recall +30.9pp）", 1),
        "重要澄清：XGBoost / LightGBM 屬 Layer 3（AI 偵測層），非 Layer 2 替代。",
        "下一步：將既有統計方法整合進儀表板，而非從零重寫 CUSUM/BOCPD/SSA。",
    ], 0.9, 1.7, 11.6, 5.0, size=17)

    # 11 App
    s = _slide(prs, "maneuver_app.py｜視覺化儀表板", "含真實資料驗證之誤報抑制案例", page=nextpg())
    _pic(s, DOCS / "fig11_gap_suppression_case.png", 6.9, 1.9, w=6.1)
    _bullets(s, [
        "10+ 分析頁籤、三種偵測管線、SSA-RAG 問答整合。",
        "軌域自動路由（LEO/MEO/GEO）、ML 模型軌域守門警告。",
        "案例 NORAD 44349：TLE 追蹤缺口致假機動，",
        ("新增缺口守門（門檻 48h）；298 視窗中 50 個正確壓制。", 1),
    ], 0.9, 2.0, 5.7, 4.4, size=16)

    # 12 本期重點：MEME vs TLE 動機方法
    s = _slide(prs, "本期重點｜MEME vs TLE：動機與方法", page=nextpg())
    _pic(s, MFIG / "fig1_concept.png", 6.7, 2.0, w=6.3)
    _bullets(s, [
        "以 SpaceX MEME 精密星曆（精度高 TLE 約一個數量級）為近似真值。",
        "洞見：每個 MEME 檔第一筆外推齡最小、最接近定軌真值。",
        "三支研究：TLE 誤差分布、MEME 自我預測、TLE 凍結曲線 + 斷點。",
        "分析 50 顆代表性子集（35 機動 + 15 靜止）。",
    ], 0.9, 2.0, 5.6, 4.4, size=16)

    # 13 MEME vs TLE 結果
    s = _slide(prs, "MEME vs TLE｜誤差量測結果", page=nextpg())
    _pic(s, MFIG / "fig2_study1_tle_error.png", 0.7, 3.5, w=6.0)
    _pic(s, MFIG / "fig3_study2_meme_self.png", 7.0, 3.2, w=6.0)
    _bullets(s, [
        "公開 TLE：位置誤差中位 2.5 km（新鮮 <0.5 天僅 1.68 km），隨齡數倍劣化，沿軌主導。",
        "MEME 自我預測：8h 0.09 km → 48h 1.6 km 後飽和，低於同時程 TLE 約一個數量級。",
    ], 0.9, 1.65, 11.6, 1.6, size=15)

    # 14 凍結 + 斷點（頭條）
    s = _slide(prs, "MEME vs TLE｜斷點驗證與機動過濾（頭條）", page=nextpg())
    _pic(s, MFIG / "fig5_study3b_gap.png", 6.7, 1.9, w=6.3)
    _kpi(s, [("246 km", "純外推 7 天 P50", BLUE), ("3,473 km", "機動衛星 P50", RED)], 0.7, 1.85)
    _bullets(s, [
        "凍結 TLE 外推 7 天；以整檔平均半長軸 robust 過濾機動衛星。",
        "純外推 vs 機動衛星中位誤差相差 ~14 倍。",
        "實證：未建模機動 ΔV 主導殘差，必須濾除方能估純外推誤差。",
    ], 0.9, 3.65, 5.6, 3.0, size=15.5)

    # 15 綜合層級
    s = _slide(prs, "MEME vs TLE｜綜合誤差層級", page=nextpg())
    _pic(s, MFIG / "fig6_hierarchy.png", 3.4, 1.9, w=6.5)
    _bullets(s, [
        "MEME 自我預測 ≪ 實務 TLE ≪ 凍結 TLE。",
        "一圖說明「以 MEME 為真相、TLE 為待驗證」的正當性。",
        "可回饋 Layer 1 自適應閾值與 Layer 3 守門特徵校準。",
    ], 0.7, 2.1, 2.6, 4.4, size=15)

    # 16 落差風險完成度
    s = _slide(prs, "已知落差、風險與完成度", page=nextpg())
    _bullets(s, [
        "Layer 2 契約具名方法未實作（最主要落差），已以等效方法替代、待整合儀表板。",
        "深度學習（LSTM/Transformer）尚未開始，依排程未逾期。",
        "Layer 3 泛化落差（外部 Recall 39.7%）為待改進項目。",
        "風險：MEME API 停服（合成資料備援）、標記資料不足（合成 + 遷移學習）、",
        ("TLE 精度限制微型機動（MEME RTN 管線下限 0.02km）。", 1),
        "整體完成度 55–60%（建議由主持人依尚未開始項目權重綜合裁定）。",
    ], 0.9, 1.7, 11.6, 5.0, size=17)

    # 17 後續交付
    s = _slide(prs, "後續工作與交付規劃", page=nextpg())
    _bullets(s, [
        "將既有統計方法（P4/3σMAD/EP）整合進儀表板，補足 Layer 2。",
        "啟動 LSTM Autoencoder / Transformer 序列模型。",
        "縮小 Layer 3 泛化落差：引入 MEME 等獨立標籤來源。",
        "MEME vs TLE 研究擴至全 285 顆，並回饋 Layer 1/3 閾值校準。",
        "交付：完整原始碼、技術說明文件、教育訓練文件（期末 2026-11-30）。",
    ], 0.9, 1.7, 11.6, 5.0, size=18)

    # 18 總結
    s = _slide(prs, "總結", page=nextpg())
    _bullets(s, [
        "Layer 1（TLE 差分 P1–P6、54 天全量）與 Layer 3（LightGBM，含標籤洩漏修正）完成且可靠。",
        "Layer 2 具名方法未實作、如實揭露，已以等效統計方法替代。",
        "本期新增 MEME vs TLE 研究，量化公開 TLE 誤差並建立精度層級，直接支撐 SSA 應用。",
        "所有量化指標均附具體檔案/函式為證據，圖表以真實計算結果為主，負向結果誠實揭露。",
    ], 0.9, 1.8, 11.6, 4.6, size=18)
    tf = _txt(s, 0.9, 6.2, 11.6, 0.7)
    _font(tf.paragraphs[0].add_run(), 15, True, BLUE)
    tf.paragraphs[0].runs[0].text = "敬請委員指教。"

    out = DOCS / "期中報告簡報_審查委員_20260712.pptx"
    prs.save(str(out))
    print(f"[pptx] {out}  （{len(prs.slides._sldIdLst)} 頁）")


if __name__ == "__main__":
    build()
