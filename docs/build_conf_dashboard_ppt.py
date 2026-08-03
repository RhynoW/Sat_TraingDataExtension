# -*- coding: utf-8 -*-
"""build_conf_dashboard_ppt.py — SatDashboard 研討會投影片(修訂版:瞬時鄰近+TCA 精化,自主新創研發,16:9)。
輸出：docs/conf_ssa_dashboard_2026.pptx"""
import sys
from pathlib import Path
try: sys.stdout.reconfigure(encoding="utf-8")
except Exception: pass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

D = Path(r"f:\GitHub\Sat_TraingDataExtension\docs")
BG=RGBColor(0x0A,0x14,0x2E); PANEL=RGBColor(0x12,0x24,0x4A); ACCENT=RGBColor(0x35,0x9E,0xE8)
INK=RGBColor(0xF2,0xF5,0xFA); MUTE=RGBColor(0x9F,0xB3,0xC8); GOLD=RGBColor(0xFF,0xD5,0x4F)
GOOD=RGBColor(0x57,0xC7,0x8E); WARN=RGBColor(0xE8,0x8A,0x3C); RED=RGBColor(0xE0,0x5A,0x50)
FONT="Microsoft JhengHei"; W,H=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]

def bg(s,c=BG):
    r=s.shapes.add_shape(1,0,0,W,H); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background()
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return r
def box(s,x,y,w,h,c):
    r=s.shapes.add_shape(1,x,y,w,h); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); return r
def txt(s,x,y,w,h,runs,size=18,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,gap=1.05):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[(runs,size,color,bold)]
    for i,item in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=gap; p.space_after=Pt(6)
        for seg in (item if isinstance(item,list) else [item]):
            t,sz,cl,bd=seg; r=p.add_run(); r.text=t; f=r.font; f.name=sz and FONT; f.size=Pt(sz); f.color.rgb=cl; f.bold=bd
    return tb
def head(s,title,sub=None):
    box(s,0,0,W,Inches(1.15),PANEL); box(s,Inches(0.5),Inches(0.28),Inches(0.13),Inches(0.6),ACCENT)
    txt(s,Inches(0.78),Inches(0.16),Inches(11.5),Inches(0.85),
        ([[(title,24,INK,True)]] + ([[(sub,14,MUTE,False)]] if sub else [])),anchor=MSO_ANCHOR.MIDDLE,gap=1.0)
    txt(s,Inches(11.2),Inches(0.35),Inches(1.85),Inches(0.5),[("🛰 SatDashboard",13,MUTE,True)],align=PP_ALIGN.RIGHT)
def foot(s,n):
    txt(s,Inches(0.5),Inches(7.08),Inches(11),Inches(0.34),
        [("互動式 SSA 儀表板:瞬時鄰近篩選與 TCA 精化　|　自主新創研發",11,MUTE,False)])
    txt(s,Inches(12.4),Inches(7.08),Inches(0.7),Inches(0.34),[(str(n),11,MUTE,False)],align=PP_ALIGN.RIGHT)
def new(): s=prs.slides.add_slide(BLANK); bg(s); return s
def img(s,path,x,y,w):
    if Path(path).exists(): s.shapes.add_picture(str(path),x,y,width=w)
def kpi(s,x,y,w,val,label,col=GOLD):
    box(s,x,y,w,Inches(1.5),PANEL)
    txt(s,x,y+Inches(0.18),w,Inches(0.8),[(val,28,col,True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x,y+Inches(1.02),w,Inches(0.42),[(label,12.5,INK,False)],align=PP_ALIGN.CENTER)
def table(s,x,y,colw,rows,rh=0.5,fs=13.5):
    yy=y
    for ri,row in enumerate(rows):
        xx=x
        for ci,cell in enumerate(row):
            hl=isinstance(cell,str) and cell.startswith("*"); c=cell[1:] if hl else cell
            bc=ACCENT if ri==0 else (PANEL if ri%2 else RGBColor(0x0E,0x1E,0x40))
            box(s,xx,yy,Inches(colw[ci]),Inches(rh),bc)
            col=INK if ri==0 else (GOLD if hl else INK)
            txt(s,xx+Inches(0.08),yy,Inches(colw[ci]-0.1),Inches(rh),[(c,fs,col,ri==0 or hl)],
                anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
            xx+=Inches(colw[ci])
        yy+=Inches(rh)

# ── 1 封面 ──
s=new()
box(s,0,Inches(2.0),W,Inches(3.15),PANEL); box(s,Inches(0.9),Inches(2.25),Inches(0.16),Inches(2.65),ACCENT)
txt(s,Inches(1.25),Inches(2.15),Inches(11.3),Inches(2.0),
    [[("模組化互動式太空情勢感知儀表板",28,INK,True)],[("瞬時鄰近篩選、TCA 精化與在地覆蓋預報",21,INK,True)]],gap=1.14)
txt(s,Inches(1.28),Inches(4.0),Inches(11),Inches(0.6),
    [("Instantaneous Proximity Screening with TCA Refinement",15,ACCENT,True)])
txt(s,Inches(1.28),Inches(5.35),Inches(11),Inches(1.1),
    [[("國內航太研討會　口頭報告　|　2026",15,MUTE,False)],
     [("SatDashboard（scenario-advanced01）　·　自主新創研發",14,MUTE,False)]],gap=1.25)

# ── 2 動機 ──
s=new(); head(s,"一、緒論：互動式 SSA 工具 —— 且須正名")
txt(s,Inches(0.6),Inches(1.35),Inches(12.4),Inches(1.3),
    [[("SSA 需能",16,INK,False),("互動探索",16,GOLD,True),
      ("的工具：全目錄鄰近熱點、地面站覆蓋/過頂、三維幾何檢視。",16,INK,False)],
     [("商用平台(STK)昂貴難客製；純 3D 檢視器擅展示但不含全目錄鄰近篩選與 Pc。",15,INK,False)]],gap=1.25)
box(s,Inches(0.6),Inches(3.1),Inches(12.1),Inches(1.55),RGBColor(0x3a,0x1a,0x1a))
txt(s,Inches(0.85),Inches(3.25),Inches(11.6),Inches(1.3),
    [[("⚠ 關鍵正名：瞬時鄰近 ≠ conjunction",16,RED,True)],
     [("LEO 相對速 7–14 km/s，兩物體停留 10 km 內僅約 1–2 秒；單快照大量納入共軌群/編隊/",14.5,INK,False)],
     [("恰好經過者。本文採「瞬時鄰近篩選」+ TCA 精化，避免方法命名超出能力。",14.5,INK,False)]],gap=1.15)
txt(s,Inches(0.6),Inches(4.9),Inches(12.4),Inches(0.5),[("四項貢獻：",16,GOLD,True)])
txt(s,Inches(0.6),Inches(5.4),Inches(12.4),Inches(1.5),
    [[("① 全目錄近即時鄰近篩選(首次0.35s/穩態52ms)　② TCA 精化分類辨真接近(8→4)",15,INK,False)],
     [("③ 受控對照量化資料時效(同星系 過期31 vs 新鮮8)　④ 單體→模組化(128測試/Docker)+SOCRATES驗證",15,INK,False)]],gap=1.3)
foot(s,2)

# ── 3 相關研究與定位 ──
s=new(); head(s,"二、相關研究與系統定位")
txt(s,Inches(0.6),Inches(1.4),Inches(12.4),Inches(4.4),
    [[("● 傳播：",16,GOLD,True),("SGP4/SDP4(Hoots 1980、Vallado 2006)；本系統用批次化 SatrecArray。",16,INK,False)],
     [("● 作業級接近篩選＝多階時間視窗濾波：",16,GOLD,True),
      ("Hoots-Crawford-Roehrich (1984) 遠/近地點+幾何+時間濾波；Alfano (2012) 環面路徑濾波。",16,INK,False)],
     [("　→ 本系統 k-d tree 為",15,ACCENT,True),("單一時刻",15,ACCENT,True),
      ("空間鄰近查詢，與時間視窗法分工：快照圈熱點 → TCA 精化沿時間求極值。",15,INK,False)],
     [("● 全目錄接近報告：",16,GOLD,True),
      ("CelesTrak SOCRATES、NASA CARA、KeepTrack —— 本文以 SOCRATES 為正確性交叉比對基準。",16,INK,False)],
     [("● Pc：",16,GOLD,True),
      ("嚴謹需相對協方差投影(Chan 2008 等效面積)；本系統採其啟發之各向同性首階近似作快速排序。",16,INK,False)]],gap=1.35)
txt(s,Inches(0.6),Inches(6.5),Inches(12.4),Inches(0.5),
    [[("定位：互動式初篩、熱點偵測與態勢展示；輸出為候選排序，非取代作業級 CDM。",14.5,GOLD,True)]])
foot(s,3)

# ── 4 架構 ──
s=new(); head(s,"三、系統架構：模組化分層","3,880 行單體 Flask → 分層套件；物理層不依賴 Flask")
table(s,Inches(0.6),Inches(1.75),[1.8,3.6,6.6],
      [["層","模組","職責"],
       ["config","settings／stations／rules","常數唯一來源；SSN 29 站與分類規則外部化"],
       ["ingestion","db／metadata／index／spacetrack","DuckDB、衛星索引與時間軸、Space-Track 降級"],
       ["*physics","*coords／propagate／coverage／conjunction","*向量化 SGP4、k-d tree 鄰近 + TCA 精化 + 近似 Pc"],
       ["services","passes_service","過頂預報背景服務"],
       ["api","Flask Blueprints","31 個 REST 端點"],
       ["web","templates／static","3D 地球儀 + 2D 台北覆蓋"]],rh=0.64,fs=13)
txt(s,Inches(0.6),Inches(6.65),Inches(12.4),Inches(0.5),
    [[("create_app() 工廠模式；Docker 封裝；設定熱重載；128 項 pytest。",13.5,MUTE,False)]])
foot(s,4)

# ── 5 方法:SGP4 + 瞬時鄰近 ──
s=new(); head(s,"四、方法：向量化 SGP4 + 瞬時鄰近篩選")
txt(s,Inches(0.6),Inches(1.35),Inches(12.4),Inches(1.6),
    [[("① 向量化 SGP4：",16,GOLD,True),("SatrecArray 一次批次傳播全目錄至同一 epoch(退回逐顆)；",16,INK,False)],
     [("　SatrecArray 首次由 TLE 解析建構，之後",15,INK,False),("跨快照重用",15,ACCENT,True),("→ 穩態延遲遠低於首次。",15,INK,False)],
     [("② 瞬時鄰近篩選：",16,GOLD,True),("k-d tree query_pairs(d) 枚舉當下時刻相互距離 < 閾值之配對。",16,INK,False)]],gap=1.2)
box(s,Inches(1.2),Inches(3.85),Inches(11),Inches(1.0),PANEL)
txt(s,Inches(1.4),Inches(3.95),Inches(10.6),Inches(0.8),
    [("暴力 O(N²) ≈ 5×10⁸ 次　→　k-d tree O(N log N) + 輸出敏感枚舉",17,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
box(s,Inches(0.6),Inches(5.2),Inches(12.1),Inches(1.3),RGBColor(0x3a,0x1a,0x1a))
txt(s,Inches(0.85),Inches(5.32),Inches(11.6),Inches(1.1),
    [[("⚠ 產出為「瞬時鄰近候選」，非逼近中的 conjunction",15.5,RED,True)],
     [("多數為共軌群/編隊/恰好經過者 —— 須經 TCA 精化辨別(下頁)。",15,INK,False)]],gap=1.2)
foot(s,5)

# ── 6 方法:TCA 精化 ──
s=new(); head(s,"四、方法：TCA 精化（讓 conjunction 站得住腳）")
txt(s,Inches(0.6),Inches(1.35),Inches(12.4),Inches(1.3),
    [[("對每一鄰近候選對，沿當下 ±1 軌道週期(±55 分)粗掃 5s + 細掃 0.25s 求相對距離極小值，",16,INK,False)],
     [("得",16,INK,False),("真正 TCA、最小 miss distance、相對速度",16,GOLD,True),
      ("；候選僅數個~數十，成本極小(約 10 ms/對)。取兩物體週期之較長者。",16,INK,False)]],gap=1.25)
txt(s,Inches(0.6),Inches(3.1),Inches(12.4),Inches(0.5),[("相對速度即判別(三類)：",16,ACCENT,True)])
for i,(t,d,c) in enumerate([
    ("快速穿越 (>1 km/s)","真實相對運動→conjunction 候選;miss<快照者為『真接近中』",WARN),
    ("中速 (0.1–1 km/s)","邊界情形，個別檢視",GOLD),
    ("低相對速/共軌 (<0.1)","多為編隊/共軌;惟低相對速為碰撞評估難案,須 3D 非線性法(非本文2D適用)",GOOD)]):
    y=Inches(3.6+i*1.02); box(s,Inches(0.6),y,Inches(3.9),Inches(0.88),PANEL)
    txt(s,Inches(0.78),y,Inches(3.7),Inches(0.88),[(t,14,c,True)],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(4.7),y,Inches(8.0),Inches(0.88),[(d,13.5,INK,False)],anchor=MSO_ANCHOR.MIDDLE,gap=1.05)
foot(s,6)

# ── 7 方法:近似 Pc ──
s=new(); head(s,"四、方法：各向同性首階近似 Pc（僅快速穿越類）")
txt(s,Inches(0.6),Inches(1.35),Inches(12.4),Inches(0.8),
    [[("對",16,INK,False),("快速穿越類",16,GOLD,True),("候選以精化 miss distance m 代入受 Chan (2008) 啟發之各向同性首階近似：",16,INK,False)]],gap=1.2)
box(s,Inches(1.2),Inches(2.3),Inches(11),Inches(1.0),PANEL)
txt(s,Inches(1.4),Inches(2.4),Inches(10.6),Inches(0.8),
    [("σ² = σ_r² + σ_t² ，　P_base = 1 − exp( −r_A² / 2σ² )",17,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
box(s,Inches(1.2),Inches(3.4),Inches(11),Inches(1.0),PANEL)
txt(s,Inches(1.4),Inches(3.5),Inches(10.6),Inches(0.8),
    [("Pc = P_base · exp( −m² / 2σ² )",17,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
box(s,Inches(0.6),Inches(4.7),Inches(12.1),Inches(1.7),RGBColor(0x2a,0x22,0x12))
txt(s,Inches(0.85),Inches(4.82),Inches(11.6),Inches(1.5),
    [[("明確界定(避免審稿誤解)：",15.5,GOLD,True)],
     [("• 取 Chan 級數首項、假設圓對稱協方差，非其各向異性等效面積解。",14.5,INK,False)],
     [("• σ_r=0.1／σ_t=0.5 km 為代表性量級假設，非個別物體真實協方差(TLE 誤差隨高度/資料齡/太陽活動變化)。",14.5,INK,False)],
     [("• 僅施於快速穿越類;共軌/低相對速類不輸出 Pc(N/A)，須改用 3D 非線性法(Coppola 2012)。",14.5,INK,False)]],gap=1.15)
foot(s,7)

# ── 8 覆蓋 + 資料層 ──
s=new(); head(s,"四、在地覆蓋 · 資料層 · 合規")
txt(s,Inches(0.6),Inches(1.35),Inches(12.4),Inches(1.5),
    [[("③ 在地覆蓋/過頂：",16,GOLD,True),("台北(25.03°N/121.57°E)相對仰角/方位，遮罩 5°；覆蓋半徑 2,000 km 作",16,INK,False),
      ("LEO 聚焦預過濾",16,ACCENT,True),("。",16,INK,False)],
     [("　±30 天時間軸回放；適用性註：GEO/MEO 尚穩，LEO 過頂時刻誤差隨天數增長。",14.5,MUTE,False)]],gap=1.2)
for i,(t,d) in enumerate([("DuckDB + 快取","索引 TTL 快取、可選 Redis"),
                          ("Space-Track 選配","缺憑證自動降級，不影響核心"),
                          ("使用者自訂資料","TLE/目錄/NORAD 監測/GeoJSON 放檔即生效"),
                          ("資料齡旗標 + 合規","依 epoch 過濾/降權；內部部署、遵循 Space-Track 再散布條款")]):
    x=Inches(0.6+(i%2)*6.25); y=Inches(3.15+(i//2)*1.35)
    box(s,x,y,Inches(6.0),Inches(1.15),PANEL)
    txt(s,x+Inches(0.2),y+0.1,Inches(5.6),Inches(0.95),[[(t,15,ACCENT,True)],[(d,13,INK,False)]],gap=1.12)
foot(s,8)

# ── 9 效能實測 ──
s=new(); head(s,"五、效能實測","單工作站 Intel Raptor Lake 24C/32T；中位數 of 5")
kpi(s,Inches(0.6),Inches(1.45),Inches(2.9),"31,481","有效衛星",ACCENT)
kpi(s,Inches(3.7),Inches(1.45),Inches(2.9),"0.35 s","首次全流程")
kpi(s,Inches(6.8),Inches(1.45),Inches(2.9),"52 ms","穩態/快照(快取後)",GOOD)
kpi(s,Inches(9.9),Inches(1.45),Inches(2.9),"O(N)","傳播近線性擴展")
img(s,D/"fig_ssa_perf.png",Inches(1.15),Inches(3.2),Inches(11.0))
foot(s,9)

# ── 10 資料時效 + TCA 實測(核心,新鮮資料) ──
s=new(); head(s,"五、核心：資料時效 + TCA 精化（新鮮 TLE）","同 Starlink 星系受控對照；行為採 TLE 中位齡 1.2 天")
kpi(s,Inches(0.6),Inches(1.45),Inches(2.9),"31→8","過期→新鮮 候選@10km",WARN)
kpi(s,Inches(3.7),Inches(1.45),Inches(2.9),"~4×","過期資料膨脹(假影)",RED)
kpi(s,Inches(6.8),Inches(1.45),Inches(2.9),"8→4","TCA 收斂為真接近",GOOD)
kpi(s,Inches(9.9),Inches(1.45),Inches(2.9),"5.5 km","快速類 miss 中位",ACCENT)
img(s,D/"fig_ssa_tca.png",Inches(1.15),Inches(3.2),Inches(11.0))
foot(s,10)

# ── 11 功能與介面 ──
s=new(); head(s,"六、系統功能與介面")
txt(s,Inches(0.6),Inches(1.5),Inches(12.4),Inches(5.0),
    [[("● 雙視圖：",16,GOLD,True),("3D 地球儀(CesiumJS) + 2D 台北覆蓋(含時間軸)。",16,INK,False)],
     [("● NORAD 監測：",16,GOLD,True),("多顆同時追蹤 + 軌道弧、每 15 秒更新、上限 50 顆。",16,INK,False)],
     [("● 鄰近/接近面板：",16,GOLD,True),("閾值可調、逐對顯示 miss distance、TCA 與 Pc。",16,INK,False)],
     [("● CDM 與衰減(decay)查詢 API、自訂向量圖層開關。",16,INK,False)],
     [("所有路由/參數/回應格式與原單體版一致 → 重構行為等價。",15,MUTE,False)]],gap=1.4)
foot(s,11)

# ── 12 工程實踐 ──
s=new(); head(s,"七、工程實踐：從研究原型到可維運服務")
for i,(t,d) in enumerate([("模組化","3,880 行單體 → 27 模組，物理層與 Flask 解耦"),
                          ("可測試性","128 項 pytest(座標/傳播/鄰近·TCA/覆蓋/API)"),
                          ("可部署性","Docker 封裝、容器平台就緒"),
                          ("可維運性","設定全外部化 + 熱重載；改 JS/CSS 不需重啟"),
                          ("可擴充性","四類使用者自訂資料放檔即生效")]):
    y=Inches(1.6+i*1.02); box(s,Inches(0.6),y,Inches(3.0),Inches(0.9),PANEL)
    txt(s,Inches(0.78),y,Inches(2.8),Inches(0.9),[(t,16,GOLD,True)],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(3.8),y,Inches(9.0),Inches(0.9),[(d,15,INK,False)],anchor=MSO_ANCHOR.MIDDLE)
foot(s,12)

# ── 13 討論與限制 ──
s=new(); head(s,"八、討論與限制（適用邊界）")
txt(s,Inches(0.6),Inches(1.5),Inches(12.4),Inches(5.0),
    [[("① 近似 Pc：",16,GOLD,True),("各向同性首階、代表性 σ，供排序；不取代作業級 CDM。",16,INK,False)],
     [("② TCA 已解快照高估，惟為兩體幾何：",16,GOLD,True),("未含機動與攝動之協方差傳播。",16,INK,False)],
     [("③ TLE 資料時效：",16,GOLD,True),("靜態展示庫資料齡中位 18 天 → 凸顯資料齡旗標與每日更新之必要。",16,INK,False)],
     [("④ 正確性外部驗證待補：",16,GOLD,True),("SOCRATES 交叉比對已定位為即將步驟。",16,INK,False)],
     [("→ 清楚界定與作業級碰撞評估之分工，不減損初篩/展示價值。",15.5,ACCENT,True)]],gap=1.5)
foot(s,13)

# ── 14 結論 ──
s=new(); head(s,"九、結論與展望")
txt(s,Inches(0.6),Inches(1.5),Inches(12.4),Inches(3.4),
    [[("● 模組化、可全目錄近即時運行之互動式 SSA 平台：",16,INK,False)],
     [("　31,481 顆單快照首次 0.35s / 穩態 52ms；傳播近線性擴展。",16,GOOD,True)],
     [("● TCA 精化(14ms/對)把 76 個 10km 瞬時鄰近收斂為 8 快速穿越、4 真接近，",16,INK,False)],
     [("　量化證明快照計數 ≠ conjunction 計數(66 共軌鄰居)。",16,GOOD,True)],
     [("● 單體 → 含 128 項測試之分層架構，邁向可維運內部服務。",16,INK,False)]],gap=1.35)
txt(s,Inches(0.6),Inches(5.2),Inches(12.4),Inches(1.3),
    [[("展望：",16,GOLD,True),("SOCRATES 正確性比對 → 獨立 ingestion+Redis Streams → WebSocket 推送/告警規則",15,INK,False)],
     [("→ 實測感測器接入、track fusion、STK 匯出。",15,INK,False)]],gap=1.3)
foot(s,14)

# ── 15 結尾 ──
s=new()
box(s,0,Inches(3.0),W,Inches(1.6),PANEL); box(s,Inches(0.9),Inches(3.2),Inches(0.16),Inches(1.2),ACCENT)
txt(s,Inches(1.25),Inches(3.15),Inches(11),Inches(1.3),[("謝謝聆聽　Q & A",38,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(1.28),Inches(4.8),Inches(11),Inches(0.6),
    [("SatDashboard　·　自主新創研發之獨立成果",15,MUTE,False)])

out=D/"conf_ssa_dashboard_2026.pptx"; prs.save(str(out))
print("saved",out,f"({len(prs.slides._sldIdLst)} slides)")
