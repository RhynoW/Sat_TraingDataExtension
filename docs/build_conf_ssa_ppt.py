# -*- coding: utf-8 -*-
"""build_conf_ssa_ppt.py — 研討會口頭報告投影片(SSA 三層機動偵測,16:9)。
輸出：docs/conf_ssa_maneuver_2026.pptx"""
import sys
from pathlib import Path
try: sys.stdout.reconfigure(encoding="utf-8")
except Exception: pass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

D = Path(r"f:\GitHub\Sat_TraingDataExtension\docs")
BG=RGBColor(0x0B,0x1E,0x45); PANEL=RGBColor(0x14,0x2A,0x55); ACCENT=RGBColor(0x2F,0x86,0xE0)
INK=RGBColor(0xF2,0xF5,0xFA); MUTE=RGBColor(0x9F,0xB3,0xC8); GOLD=RGBColor(0xFF,0xD5,0x4F)
GOOD=RGBColor(0x57,0xC7,0x8E); WARN=RGBColor(0xE8,0x8A,0x3C)
FONT="Microsoft JhengHei"; W,H=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]

def bg(s,c=BG):
    r=s.shapes.add_shape(1,0,0,W,H); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background()
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return r

def box(s,x,y,w,h,c):
    r=s.shapes.add_shape(1,x,y,w,h); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); return r

def txt(s,x,y,w,h,runs,size=18,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,gap=1.05):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[(runs,size,color,bold)]
    for i,item in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=gap; p.space_after=Pt(6)
        for seg in (item if isinstance(item,list) else [item]):
            t,sz,cl,bd=seg; r=p.add_run(); r.text=t; f=r.font; f.name=sz and FONT; f.size=Pt(sz); f.color.rgb=cl; f.bold=bd
    return tb

def head(s,title,sub=None):
    box(s,0,0,W,Inches(1.15),PANEL); box(s,Inches(0.5),Inches(0.28),Inches(0.13),Inches(0.6),ACCENT)
    txt(s,Inches(0.78),Inches(0.16),Inches(11.6),Inches(0.85),
        ([[(title,25,INK,True)]] + ([[(sub,14,MUTE,False)]] if sub else [])),anchor=MSO_ANCHOR.MIDDLE,gap=1.0)
    txt(s,Inches(11.3),Inches(0.35),Inches(1.7),Inches(0.5),[("✷ TASA",15,MUTE,True)],align=PP_ALIGN.RIGHT)

def foot(s,n):
    txt(s,Inches(0.5),Inches(7.08),Inches(11),Inches(0.34),
        [("低軌巨型星系軌道機動偵測之三層架構　|　TASA-S-1150268",11,MUTE,False)])
    txt(s,Inches(12.4),Inches(7.08),Inches(0.7),Inches(0.34),[(str(n),11,MUTE,False)],align=PP_ALIGN.RIGHT)

def new(): s=prs.slides.add_slide(BLANK); bg(s); return s
def img(s,path,x,y,w):
    if Path(path).exists(): s.shapes.add_picture(str(path),x,y,width=w)

def kpi(s,x,y,w,val,label,col=GOLD):
    box(s,x,y,w,Inches(1.5),PANEL)
    txt(s,x,y+Inches(0.18),w,Inches(0.8),[(val,30,col,True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x,y+Inches(1.02),w,Inches(0.42),[(label,13,INK,False)],align=PP_ALIGN.CENTER)

def table(s,x,y,colw,rows,header_c=ACCENT,rh=0.46,fs=13.5):
    """rows[0]=表頭；highlight 若 cell 前置 '*' 則以 GOLD 粗體。"""
    yy=y
    for ri,row in enumerate(rows):
        xx=x
        for ci,cell in enumerate(row):
            hl = isinstance(cell,str) and cell.startswith("*")
            c = cell[1:] if hl else cell
            bc = header_c if ri==0 else (PANEL if ri%2 else RGBColor(0x0F,0x24,0x4C))
            box(s,xx,yy,Inches(colw[ci]),Inches(rh),bc)
            col = INK if ri==0 else (GOLD if hl else INK)
            txt(s,xx+Inches(0.08),yy,Inches(colw[ci]-0.1),Inches(rh),
                [(c,fs, col, ri==0 or hl)],anchor=MSO_ANCHOR.MIDDLE,
                align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
            xx+=Inches(colw[ci])
        yy+=Inches(rh)

# ── 1 封面 ──
s=new()
box(s,0,Inches(2.0),W,Inches(3.1),PANEL); box(s,Inches(0.9),Inches(2.25),Inches(0.16),Inches(2.6),ACCENT)
txt(s,Inches(1.25),Inches(2.25),Inches(11.2),Inches(1.9),
    [[("低軌巨型星系軌道機動偵測之三層架構",30,INK,True)],[("方法與驗證",26,INK,True)]],gap=1.12)
txt(s,Inches(1.28),Inches(3.95),Inches(11),Inches(0.6),
    [("A Three-Layer Architecture for Orbital Maneuver Detection of LEO Mega-Constellations",16,ACCENT,True)])
txt(s,Inches(1.28),Inches(5.35),Inches(11),Inches(1.1),
    [[("國內航太研討會　口頭報告　|　2026",15,MUTE,False)],
     [("社團法人中華民國國防科技學術研究學會　|　TASA-S-1150268",14,MUTE,False)]],gap=1.25)

# ── 2 動機 ──
s=new(); head(s,"一、緒論：為何需要自主機動偵測")
txt(s,Inches(0.6),Inches(1.4),Inches(12.4),Inches(2.2),
    [[("● LEO 巨型星系爆發部署 ",17,GOLD,True),("——每週多批、在軌數萬顆，軌道擁擠與碰撞風險同步升高。",17,INK,False)],
     [("● 機動計畫多不公開 ",17,GOLD,True),("——偵測須建立於公開資料（TLE／MEME）與可驗證真值之上。",17,INK,False)],
     [("● 自主 SSA 為國防與民生韌性所需 ",17,GOLD,True),("——不能只依賴外購編目。",17,INK,False)]],gap=1.3)
txt(s,Inches(0.6),Inches(4.0),Inches(12.4),Inches(0.5),[("三項根本挑戰：",18,ACCENT,True)])
for i,(t,d) in enumerate([("真值稀缺且異質","來源可信度差異極大"),
                          ("訊雜比受限","TLE 雜訊達數十公尺，小機動被淹沒"),
                          ("泛化風險","模型易只認得訓練過的衛星／時段")]):
    x=Inches(0.6+i*4.15); box(s,x,Inches(4.6),Inches(3.9),Inches(1.7),PANEL)
    txt(s,x+Inches(0.2),Inches(4.75),Inches(3.6),Inches(0.6),[(f"({i+1}) {t}",16,GOLD,True)])
    txt(s,x+Inches(0.2),Inches(5.4),Inches(3.6),Inches(0.8),[(d,14,INK,False)],gap=1.15)
foot(s,2)

# ── 3 貢獻 ──
s=new(); head(s,"本文貢獻")
items=[("三層互補架構","規則廣掃 → 統計變點 → 梯度提升融合 + 物理路由"),
       ("真值可信度分級","效能宣稱僅採嚴格真值（MEME／ILRS-IDS）"),
       ("三種嚴格泛化協定","GroupKFold OOF、時間外推、未見衛星 hold-out（AUC 0.980）"),
       ("真實全星系部署","67 天 284 星，量化營運負荷（每日 ~1.1 人時）"),
       ("SNR 物理偵測下限","誠實界定純 TLE 對小型機動之能力邊界")]
y=1.5
for i,(t,d) in enumerate(items):
    box(s,Inches(0.6),Inches(y),Inches(0.55),Inches(0.9),ACCENT)
    txt(s,Inches(0.6),Inches(y),Inches(0.55),Inches(0.9),[(str(i+1),22,INK,True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    box(s,Inches(1.25),Inches(y),Inches(11.5),Inches(0.9),PANEL)
    txt(s,Inches(1.45),Inches(y+0.08),Inches(11.1),Inches(0.8),
        [[(t+"　",16,GOLD,True),(d,15,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    y+=1.02
foot(s,3)

# ── 4 架構 ──
s=new(); head(s,"三、系統架構")
img(s,D/"fig3_flowchart.png",Inches(0.6),Inches(1.35),Inches(8.3))
txt(s,Inches(9.1),Inches(1.5),Inches(3.9),Inches(5.4),
    [[("L1 規則廣掃",16,ACCENT,True)],[("P1–P6，高召回、可稽核",13.5,MUTE,False)],
     [("L2 統計變點",16,ACCENT,True)],[("CUSUM／BOCPD／SSA／3σ-MAD，σ 正規化 SNR",13.5,MUTE,False)],
     [("L3 融合",16,ACCENT,True)],[("五通道→15 特徵→HistGB，GroupKFold OOF",13.5,MUTE,False)],
     [("物理路由",16,ACCENT,True)],[("NRLMSIS 阻力殘差、再入守門、域外改道",13.5,MUTE,False)],
     [("三層互相印證，非互相取代。",14,GOLD,True)]],gap=1.15)
foot(s,4)

# ── 5 資料與真值 ──
s=new(); head(s,"三、資料與真值：可信度分級")
table(s,Inches(0.6),Inches(1.5),[3.0,4.2,1.8,3.2],
      [["真值類別","來源","可信度","角色"],
       ["*嚴格真值","MEME 星曆、ILRS/IDS 點火日誌","*高","效能宣稱唯一依據"],
       ["代理真值","推進能力（有無推進器）","低","僅描述覆蓋規模"],
       ["合成真值","已知 Δa 注入真實序列","受控","量化偵測下限"]],rh=0.62)
txt(s,Inches(0.6),Inches(4.5),Inches(12.4),Inches(2.2),
    [[("評估資料集：",17,GOLD,True),("284 顆 Starlink × 67 天全星系密集窗（2026-05-02～07-08、19,066 星日）。",16,INK,False)],
     [("597 正 unit ／ 6,636 負 unit（MEME medium+ 真值）。",16,INK,False)],
     [("● 關鍵：unit（事件）級對齊 ",15,ACCENT,True),
      ("——MEME 8h 網格 vs TLE 不規則 epoch，點級對齊會使 AUC 塌至約 0.55。",15,INK,False)]],gap=1.25)
foot(s,5)

# ── 6 L1+L2 方法 ──
s=new(); head(s,"四、方法：L1 規則 + L2 統計變點")
txt(s,Inches(0.6),Inches(1.35),Inches(12.4),Inches(1.2),
    [[("L1：",16,GOLD,True),("軌道分類 + 相鄰 TLE 元素轉移之六規則 P1–P6，逐 epoch 合併旗標（高召回、完全可稽核）。",16,INK,False)],
     [("L2：",16,GOLD,True),("四統計通道逐點分數，核心為每衛星自身雜訊 σ 之正規化（以 SNR 而非絕對量判定）。",16,INK,False)]],gap=1.25)
box(s,Inches(1.2),Inches(3.4),Inches(11),Inches(1.15),PANEL)
txt(s,Inches(1.4),Inches(3.5),Inches(10.6),Inches(0.95),
    [("SNR(t) = |Δa(t)| / σ_a ，　σ_a = 1.4826 · median( |Δa − median(Δa)| )",18,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
box(s,Inches(1.2),Inches(4.75),Inches(11),Inches(1.0),PANEL)
txt(s,Inches(1.4),Inches(4.85),Inches(10.6),Inches(0.8),
    [("CUSUM：　S_t = max( 0 , S_(t−1) + (x_t − μ) − k )",18,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(0.6),Inches(6.1),Inches(12.4),Inches(0.8),
    [[("σ 正規化使同一門檻可跨不同追蹤品質衛星一致套用",15,ACCENT,True),("（見雜訊底分析）。",15,INK,False)]])
foot(s,6)

# ── 7 L3 融合 ──
s=new(); head(s,"四、方法：L3 梯度提升融合")
txt(s,Inches(0.6),Inches(1.35),Inches(12.4),Inches(1.4),
    [[("五通道（cusum／bocpd／ssa／mad3sig／NRLMSIS drag）→ 每 unit 聚合 15 維特徵（max／mean／p90）。",16,INK,False)],
     [("HistGradientBoosting 融合；GroupKFold OOF（同一衛星不跨 train／test）杜絕洩漏。",16,INK,False)]],gap=1.25)
box(s,Inches(1.2),Inches(3.15),Inches(11),Inches(1.05),PANEL)
txt(s,Inches(1.4),Inches(3.25),Inches(10.6),Inches(0.85),
    [("p = f_HGB(x) ，　alert if p ≥ τ ，　τ = τ(FPR ≤ 0.05)",18,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(0.6),Inches(4.5),Inches(12.4),Inches(2.2),
    [[("● 操作門檻 τ 以「誤報率下限（FPR floor）」法求取 ",15,ACCENT,True),("→ 嚴格控制誤報預算。",15,INK,False)],
     [("● 去除 unit 長度特徵以避免長度洩漏；負窗與正 unit 等寬以避免 AUC 虛高。",15,INK,False)],
     [("● 物理路由：域外／再入目標改走非監督 Model 2，不作量化召回宣稱。",15,INK,False)]],gap=1.3)
foot(s,7)

# ── 8 實驗設計 ──
s=new(); head(s,"五、實驗設計")
txt(s,Inches(0.6),Inches(1.35),Inches(12.4),Inches(1.0),
    [[("主指標：",16,GOLD,True),("ROC-AUC（整體判別力）＋ FPR≤0.05 操作點下之召回率；依量級分層（small 1–5／medium 5–10／large ≥10 km）。",16,INK,False)]],gap=1.2)
txt(s,Inches(0.6),Inches(2.7),Inches(12.4),Inches(0.5),[("三種漸次嚴格之泛化協定：",17,ACCENT,True)])
for i,(t,d) in enumerate([("GroupKFold OOF","同一衛星不跨折；沿族群／高度／時間／品質四軸切片"),
                          ("out-of-time","前 60% 時間訓、後 40% 從未見時段盲測"),
                          ("unseen-satellite","保留 56 顆整組、完全不參與訓練（可部署情境）")]):
    y=Inches(3.3+i*1.15); box(s,Inches(0.6),y,Inches(3.4),Inches(0.95),PANEL)
    txt(s,Inches(0.78),y,Inches(3.2),Inches(0.95),[(chr(97+i)+") "+t,15,GOLD,True)],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(4.2),y,Inches(8.6),Inches(0.95),[(d,15,INK,False)],anchor=MSO_ANCHOR.MIDDLE,gap=1.05)
foot(s,8)

# ── 9 結果：擂台 ──
s=new(); head(s,"六、結果：同一擂台三層比較","284 星 / 67 天 / FPR≤0.05；相同測試集、真值、評估單元")
table(s,Inches(0.6),Inches(1.7),[4.6,2.1,1.9,2.0,1.6],
      [["方法","ROC-AUC","精確率","large 召回","總召回"],
       ["L1 規則 P1–P6","—","0.296","0.333","0.251"],
       ["L2 cusum（最佳單通道）","0.901","0.472","0.469","0.496"],
       ["基線 σ 正規化 |Δa|","0.878","0.337","0.365","0.281"],
       ["*L3 融合（本系統）","*0.985","*0.631","*0.970","*0.946"],
       ["naive 隨機（對照）","—","0.084","0.053","0.051"]],rh=0.6)
txt(s,Inches(0.6),Inches(5.7),Inches(12.4),Inches(1.1),
    [[("L3 AUC 0.985 ≫ 最佳單通道 0.901 ≫ 基線 ≫ 隨機 0.05。",16,GOOD,True)],
     [("三層增量：L3 相對 L1∪L2 淨補漏 202 真機動、淨除誤 333 假警報 → 融合具實質增量。",14.5,INK,False)]],gap=1.2)
foot(s,9)

# ── 10 消融 ──
s=new(); head(s,"六、結果：模型消融")
table(s,Inches(1.2),Inches(1.7),[4.6,2.4,3.2,2.0],
      [["分類器","ROC-AUC","large 召回@FPR≤.05","總召回"],
       ["*HistGB（本系統）","*0.985","*0.970","*0.946"],
       ["LightGBM","0.985","0.970","0.948"],
       ["Logistic（線性基線）","0.975","0.953","0.933"]],rh=0.62)
txt(s,Inches(0.6),Inches(4.4),Inches(12.4),Inches(2.0),
    [[("LightGBM ≈ HistGB ",17,GOLD,True),("→ 0.97 之增益來自「episode-native 資料 + 聚合特徵框架」，非特定演算法。",16,INK,False)],
     [("即便線性分類器亦達 0.975 ",16,ACCENT,True),("→ 進一步佐證框架之主導性。",16,INK,False)]],gap=1.35)
foot(s,10)

# ── 11 泛化(重點) ──
s=new(); head(s,"六、結果：泛化驗證（三協定）")
table(s,Inches(0.6),Inches(1.6),[5.2,2.1,2.1,1.5,1.5],
      [["協定","ROC-AUC","large 召回","FPR","測試 unit"],
       ["GroupKFold OOF（全域）","0.985","0.970","0.050","7,233"],
       ["out-of-time（後 40% 盲測）","0.94","0.800","0.044","2,893"],
       ["*unseen-satellite hold-out","*0.980","*1.000","0.070","1,457"]],rh=0.6)
kpi(s,Inches(1.2),Inches(4.4),Inches(3.4),"0.980","未見衛星 hold-out AUC")
kpi(s,Inches(5.0),Inches(4.4),Inches(3.4),"1.000","未見衛星 large 召回",GOOD)
kpi(s,Inches(8.8),Inches(4.4),Inches(3.4),"0.044","四軸 OOF AUC 全距",ACCENT)
txt(s,Inches(0.6),Inches(6.15),Inches(12.4),Inches(0.9),
    [[("56 顆從未訓練之衛星仍完整偵測大型機動 → 泛化非源於記憶特定衛星（最強證據）。",15.5,GOLD,True)]])
foot(s,11)

# ── 12 部署 ──
s=new(); head(s,"六、結果：全星系 67 天部署營運","真實窗實測（非線性投影），19,066 星日")
kpi(s,Inches(0.6),Inches(1.5),Inches(2.9),"894","67 天總告警")
kpi(s,Inches(3.7),Inches(1.5),Inches(2.9),"13.3","每日告警數",ACCENT)
kpi(s,Inches(6.8),Inches(1.5),Inches(2.9),"1.1","分析員人時／日",GOOD)
kpi(s,Inches(9.9),Inches(1.5),Inches(2.9),"17.4","FAR／千星日",WARN)
txt(s,Inches(0.6),Inches(3.4),Inches(12.4),Inches(2.6),
    [[("● 精確率 0.631；TP 564、FP 330。",16,INK,False)],
     [("● 高信心（p≥0.9）自動升級佔 57% ",16,ACCENT,True),("→ 值勤員優先審少量關鍵告警，其餘批次處理。",16,INK,False)],
     [("● 平均告警延遲中位 0.1–0.5 h ",16,INK,False),("（遠優於 24 h）。",16,MUTE,False)],
     [("結論：單一分析員即可值守全 284 星星座，系統可落地。",17,GOLD,True)]],gap=1.35)
foot(s,12)

# ── 13 雜訊底 + SNR 下限 ──
s=new(); head(s,"六、結果：雜訊底與 SNR 偵測下限")
img(s,D/"fig_sigma_vs_altitude.png",Inches(0.5),Inches(1.35),Inches(6.5))
img(s,D/"fig_gradual_arc.png",Inches(7.2),Inches(1.35),Inches(5.7))
txt(s,Inches(0.6),Inches(6.2),Inches(12.4),Inches(1.0),
    [[("13 顆 DORIS 測高衛星實測 σ≈0.2 m，較 Starlink 級假設 50 m 小兩個數量級 → σ 正規化必要。",14.5,INK,False)],
     [("逐步 SNR < 2 之小機動沒入雜訊底 → 純 TLE 之物理偵測下限（非模型缺陷）；精密星曆為根治路徑。",14.5,GOLD,True)]],gap=1.2)
foot(s,13)

# ── 14 討論：負面結果 ──
s=new(); head(s,"七、討論：能力邊界與負面結果")
txt(s,Inches(0.6),Inches(1.4),Inches(12.4),Inches(1.3),
    [[("能力邊界：",16,GOLD,True),("量化成績為 episode 級；epoch（單筆 TLE）級 large TPR 僅 0.383、AUC 0.572（8h 網格 + SGP4 限制）。",16,INK,False)],
     [("系統定位為大尺度行為改變之候選事件排序器，非單筆 TLE 高精度推力分類器。",15,MUTE,False)]],gap=1.25)
txt(s,Inches(0.6),Inches(3.15),Inches(12.4),Inches(0.5),[("負面結果（模型選型參考）：深度序列模型全面輸給工程特徵 GBM",17,ACCENT,True)])
table(s,Inches(1.6),Inches(3.75),[5.0,3.0],
      [["模型","unit 級 AUC"],
       ["*工程特徵 GBM 融合（本系統）","*0.98"],
       ["LSTM-AutoEncoder","0.70"],
       ["PatchTST","0.67"],
       ["bi-GRU","無有效判別力"]],rh=0.5,fs=14)
txt(s,Inches(0.6),Inches(6.6),Inches(12.4),Inches(0.6),
    [[("機動為稀疏點事件 → 變化點統計之聚合特徵優於有限資料上的原始序列深度學習。",14.5,GOLD,True)]])
foot(s,14)

# ── 15 結論 ──
s=new(); head(s,"八、結論")
txt(s,Inches(0.6),Inches(1.5),Inches(12.4),Inches(4.8),
    [[("① 三層架構於 284 星 / 67 天嚴格真值擂台達 ",17,INK,False),("ROC-AUC 0.985、large 召回 0.970。",17,GOOD,True)],
     [("② 三種泛化協定（未見衛星 hold-out 0.980 / 召回 1.000）證明非記憶式泛化。",17,INK,False)],
     [("③ 真實 67 天全星系部署顯示每日約 1.1 人時之可負荷營運。",17,INK,False)],
     [("④ 量化純 TLE 對小型機動之 SNR 物理下限，誠實界定能力邊界。",17,INK,False)],
     [("",8,INK,False)],
     [("後續：",16,GOLD,True),("以 MEME 原生標籤重訓單衛星模型；擴展精密星曆輸入以突破小型機動偵測下限。",16,INK,False)]],gap=1.4)
foot(s,15)

# ── 16 結尾 ──
s=new()
box(s,0,Inches(3.0),W,Inches(1.6),PANEL); box(s,Inches(0.9),Inches(3.2),Inches(0.16),Inches(1.2),ACCENT)
txt(s,Inches(1.25),Inches(3.15),Inches(11),Inches(1.3),[("謝謝聆聽　Q & A",38,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(1.28),Inches(4.8),Inches(11),Inches(0.6),
    [("TASA-S-1150268　|　社團法人中華民國國防科技學術研究學會",15,MUTE,False)])

out=D/"conf_ssa_maneuver_2026.pptx"; prs.save(str(out))
print("saved",out,f"({len(prs.slides._sldIdLst)} slides)")
