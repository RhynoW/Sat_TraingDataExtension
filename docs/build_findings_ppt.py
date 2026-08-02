# -*- coding: utf-8 -*-
"""build_findings_ppt.py — 期中報告附帶發現簡報（4 則）
產出：docs/期中報告_附帶發現_20260716.pptx
四則發現：神龍 RPO 高度交換／Starlink MEME 預寫推力弧／福衛七星群三發現／LEO-PNT TLE vs MEME
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

DOCS = Path(__file__).parent
FONT = "Microsoft JhengHei"

# ── 配色（深藍太空主題）────────────────────────────────────────────────────
BG      = RGBColor(0x0E, 0x1B, 0x2E)   # 深海軍藍
PANEL   = RGBColor(0x16, 0x27, 0x40)
ACCENT  = RGBColor(0x4F, 0xC3, 0xF7)   # 亮天藍
ACCENT2 = RGBColor(0xFF, 0xD5, 0x4F)   # 琥珀（強調數字）
INK     = RGBColor(0xEC, 0xF2, 0xF9)   # 近白
MUTE    = RGBColor(0x9F, 0xB3, 0xC8)   # 灰藍
GOOD    = RGBColor(0x66, 0xBB, 0x6A)

W, H = Inches(13.333), Inches(7.5)     # 16:9


def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def _bg(slide):
    r = slide.shapes.add_shape(1, 0, 0, W, H); _fill(r, BG)
    r.shadow.inherit = False
    return r


def _txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.0):
    """runs: list of (text, size, color, bold) 或 list of lines(list of runs)。"""
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp
        for (text, size, color, bold) in line:
            r = p.add_run(); r.text = text
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = bold
    return tb


def _chip(slide, x, y, text, color=ACCENT):
    c = slide.shapes.add_shape(1, x, y, Inches(0.14), Inches(0.5)); _fill(c, color)
    _txt(slide, x + Inches(0.28), y - Inches(0.04), Inches(9), Inches(0.6),
         [(text, 24, INK, True)], anchor=MSO_ANCHOR.MIDDLE)


def _pic(slide, path, x, y, w):
    p = DOCS / path
    if p.exists():
        return slide.shapes.add_picture(str(p), x, y, width=w)
    return None


def _footer(slide, n):
    _txt(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.35),
         [("TASA-S-1150268　智慧化低軌通訊衛星軌道異常及太空事件偵測", 10, MUTE, False)])
    _txt(slide, Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.35),
         [(str(n), 10, MUTE, False)], align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, items, size=15, gap=1.28):
    """items: list of (marker_color, [runs])。"""
    tb = slide.shapes.add_textbox(x, y, w, Inches(4.6)); tf = tb.text_frame
    tf.word_wrap = True
    for i, (mc, runs) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = gap; p.space_after = Pt(4)
        b = p.add_run(); b.text = "▍ "; b.font.name = FONT; b.font.size = Pt(size); b.font.color.rgb = mc
        for (text, sz, color, bold) in runs:
            r = p.add_run(); r.text = text
            r.font.name = FONT; r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return tb


prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK); _bg(s); return s


# ══ 封面 ══════════════════════════════════════════════════════════════════
s = slide()
band = s.shapes.add_shape(1, 0, Inches(2.5), W, Inches(2.4)); _fill(band, PANEL)
_txt(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.2),
     [("期中報告　附帶發現彙整", 40, INK, True)])
_txt(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.7),
     [("從公開軌道資料，看見他人的操作意圖與系統的延伸價值", 20, ACCENT, False)])
_txt(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(1.5),
     [[("五則發現：", 16, MUTE, False)],
      [("① 神龍太空梭在軌釋放與高度交換　② Starlink 於 MEME 預寫推力弧", 15, INK, False)],
      [("③ 福衛七星群的艦隊級操作特徵　④ LEO-PNT：TLE 與 MEME 的定位天花板", 15, INK, False)],
      [("⑤ TLE 品質不是常數：偵測與否取決於追蹤品質、如何量化", 15, INK, False)]], sp=1.25)
_txt(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.4),
     [("計畫案號 TASA-S-1150268　｜　社團法人中華民國國防科技學術研究學會　｜　2026-07-16", 11, MUTE, False)])

# ══ 目錄 ══════════════════════════════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.5), "本簡報導覽")
cards = [
    ("①", "神龍 RPO 事件", "58573 → 物體 G(59884)：釋放後約 1 km 量級高度差與對稱交換，加四則國際同型案例"),
    ("②", "MEME 預寫推力弧", "Starlink 把計畫機動預先編入精密星曆；連續數日剖面揭示連續電推機動作業"),
    ("③", "福衛七星群觀點", "六星 95% 機動為艦隊同步機動作業、FS7-3 停止維持、面幾何緩慢重排"),
    ("④", "LEO-PNT 定位天花板", "2,600 萬點實測 TLE(km 級) vs MEME(公尺級)；水平/垂直精度比較與三項創見"),
    ("⑤", "TLE 品質量化", "ILRS 真值校準：追蹤品質決定 σ（差約 100×）；如何量化每顆星的 TLE 品質"),
]
y = Inches(1.28)
for tag, title, desc in cards:
    card = s.shapes.add_shape(1, Inches(0.6), y, Inches(12.1), Inches(1.02)); _fill(card, PANEL)
    _txt(s, Inches(0.85), y + Inches(0.10), Inches(1.1), Inches(0.82),
         [(tag, 36, ACCENT2, True)], anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(2.0), y + Inches(0.08), Inches(3.4), Inches(0.82),
         [(title, 19, INK, True)], anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(5.3), y + Inches(0.08), Inches(7.2), Inches(0.86),
         [(desc, 13.5, MUTE, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.16)

# ══ 發現 ① 神龍 RPO — 高度交換 ═══════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.45), "發現①　神龍太空梭在軌釋放：約 1 公里量級高度交換")
_pic(s, "fig_r6_shenlong_height.png", Inches(0.5), Inches(1.35), Inches(7.7))
bullets(s, Inches(8.5), Inches(1.5), Inches(4.4), [
    (ACCENT2, [("釋放（2024-05-24）", 15, INK, True)]),
    (MUTE, [("物體 G 進入", 13, INK, False), ("約高 1 km", 13, ACCENT2, True), ("的軌道", 13, INK, False)]),
    (ACCENT2, [("高度交換（05-27~28）", 15, INK, True)]),
    (MUTE, [("58573 升至 605.77、G 降至 604.77，", 13, INK, False), ("對稱互換", 13, ACCENT2, True)]),
    (ACCENT2, [("再會合（06-06）", 15, INK, True)]),
    (MUTE, [("交換＝快慢互換→漂移反轉→受控再接近", 13, INK, False)]),
    (GOOD, [("判讀指紋", 15, INK, True)]),
    (MUTE, [("約 1 km 量級高度差＋對稱交換＋漂移反轉＝", 13, INK, False),
            ("受控 RPO（非碎片分離）", 13, GOOD, True)]),
])
_txt(s, Inches(8.5), Inches(6.35), Inches(4.4), Inches(0.7),
     [("註：0.01 km 層級已低於半長軸雜訊 σ；取「約 1 km 量級、穩定維持」之穩健判讀", 10.5, MUTE, False)])
_footer(s, 3)

# ══ 發現 ① 相對幾何重建 + 延伸案例 ════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.45), "發現①　相對運動重建與國際同型案例（皆由公開 TLE 重建）")
_pic(s, "fig_r4_shenlong_release.png", Inches(0.5), Inches(1.35), Inches(6.7))
bullets(s, Inches(7.4), Inches(1.45), Inches(5.5), [
    (ACCENT, [("重建結果", 14, INK, True)]),
    (MUTE, [("05-28 最近距離 ", 12.5, INK, False), ("約 0.2 km", 12.5, ACCENT2, True),
            ("→沿軌分離 1,700 km→06-06 主動再接近", 12.5, INK, False)]),
    (ACCENT, [("(a) 神龍第二飛 (2022)", 13.5, INK, True)]),
    (MUTE, [("釋放 OBJECT J(54218)，多次分離-再接近循環", 12.5, INK, False)]),
    (ACCENT, [("(b) 實踐-21 / SJ-21 (2021)", 13.5, INK, True)]),
    (MUTE, [("GEO 拖曳失效北斗 G2 至墓地軌道", 12.5, INK, False)]),
    (ACCENT, [("(c) Cosmos 2542/2543 (2019)", 13.5, INK, True)]),
    (MUTE, [("套娃釋放子星、伴飛 USA-245", 12.5, INK, False)]),
    (ACCENT, [("(d) Luch/Olymp-K (2014-)", 13.5, INK, True)]),
    (MUTE, [("GEO 長期遊走、停泊各國通訊衛星旁", 12.5, INK, False)]),
])
_txt(s, Inches(7.4), Inches(6.4), Inches(5.5), Inches(0.6),
     [("共同點：全部在公開 TLE 留下可重建的運動學指紋；本計畫工具皆可直接套用", 12, GOOD, True)])
_footer(s, 4)

# ══ 發現 ② MEME 預寫推力弧 ═══════════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.45), "發現②　Starlink 把計畫機動「預先」寫入精密星曆(MEME)")
_pic(s, "fig_r9_planupload.png", Inches(0.5), Inches(1.4), Inches(7.9))
bullets(s, Inches(8.6), Inches(1.5), Inches(4.3), [
    (ACCENT2, [("關鍵觀察", 15, INK, True)]),
    (MUTE, [("相鄰 MEME 檔重疊軌跡平日僅差 ", 13, INK, False), ("91 m", 13, ACCENT2, True)]),
    (MUTE, [("機動已預先編入→無「有無機動」大分歧", 13, INK, False)]),
    (ACCENT2, [("轉為預警前哨", 15, INK, True)]),
    (MUTE, [("機動計畫上傳當下位置差躍升 ", 13, INK, False), ("1,892 m", 13, ACCENT2, True)]),
    (GOOD, [("價值", 15, INK, True)]),
    (MUTE, [("能在機動", 13, INK, False), ("尚未發生、僅剛被編入星曆", 13, GOOD, True),
            ("時即預警——比偵測機動本身更早一步", 13, INK, False)]),
])
_footer(s, 5)

# ══ 發現 ② 連續剖面 ══════════════════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.45), "發現②　連續數日 MEME 剖面：電推升軌是「連續機動作業」而非單次點火")
bullets(s, Inches(0.7), Inches(1.6), Inches(7.2), [
    (ACCENT, [("方法：檔內逐分鐘 a(t) 剖面（burn-arc profiler）", 16, INK, True)]),
    (MUTE, [("由每檔逐分鐘 r,v 以活力積分算瞬時半長軸，", 14, INK, False),
            ("一個軌道週期滾動中位消 J2 短週期", 14, INK, False)]),
    (ACCENT, [("實測：STARLINK-37471 之 +50 km 升軌", 16, INK, True)]),
    (MUTE, [("解析為橫跨數日、含 ", 14, INK, False), ("422 段每軌微點火", 14, ACCENT2, True),
            ("的連續機動作業", 14, INK, False)]),
    (ACCENT, [("時刻定位精度", 16, INK, True)]),
    (MUTE, [("機動作業起點與最陡爬升可定位至", 14, INK, False), ("分鐘級", 14, ACCENT2, True),
            ("（vs MEME 8 小時網格）", 14, INK, False)]),
    (GOOD, [("物理意涵", 16, INK, True)]),
    (MUTE, [("Starlink 電推推力弧長達數小時、逐圈分段——", 14, INK, False),
            ("「單一點火時刻」本身是理想化", 14, GOOD, True)]),
])
# 右側數據卡
card = s.shapes.add_shape(1, Inches(8.4), Inches(1.6), Inches(4.4), Inches(4.5)); _fill(card, PANEL)
_txt(s, Inches(8.7), Inches(1.85), Inches(3.9), Inches(0.5), [("STARLINK-37471 剖面摘要", 15, ACCENT, True)])
stats = [("升軌總量", "+50 km"), ("微點火段數", "422 段"), ("推力弧型態", "連續電推"),
         ("時刻定位", "分鐘級"), ("計畫上傳偵測", "機動日 1,892 m")]
yy = Inches(2.5)
for k, v in stats:
    _txt(s, Inches(8.7), yy, Inches(2.2), Inches(0.5), [(k, 13, MUTE, False)], anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(10.7), yy, Inches(1.9), Inches(0.5), [(v, 15, ACCENT2, True)], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    yy += Inches(0.68)
_footer(s, 6)

# ══ 發現 ③ 福衛七星群 ════════════════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.45), "發現③　福衛七/COSMIC-2 星群觀點：三個發現")
_pic(s, "fig_r7_fs7_fleet.png", Inches(0.5), Inches(1.4), Inches(7.6))
bullets(s, Inches(8.3), Inches(1.5), Inches(4.6), [
    (ACCENT2, [("① 艦隊級同步", 15, INK, True)]),
    (MUTE, [("64 個機動 episode 中 ", 12.5, INK, False), ("95% 在 ±5 天內跨星同步", 12.5, ACCENT2, True),
            ("——整批機動作業非各自為政", 12.5, INK, False)]),
    (ACCENT2, [("② FS7-3 停止維持", 15, INK, True)]),
    (MUTE, [("2025-05 後 14 個月無機動、衰至 536 km，", 12.5, INK, False),
            ("距艦隊 41 km", 12.5, ACCENT2, True)]),
    (ACCENT2, [("③ 面幾何緩慢重排", 15, INK, True)]),
    (MUTE, [("機動關聯在", 12.5, INK, False), ("時間軸", 12.5, ACCENT2, True),
            ("、非鎖定面間距；適配 RO 任務", 12.5, INK, False)]),
    (GOOD, [("對 TASA 的意義", 14, INK, True)]),
    (MUTE, [("全部僅用公開 TLE→", 12, INK, False),
            ("可用內部真值反向驗收", 12, GOOD, True)]),
])
_footer(s, 7)

# ══ 發現 ④ LEO-PNT 天花板 ════════════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.45), "發現④　LEO-PNT：Starlink 訊號替代 GPS 的定位天花板")
_pic(s, "fig_r10_leopnt.png", Inches(0.4), Inches(1.5), Inches(7.9))
bullets(s, Inches(8.5), Inches(1.55), Inches(4.4), [
    (ACCENT, [("核心：精度＝星曆精度", 15, INK, True)]),
    (MUTE, [("Starlink 訊號不含軌道→衛星位置差多少、使用者就差多少", 12.5, INK, False)]),
    (ACCENT2, [("TLE 路線（KOC 實測）", 14.5, INK, True)]),
    (MUTE, [("新鮮 ", 12.5, INK, False), ("1.5 km", 12.5, ACCENT2, True),
            ("（沿軌主導）→72h 惡化 32 km", 12.5, INK, False)]),
    (ACCENT2, [("MEME 路線（本計畫）", 14.5, INK, True)]),
    (MUTE, [("公尺級——三個數量級之差", 12.5, INK, False)]),
    (GOOD, [("印證 KOC：可行但非商用替代", 12.5, GOOD, True)]),
])
_footer(s, 8)

# ══ 發現 ④ 水平/垂直精度提升 ═════════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.45), "發現④　TLE vs MEME：水平／垂直定位精度的提升比較")
_pic(s, "fig_r11_pnt_hv.png", Inches(0.4), Inches(1.35), Inches(7.7))
bullets(s, Inches(8.3), Inches(1.5), Inches(4.6), [
    (ACCENT, [("星曆誤差是「沿軌主導」", 15, INK, True)]),
    (MUTE, [("新鮮 TLE 三軸：沿軌 ", 12.5, INK, False), ("1,524 m", 12.5, ACCENT2, True),
            (" 遠大於徑向 143 m、法向 181 m", 12.5, INK, False)]),
    (ACCENT, [("水平（沿軌）＝定位瓶頸", 15, INK, True)]),
    (MUTE, [("沿軌誤差落在水平面→水平受害最重；MEME 壓到公尺級→", 12.5, INK, False),
            ("水平改善約 300×", 12.5, ACCENT2, True)]),
    (ACCENT, [("垂直（徑向）較能保留", 15, INK, True)]),
    (MUTE, [("徑向誤差小一個量級→高度分量較準（呼應 KOC「水平可用、整體仍 km 級」）；MEME ", 12.5, INK, False),
            ("垂直改善約 30×", 12.5, ACCENT2, True)]),
    (GOOD, [("結論", 15, INK, True)]),
    (MUTE, [("MEME 對", 12.5, INK, False), ("水平精度", 12.5, GOOD, True),
            ("的貢獻遠大於垂直；陳舊 TLE（72h）水平更達約 6,400×", 12.5, INK, False)]),
])
_footer(s, 9)

# ══ 發現 ⑤ TLE 品質不是常數 ══════════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.45), "發現⑤　TLE 品質不是常數：偵測與否取決於追蹤品質")
_pic(s, "fig_r12_tle_quality.png", Inches(0.4), Inches(1.35), Inches(7.9))
bullets(s, Inches(8.5), Inches(1.5), Inches(4.4), [
    (ACCENT2, [("修正一個常見誤解", 15, INK, True)]),
    (MUTE, [("「小機動＝TLE 看不見」是錯的——", 12.5, INK, False), ("取決於該星追蹤品質 σ", 12.5, ACCENT2, True)]),
    (ACCENT, [("σ 校準（認證安靜期實測）", 14.5, INK, True)]),
    (MUTE, [("精密測高星 σ≈", 12.5, INK, False), ("0.1–0.8 m", 12.5, ACCENT2, True),
            ("（SLR/DORIS/GPS）；Starlink 級 σ≈50 m（雷達）", 12.5, INK, False)]),
    (ACCENT, [("差約 100–250×", 14.5, INK, True)]),
    (MUTE, [("同一次 22 m 機動：Starlink 級上 0.4σ（埋雜訊）、精密星上數十σ（清楚可見）", 12.5, INK, False)]),
    (GOOD, [("意涵", 14.5, INK, True)]),
    (MUTE, [("表 11 的 σ 是「Starlink 級」非普適；MEME 的 100× 優勢專屬雷達追蹤的雜訊目標", 12, INK, False)]),
])
_footer(s, 10)

# ══ 發現 ⑤ 如何量化 TLE 品質 ═════════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.45), "發現⑤　下一個問題：如何量化一顆衛星的 TLE 品質？")
metrics = [
    ("σ_sma　半長軸雜訊", "認證/穩健安靜期去趨勢殘差 MAD。直接決定機動偵測下限（Δa_min≈4σ）。本計畫 ids_sigma_calibrate.py 已實作", ACCENT, "★ 主指標"),
    ("f＝N/T_D　更新頻率", "每日 TLE 產出數，追蹤強度代理（LEO 1–5/日、GEO ~1/日）。高頻＝追蹤密、外推新鮮（NUDT 自適應視窗即用此）", ACCENT2, "代理"),
    ("重疊 TLE 一致性", "相鄰兩 TLE 傳播至同一時刻的位置差（Lemmens 2014）。小＝一致＝品質好，不需真值", ACCENT2, "代理"),
    ("外推誤差增長率", "SGP4 位置誤差隨時間發散速率（沿軌主導，報告表 18/圖 13）。決定 LEO-PNT 定位天花板", ACCENT2, "代理"),
]
y = Inches(1.35)
for title, desc, col, tag in metrics:
    card = s.shapes.add_shape(1, Inches(0.6), y, Inches(12.1), Inches(1.02)); _fill(card, PANEL)
    bar = s.shapes.add_shape(1, Inches(0.6), y, Inches(0.14), Inches(1.02)); _fill(bar, col)
    _txt(s, Inches(0.95), y + Inches(0.10), Inches(3.6), Inches(0.85), [(title, 16, INK, True)], anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(4.6), y + Inches(0.08), Inches(6.7), Inches(0.9), [(desc, 12, MUTE, False)], anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(11.5), y + Inches(0.10), Inches(1.2), Inches(0.85), [(tag, 12, col, True)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.12)
_txt(s, Inches(0.6), Inches(6.02), Inches(12.1), Inches(1.05),
     [[("綜合 → 每顆衛星的「TLE 品質分數」→ 自適應偵測門檻", 15, GOOD, True)],
      [("取代現行固定高度帶門檻；非合作目標無 operator 日誌時，用穩健滾動 σ（各滑窗 MAD 取低百分位，機動污染窗被多數決淘汰）。列為期末方向。", 12, MUTE, False)]], sp=1.2)
_footer(s, 11)

# ══ 發現 ④ 三創見 ════════════════════════════════════════════════════════
s = slide()
_chip(s, Inches(0.6), Inches(0.45), "發現④　三項創見：本計畫對 LEO-PNT 的獨特貢獻")
cards = [
    ("以機動偵測作星曆可信度即時把關", "剛機動的衛星 TLE 失效達數十 km；burn-arc／計畫上傳偵測器可即時標記「暫不可信」供定位端剔除——定位團隊自身不會建的安全層", ACCENT),
    ("以阻力殘差改善 TLE 沿軌預測", "TLE 誤差幾乎全在沿軌；NRLMSIS 逐衛星阻力殘差正為修正沿軌漂移而生，可在不需精密星曆下把 1.5 km 壓向數百公尺", ACCENT2),
    ("TASA 的戰略縱深", "自有星座(FS-8)＋自主精密定軌＋機動監測的完整鏈條，是少數國家級單位才有的 LEO-PNT 自主能力", GOOD),
]
y = Inches(1.6)
for i, (title, desc, col) in enumerate(cards):
    card = s.shapes.add_shape(1, Inches(0.6), y, Inches(12.1), Inches(1.55)); _fill(card, PANEL)
    bar = s.shapes.add_shape(1, Inches(0.6), y, Inches(0.16), Inches(1.55)); _fill(bar, col)
    _txt(s, Inches(1.0), y + Inches(0.12), Inches(0.9), Inches(1.3),
         [(f"創見{i+1}", 15, col, True)], anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, Inches(2.2), y + Inches(0.14), Inches(10.2), Inches(0.55),
         [(title, 18, INK, True)])
    _txt(s, Inches(2.2), y + Inches(0.72), Inches(10.2), Inches(0.75),
         [(desc, 13, MUTE, False)])
    y += Inches(1.72)
_footer(s, 12)

# ══ 結語 ══════════════════════════════════════════════════════════════════
s = slide()
band = s.shapes.add_shape(1, 0, Inches(2.6), W, Inches(2.2)); _fill(band, PANEL)
_txt(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(1.0),
     [("一條主線貫穿五則發現", 32, INK, True)])
_txt(s, Inches(0.9), Inches(3.95), Inches(11.5), Inches(0.7),
     [("僅憑公開軌道資料與精密星曆，即可看見他人的操作意圖、界定自身能力邊界、並延伸出應用價值", 17, ACCENT, False)])
_txt(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(1.6),
     [[("・神龍案例展示「受控 RPO 的判讀指紋」　・MEME 剖面揭示「電推連續機動作業與計畫預寫」", 14, INK, False)],
      [("・福衛七示範「星系級操作特徵萃取」　・LEO-PNT 量化「星曆即定位天花板」", 14, INK, False)],
      [("・TLE 品質量化：偵測下限隨追蹤品質縮放（ILRS 真值校準，差約 100×）", 14, INK, False)],
      [("五者共用同一套「精密星曆＋機器學習機動偵測」能力（融合評分器 large 召回 0.973，勝過最佳統計通道 0.67）——這正是本計畫的核心資產", 13, GOOD, True)]], sp=1.35)

out = DOCS / "期中報告_附帶發現_20260716_v2.pptx"
prs.save(str(out))
print(f"saved {out}  ({out.stat().st_size//1024} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
