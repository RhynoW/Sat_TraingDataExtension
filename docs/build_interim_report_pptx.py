# -*- coding: utf-8 -*-
"""build_interim_report_pptx.py — 期中報告「正式版」簡報（自附帶發現 v2 擴充）
產出：docs/期中報告_正式簡報_20260716.pptx
結構：封面→大綱→研究背景與價值定位(四大價值)→契約目標驗收→系統架構→核心成果(六)
      →應用價值(四點對映)→附帶發現(五)→期末待辦與roadmap→建議→結語
主題色沿用 v2（深海軍藍）。
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

DOCS = Path(__file__).parent
FONT = "Microsoft JhengHei"
BG      = RGBColor(0x0E, 0x1B, 0x2E)
PANEL   = RGBColor(0x16, 0x27, 0x40)
PANEL2  = RGBColor(0x1E, 0x33, 0x52)
ACCENT  = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT2 = RGBColor(0xFF, 0xD5, 0x4F)
INK     = RGBColor(0xEC, 0xF2, 0xF9)
MUTE    = RGBColor(0x9F, 0xB3, 0xC8)
GOOD    = RGBColor(0x66, 0xBB, 0x6A)
WARN    = RGBColor(0xEF, 0x9A, 0x9A)
W, H = Inches(13.333), Inches(7.5)


def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background(); shape.shadow.inherit = False


def _bg(slide):
    r = slide.shapes.add_shape(1, 0, 0, W, H); _fill(r, BG); r.shadow.inherit = False


def _txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp
        for (text, size, color, bold) in line:
            r = p.add_run(); r.text = text
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    return tb


def _chip(slide, text, color=ACCENT):
    c = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.14), Inches(0.5)); _fill(c, color)
    _txt(slide, Inches(0.88), Inches(0.41), Inches(11.8), Inches(0.6),
         [(text, 23, INK, True)], anchor=MSO_ANCHOR.MIDDLE)


def _pic(slide, path, x, y, w):
    p = DOCS / path
    if p.exists():
        return slide.shapes.add_picture(str(p), x, y, width=w)
    _txt(slide, x, y, w, Inches(0.4), [(f"[缺圖 {path}]", 10, MUTE, False)])
    return None


def _footer(slide, n):
    _txt(slide, Inches(0.5), Inches(7.06), Inches(10), Inches(0.32),
         [("TASA-S-1150268　智慧化低軌通訊衛星軌道異常及太空事件偵測　｜　期中報告", 9.5, MUTE, False)])
    _txt(slide, Inches(12.4), Inches(7.06), Inches(0.7), Inches(0.32),
         [(str(n), 9.5, MUTE, False)], align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, items, size=14.5, gap=1.22, h=Inches(4.8)):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, (mc, runs) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = gap; p.space_after = Pt(3)
        b = p.add_run(); b.text = "▍ "; b.font.name = FONT; b.font.size = Pt(size); b.font.color.rgb = mc
        for (text, sz, color, bold) in runs:
            r = p.add_run(); r.text = text
            r.font.name = FONT; r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return tb


def cards_rows(slide, rows, y0=Inches(1.45), ch=Inches(1.02), gap=Inches(1.14), tag_w=Inches(2.7)):
    """rows: list of (tag, title, desc, color)。左標籤色條 + 標題 + 說明。"""
    y = y0
    for tag, title, desc, col in rows:
        card = slide.shapes.add_shape(1, Inches(0.6), y, Inches(12.1), ch); _fill(card, PANEL)
        bar = slide.shapes.add_shape(1, Inches(0.6), y, Inches(0.16), ch); _fill(bar, col)
        _txt(slide, Inches(0.95), y + Inches(0.08), tag_w, ch - Inches(0.16),
             [(tag, 15.5, col, True)], anchor=MSO_ANCHOR.MIDDLE)
        _txt(slide, Inches(0.95) + tag_w, y + Inches(0.06), Inches(12.5) - tag_w - Inches(0.55), ch - Inches(0.12),
             [(title, 14.5, INK, True)] if not desc else
             [[(title, 14.5, INK, True)], [(desc, 12, MUTE, False)]], anchor=MSO_ANCHOR.MIDDLE, sp=1.08)
        y += gap


prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
def slide():
    s = prs.slides.add_slide(BLANK); _bg(s); return s


def divider(title, sub):
    s = slide()
    band = s.shapes.add_shape(1, 0, Inches(2.7), W, Inches(2.0)); _fill(band, PANEL)
    bar = s.shapes.add_shape(1, Inches(0.9), Inches(2.9), Inches(0.18), Inches(1.6)); _fill(bar, ACCENT)
    _txt(s, Inches(1.3), Inches(3.0), Inches(11), Inches(1.0), [(title, 34, INK, True)])
    _txt(s, Inches(1.3), Inches(4.05), Inches(11), Inches(0.6), [(sub, 16, ACCENT, False)])
    return s


# ══ 1 封面 ══════════════════════════════════════════════════════════════════
s = slide()
band = s.shapes.add_shape(1, 0, Inches(2.3), W, Inches(2.7)); _fill(band, PANEL)
_txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.6),
     [("期　中　報　告", 22, ACCENT, True)])
_txt(s, Inches(0.9), Inches(2.55), Inches(11.6), Inches(1.3),
     [("智慧化低軌通訊衛星軌道異常", 38, INK, True)])
_txt(s, Inches(0.9), Inches(3.5), Inches(11.6), Inches(0.9),
     [("及太空事件偵測", 38, INK, True)])
_txt(s, Inches(0.9), Inches(5.15), Inches(11.6), Inches(0.6),
     [("以 TLE 公開星曆＋MEME 精密真值，融合物理模型與人工智慧之可解釋偵測架構", 16, ACCENT, False)])
_txt(s, Inches(0.9), Inches(6.55), Inches(11.6), Inches(0.5),
     [("計畫案號 TASA-S-1150268　｜　社團法人中華民國國防科技學術研究學會　｜　2026-07-16", 12, MUTE, False)])

# ══ 2 大綱 ══════════════════════════════════════════════════════════════════
s = slide(); _chip(s, "簡報大綱")
cards_rows(s, [
    ("一", "研究背景與價值定位", "低軌通訊衛星爆發性部署下的營運安全與碰撞風險；本計畫四大價值", ACCENT),
    ("二", "契約目標與驗收", "契約表 8 共 14 項指標：12 項直接達標、1 項調整後達標、1 項待複核", ACCENT2),
    ("三", "系統架構", "規則—統計—機器學習三層＋物理模型交叉驗證（可解釋 AI）", ACCENT),
    ("四", "核心成果", "ML 主力偵測、機動時刻、星系級異常、偵測下限、可解釋 AI 改進、全庫驗證", GOOD),
    ("五", "應用價值", "四大應用價值對映實測成果", ACCENT2),
    ("六", "附帶發現", "五則從公開軌道資料看見的操作意圖", ACCENT2),
    ("七", "期末待辦、roadmap 與建議", "ML 解鎖路徑、精密星曆、延伸研究基礎", ACCENT),
], y0=Inches(1.25), ch=Inches(0.70), gap=Inches(0.79), tag_w=Inches(1.0))
_footer(s, 2)

# ══ 一、研究背景 ══════════════════════════════════════════════════════════
divider("一、研究背景與價值定位", "為何要做低軌衛星異常偵測")
_footer(prs.slides[-1], 3)

s = slide(); _chip(s, "研究背景：低軌爆發性部署下的安全需求")
bullets(s, Inches(0.7), Inches(1.45), Inches(12.0), [
    (ACCENT2, [("低軌通訊衛星進入巨型星系時代", 17, INK, True)]),
    (MUTE, [("Starlink 逾 1 萬顆、OneWeb／Kuiper／千帆持續部署；在軌物件逾 3.3 萬。", 14, INK, False)]),
    (ACCENT2, [("軌道擁擠 → 碰撞與異常風險升高", 17, INK, True)]),
    (MUTE, [("非預期軌道偏移、異常推力事件、任務失效若無法即時辨識，將直接威脅營運安全與太空環境。", 14, INK, False)]),
    (ACCENT2, [("公開資料的挑戰與機會", 17, INK, True)]),
    (MUTE, [("TLE 公開但為公里級、雜訊大；MEME 精密星曆公尺級但僅 Starlink。", 14, INK, False),
            ("本計畫以 MEME 為真值訓練，讓「只用公開 TLE」也能高精度偵測。", 14, ACCENT, True)]),
    (GOOD, [("目標", 17, INK, True)]),
    (MUTE, [("建立可對單一衛星與整個星系運作的機動／異常偵測分析環境，並誠實界定自身能力邊界。", 14, INK, False)]),
])
_footer(s, 4)

# ══ 四大價值定位（核心） ═══════════════════════════════════════════════════
s = slide(); _chip(s, "本計畫四大價值定位")
cards_rows(s, [
    ("① 營運安全", "支援低軌通訊衛星營運安全", "機動偵測＋時刻定位＋星系級異常＋再入守門，即時掌握衛星是否正常運作", ACCENT),
    ("② 風險管理", "碰撞風險升高下的異常辨識模型", "非預期軌道偏移監測、異常推力事件判識、任務失效早期診斷 → 提升穩定性與風險管理", ACCENT2),
    ("③ AI 示範", "AI 於太空監控領域之示範應用", "監督式 LightGBM＋五通道融合（AUC 0.982），全庫 1 萬星實測驗證", GOOD),
    ("④ 可解釋 AI", "融合物理＋AI 之可解釋偵測架構", "作為太空碎片監測、衛星接近行為分析、軌道預測不確定性評估之延伸研究基礎", ACCENT),
], y0=Inches(1.4), ch=Inches(1.16), gap=Inches(1.3), tag_w=Inches(2.5))
_footer(s, 5)

# ══ 二、契約目標與驗收 ═════════════════════════════════════════════════════
divider("二、契約目標與驗收", "契約表 8：14 項指標，12 直接達標 ＋ 1 調整後 ＋ 1 待複核")
_footer(prs.slides[-1], 6)

s = slide(); _chip(s, "契約表 8 驗收總覽：12 直接 ＋ 1 調整後 ＋ 1 待複核")
# 左：達標統計卡
card = s.shapes.add_shape(1, Inches(0.6), Inches(1.5), Inches(4.5), Inches(4.7)); _fill(card, PANEL)
_txt(s, Inches(0.9), Inches(1.75), Inches(4.0), Inches(0.6), [("驗收指標現況", 17, ACCENT, True)])
_txt(s, Inches(0.9), Inches(2.5), Inches(4.0), Inches(1.4),
     [[("12 ＋ 1 ＋ 1", 40, GOOD, True)], [("直接達標／調整後／待複核", 13, INK, False)]], sp=1.0)
_txt(s, Inches(0.9), Inches(4.25), Inches(4.0), Inches(1.9),
     [[("情境① TLE 單衛星：5 直接＋1 待對齊", 13, INK, False)],
      [("情境② MEME 星系級：7 直接＋1 待複核", 13, INK, False)],
      [("＋三項星系級分析全建置", 13, INK, False)]], sp=1.5)
bullets(s, Inches(5.4), Inches(1.5), Inches(7.3), [
    (ACCENT2, [("情境①#2：調整後達標，惟數值來源待對齊", 15, INK, True)]),
    (MUTE, [("原引 0.397→", 12.5, INK, False), ("0.973", 13, ACCENT2, True),
            (" 係融合評分器於星系級 unit 集之成績，非 Model 1 外部 TPR；期末補量以對齊定義", 12.5, INK, False)]),
    (ACCENT2, [("情境②#7：達標證據已撤回，改列待複核", 15, INK, True)]),
    (MUTE, [("原引「episode 級 small 召回 0.965」經 ", 12.5, INK, False),
            ("naive 隨機對照證實不具鑑別力", 13, WARN, True),
            ("（隨機分數得 0.999）；期末改以窗級指標於契約集複核", 12.5, INK, False)]),
    (GOOD, [("核心方法論（誠實透明）", 15, INK, True)]),
    (MUTE, [("原始實測值（0.397／0.000）全數揭露；", 12.5, INK, False),
            ("寧可誠實列待複核，不以不具鑑別力之指標宣稱達標。", 12.5, GOOD, True)]),
])
_footer(s, 7)

# ══ 三、系統架構 ═══════════════════════════════════════════════════════════
divider("三、系統架構", "三層偵測＋物理交叉驗證＝可解釋 AI")
_footer(prs.slides[-1], 8)

s = slide(); _chip(s, "系統架構：規則—統計—機器學習＋物理")
cards_rows(s, [
    ("Layer 1 規則", "白箱規則 P1–P6", "人看得懂、每個旗標附理由；可解釋、可稽核（14,090 顆全量驗證）", ACCENT),
    ("Layer 2 統計", "CUSUM／BOCPD／SSA／3σ-MAD", "不預設機動樣貌，逐點變化點偵測；高召回候選產生器", ACCENT2),
    ("Layer 3 機器學習", "LightGBM Model 1 ＋ 五通道融合", "本專案偵測主力；綜合多通道給出可操作判別（AUC 0.982）", GOOD),
    ("物理交叉驗證", "NRLMSIS 阻力殘差（Model 2）", "把大氣阻力從嫌疑名單排除，跨軌域通用；再入守門", ACCENT),
    ("自動路由", "依軌域/域自動選路", "域外目標走物理路徑，避免監督式模型 OOD 失效", ACCENT2),
], y0=Inches(1.4), ch=Inches(0.92), gap=Inches(1.04), tag_w=Inches(2.9))
_txt(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.5),
     [("可解釋性：每個警報都能溯源到「哪一層、哪條規則、哪個物理量」——這正是可解釋 AI 偵測架構（價值④）", 12.5, GOOD, True)])
_footer(s, 9)

# ══ 四、核心成果 ═══════════════════════════════════════════════════════════
divider("四、核心成果", "六項實測成果")
_footer(prs.slides[-1], 10)

# 指標判讀速查（看數字前先讀這頁）
s = slide(); _chip(s, "先讀這頁：指標的物理意義與好壞方向")
cards_rows(s, [
    ("召回率 Recall ↑", "100 次真的有機動，抓到幾次——「抓得全不全」",
     "越大越好（1.0＝一次都沒漏）；本計畫融合 large 召回 0.973", GOOD),
    ("精確率 Precision ↑", "發 100 次警報，幾次是真的——「報得準不準」",
     "越大越好（1.0＝零誤報）；本計畫融合星系級精確率高", GOOD),
    ("ROC-AUC ↑", "隨機抽一真一假，給真者較高分的機率——「排序能力」",
     "越大越好（0.5＝亂猜、1.0＝完美）；本計畫融合 0.982", GOOD),
    ("FPR 誤報率 ↓", "100 個沒機動的時段，誤標幾個——「亂叫的程度」",
     "越小越好（0＝從不誤報）；本計畫控制在 ≤0.05 預算內（0.0498）", ACCENT),
    ("延遲 / 時刻誤差 ↓", "多快知道有事（延遲）、點火時刻估得多準（誤差）",
     "越小越好；延遲中位 0.1h、時刻誤差中位 2.74h", ACCENT),
    ("large / medium / small", "機動的大小分級（依半長軸變化 Δa 量級）",
     "large＝明顯變軌(好抓)、medium＝一般維持、small＝微動(貼近雜訊底、三層皆難，屬資料物理限制)", ACCENT2),
], y0=Inches(1.3), ch=Inches(0.8), gap=Inches(0.88), tag_w=Inches(3.0))
_txt(s, Inches(0.6), Inches(6.72), Inches(12.1), Inches(0.4),
     [("一句話：召回／精確／AUC 越大越好；FPR／延遲越小越好。理想偵測器＝抓得全、報得準、反應快、不亂叫，且連 small 都抓得到。",
       12, GOOD, True)])
_footer(s, 11)

# 成果① ML 主力
s = slide(); _chip(s, "成果①　機器學習為偵測主力（AI 示範・價值③）")
_pic(s, "../Orbital_Maneuver_V2/output/roc_comparison.png", Inches(0.5), Inches(1.5), Inches(6.7))
bullets(s, Inches(7.5), Inches(1.5), Inches(5.4), [
    (ACCENT2, [("五通道融合評分器（梯度提升樹）", 15, INK, True)]),
    (MUTE, [("ROC-AUC ", 13, INK, False), ("0.982", 14, GOOD, True),
            ("、large 召回 0.973、FPR 0.0498", 13, INK, False)]),
    (ACCENT2, [("勝過任何單一統計通道", 15, INK, True)]),
    (MUTE, [("統計層最佳 3σ-MAD 召回 0.67；ML 才給出可操作的機動判別", 13, INK, False)]),
    (ACCENT2, [("單衛星 LightGBM（Layer 3）", 15, INK, True)]),
    (MUTE, [("內部交叉驗證 TPR 97.5%、AUC 0.996；", 13, INK, False),
            ("全庫外部旗標率 21.6%（域外落差，詳成果⑥）", 12.5, MUTE, False)]),
    (GOOD, [("ML 價值已交付；成果⑤ 進一步改善其決策依據", 13.5, GOOD, True)]),
])
_footer(s, 12)

# 成果② 機動時刻
s = slide(); _chip(s, "成果②　機動偵測與時刻定位（8 小時分辨率達標）")
_pic(s, "fig_r3_joint_element.png", Inches(0.5), Inches(1.5), Inches(7.6))
bullets(s, Inches(8.3), Inches(1.5), Inches(4.6), [
    (ACCENT2, [("機動時刻（burn epoch）偵測", 15, INK, True)]),
    (MUTE, [("中位誤差 2.74 小時、86.9% 落在 8 小時內 → 達契約分辨率", 13, INK, False)]),
    (ACCENT2, [("警報延遲中位 0.1 小時", 15, INK, True)]),
    (MUTE, [("半數事件約 6 分鐘內亮旗，遠優於契約 24 小時", 13, INK, False)]),
    (ACCENT2, [("元素／增量聯合圖", 15, INK, True)]),
    (MUTE, [("長期趨勢＋單次事件突刺同框，與 MEME 真值時刻對齊", 13, INK, False)]),
    (GOOD, [("把「某天動過」收窄到「當天哪個時段動的」", 13, GOOD, True)]),
])
_footer(s, 13)

# 成果③ 星系級（碰撞/風險）
s = slide(); _chip(s, "成果③　星系級異常分析（碰撞風險管理・價值②）")
_pic(s, "fig_r2_starlink_planes.png", Inches(0.5), Inches(1.5), Inches(7.4))
bullets(s, Inches(8.1), Inches(1.5), Inches(4.8), [
    (ACCENT2, [("契約三項星系級分析全建置", 15, INK, True)]),
    (MUTE, [("軌道面一致性／批量機動識別／陣型誤差", 13, INK, False)]),
    (ACCENT2, [("該響時響、不該響時安靜", 15, INK, True)]),
    (MUTE, [("Starlink 偵得 7 個異常軌道面；OneWeb 654／千帆 238 正確回報零異常", 13, INK, False)]),
    (GOOD, [("對接價值②：異常辨識模型", 14, INK, True)]),
    (MUTE, [("非預期軌道偏移監測、批量部署／重組識別、相位保持失效偵測", 13, INK, False)]),
])
_footer(s, 14)

# 成果④ 偵測下限
s = slide(); _chip(s, "成果④　偵測下限量化與精密星曆價值")
bullets(s, Inches(0.7), Inches(1.55), Inches(12.0), [
    (ACCENT2, [("量化「多小的機動仍可偵測」", 16, INK, True)]),
    (MUTE, [("TLE 偵測下限 ΔV ≈ 0.1–0.57 m/s（合成注入實測）；", 14, INK, False),
            ("MEME 理論下限 ≈ 0.011 m/s，靈敏約 10–50 倍。", 14, ACCENT2, True)]),
    (ACCENT2, [("獨立外部驗證（ILRS／IDS 第二真值集）", 16, INK, True)]),
    (MUTE, [("14 顆測高衛星 1,651 次 operator 點火；其例行機動中位 0.0099／0.0115 m/s ", 14, INK, False),
            ("位於 MEME 推導門檻附近，屬邊際可偵測", 14, GOOD, True)]),
    (ACCENT2, [("TLE 品質不是常數（σ 校準）", 16, INK, True)]),
    (MUTE, [("精密測高星實測 σ≈0.3 m、Starlink 級 ≈50 m，差約 100 倍 → 偵測下限隨追蹤品質縮放", 14, INK, False)]),
    (GOOD, [("政策含義：對重點目標爭取精密星曆，是解鎖微型（電推）機動偵測的唯一途徑", 13.5, GOOD, True)]),
])
_footer(s, 15)

# 成果⑤ 可解釋 AI 改進（18.3）
s = slide(); _chip(s, "成果⑤　可解釋 AI 改進實證：域不變模型（價值④）")
_pic(s, "fig_r12_tle_quality.png", Inches(0.4), Inches(1.5), Inches(6.6))
bullets(s, Inches(7.2), Inches(1.5), Inches(5.6), [
    (ACCENT2, [("① 改用 MEME 真值原生標籤", 14.5, INK, True)]),
    (MUTE, [("破「自標籤循環」，改由外部精密星曆給答案；配合 episode 級評估，外部召回 0.397→0.973（詳 p.7）", 12.5, INK, False)]),
    (ACCENT2, [("② 域不變特徵：拿掉「身分」、只看「行為」（本頁主結果）", 14.5, INK, True)]),
    (MUTE, [("移除「這是哪一殼」的身分特徵（高度/傾角/離心率），Δa 改成「自身雜訊的幾倍(SNR)」；特徵 40→36，", 12.5, INK, False),
            ("效能幾乎不變", 13, GOOD, True), ("（AUC 0.554→0.560、large 召回 0.182→0.165）＝消除身分零代價", 12.5, INK, False)]),
    (ACCENT2, [("最硬證據——反事實驗證：行為不動、只換身分證，看分數變多少", 14.5, INK, True)]),
    (MUTE, [("舊模型：光換身分、行為沒變，分數最多跳 0.79（＝它在看身分，非行為）；", 12.5, INK, False),
            ("新模型 Δp＝0.000", 13, GOOD, True), ("（只看行為，域先驗〔靠軌道身分而非機動行為給分的偏誤〕結構性歸零）", 12.5, INK, False)]),
    (GOOD, [("結論：窗級效能持平下結構性消除域先驗——靠真值標籤＋SNR，非加大模型", 12.5, GOOD, True)]),
    (MUTE, [("註：效能一律以窗級指標為據；episode 級召回因單顆多窗任一命中即算抓到、", 11.5, MUTE, False),
            ("naive 隨機亦得 1.000 不具鑑別力，故不採為效能指標", 11.5, MUTE, False)]),
])
_footer(s, 16)

# 成果⑥ 全庫驗證
s = slide(); _chip(s, "成果⑥　全庫 1 萬星比較：Model 1 vs Model 2（互補而非替代）")
bullets(s, Inches(0.7), Inches(1.42), Inches(12.0), [
    (MUTE, [("隨機 1 萬顆酬載、三方法並跑（可重現）：統計層 99.4%（候選產生器）｜Model 1 21.6%（保守）｜Model 2 41.7%（跨域）", 13.5, INK, False)]),
])
cards_rows(s, [
    ("模型", "Model 1：監督式 LightGBM", "Model 2：無監督 Isolation Forest", ACCENT),
    ("訓練／輸入", "Model 1：Starlink＋MEME 真值、絕對特徵（含域先驗）", "Model 2：免真值、物理殘差 z（阻力／Δi／Δe／ΔRAAN，扣自然演化）", ACCENT2),
    ("適用／弱點", "Model 1：Starlink 域內精準；出域（衰減軌）誤報", "Model 2：任何軌域通用；較敏感、旗標率偏高", ACCENT),
    ("全庫旗標率", "Model 1：~21%（保守）", "Model 2：~41%（跨域）", GOOD),
], y0=Inches(2.15), ch=Inches(0.86), gap=Inches(0.98), tag_w=Inches(2.2))
bullets(s, Inches(0.7), Inches(6.25), Inches(12.0), [
    (GOOD, [("三輪抽樣 κ≈0.40 互補一致；有 MEME 真值走 Model 1、其餘走 Model 2＝軌域路由依據", 13, GOOD, True),
            ("（「僅 Model 1 報」高度中位恆為 483 km＝域先驗指紋三輪重現）", 11.5, MUTE, False)]),
])
_footer(s, 17)

# 成果⑦ L3 泛化穩定度（§13.4）——過擬合反駁
s = slide(); _chip(s, "成果⑦　L3 泛化穩定度：增益非過擬合特定分布")
bullets(s, Inches(0.7), Inches(1.32), Inches(12.0), [
    (ACCENT2, [("問題：L3（AUC 0.982）明顯優於 L1／L2，是否只是過擬合特定分布？", 15, INK, True)]),
    (MUTE, [("做法：同一擂台沿四軸切片；族群/高度/品質用無洩漏 OOF、時間做真正 out-of-time holdout（前 60% 訓練、後 40% 測）", 12.5, INK, False)]),
])
cards_rows(s, [
    ("族群／傾角 shell", "53.2° / 70° / 97.6° SSO", "AUC 0.981 / 0.970 / 0.984——跨族群穩定", ACCENT),
    ("軌域高度帶", "低 / 中 / 高 三帶", "AUC 0.989 / 0.968 / 0.983——跨高度穩定", ACCENT2),
    ("資料品質（σ、更新頻率）", "佳 / 中 / 差 三級", "AUC 0.978–0.986——對域內品質變異不敏感", ACCENT),
    ("時間 out-of-time", "只見前 60%、預測後 40%", "AUC 0.945（最強抗過擬合：不可能記住未見時間段）", GOOD),
], y0=Inches(2.75), ch=Inches(0.92), gap=Inches(1.02), tag_w=Inches(3.4))
bullets(s, Inches(0.7), Inches(6.5), Inches(12.0), [
    (GOOD, [("四軸全部切片 AUC 全距僅 0.044（0.945–0.989）", 15, GOOD, True),
            ("——無單一切片撐盤，證明 L3 是普遍判別力、非過擬合", 13, INK, False)]),
])
_footer(s, 18)

# ══ 五、應用價值（四點對映） ═══════════════════════════════════════════════
divider("五、應用價值", "四大價值對映實測成果")
_footer(prs.slides[-1], 19)

s = slide(); _chip(s, "價值①②　營運安全與碰撞風險管理")
cards_rows(s, [
    ("非預期軌道偏移監測", "阻力殘差（Model 2）＋星系級陣型誤差", "扣除大氣阻力後的異常 Δa；相位保持失效偵測", ACCENT),
    ("異常推力事件判識", "機動偵測＋burn-arc 剖面器＋計畫上傳偵測", "分鐘級推力弧定位；機動「剛編入星曆」即預警", ACCENT2),
    ("任務失效早期診斷", "衰減軌辨識＋再入守門＋停止維持偵測", "FS7-3 停止維持案例：14 個月無機動、衰至 536 km 可判", GOOD),
    ("接近／碰撞風險", "交會 TCA／Pc＋RPO 事件重建", "受控伴飛 vs 碎片分離之運動學指紋辨識", ACCENT),
], y0=Inches(1.45), ch=Inches(1.12), gap=Inches(1.26), tag_w=Inches(3.4))
_footer(s, 20)

s = slide(); _chip(s, "價值③④　AI 示範與可解釋 AI 之延伸研究基礎")
bullets(s, Inches(0.7), Inches(1.5), Inches(12.0), [
    (GOOD, [("價值③　AI 於太空監控之示範應用", 16, INK, True)]),
    (MUTE, [("監督式 LightGBM＋五通道融合（AUC 0.982）＋全庫 1 萬星實測，", 14, INK, False),
            ("是 AI 在 SDA 機動偵測落地的完整示範。", 14, ACCENT, True)]),
    (GOOD, [("價值④　融合物理＋AI 之可解釋偵測架構", 16, INK, True)]),
    (MUTE, [("三層架構每個警報可溯源；域不變模型（成果⑤）依「訊噪比」而非「域記憶」決策（反事實 Δp=0）。", 14, INK, False)]),
    (ACCENT2, [("可延伸之研究基礎（本架構為地基）", 15.5, INK, True)]),
    (MUTE, [("・太空碎片監測：星系級異常＋阻力殘差跨軌域通用", 13.5, INK, False)]),
    (MUTE, [("・衛星接近行為分析：RPO 重建＋交會 TCA/Pc（第十五節已實測）", 13.5, INK, False)]),
    (MUTE, [("・軌道預測不確定性評估：TLE vs MEME 誤差預算＋σ 品質校準（LEO-PNT 章）", 13.5, INK, False)]),
])
_footer(s, 21)

# ══ 六、附帶發現（自 v2） ══════════════════════════════════════════════════
divider("六、附帶發現", "五則從公開軌道資料看見的操作意圖")
_footer(prs.slides[-1], 22)

s = slide(); _chip(s, "附帶發現五則速覽")
cards_rows(s, [
    ("①", "神龍 RPO 高度交換", "58573→物體 G：約 1 km 量級高度差與對稱交換，加四則國際同型案例", ACCENT),
    ("②", "MEME 預寫推力弧", "Starlink 把計畫機動預先編入精密星曆；電推升軌為連續機動作業（422 微點火）", ACCENT2),
    ("③", "福衛七星群觀點", "六星 95% 機動為艦隊同步作業、FS7-3 停止維持、面幾何緩慢重排", GOOD),
    ("④", "LEO-PNT 定位天花板", "2,600 萬點實測 TLE(km 級) vs MEME(公尺級)；水平/垂直精度比較", ACCENT),
    ("⑤", "TLE 品質量化", "追蹤品質決定 σ（差約 100×）；如何量化每顆星的 TLE 品質", ACCENT2),
], y0=Inches(1.3), ch=Inches(0.9), gap=Inches(1.02), tag_w=Inches(0.9))
_txt(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
     [("共同點：全部僅用公開 TLE 重建 → 展示「從公開資料看見操作意圖」的 SDA 能力（詳附帶發現彙整簡報）", 12, GOOD, True)])
_footer(s, 23)

# 附帶發現：LEO-PNT 水平/垂直（放一張代表圖）
s = slide(); _chip(s, "附帶發現④　TLE vs MEME：水平／垂直定位精度")
_pic(s, "fig_r11_pnt_hv.png", Inches(0.5), Inches(1.4), Inches(7.7))
bullets(s, Inches(8.4), Inches(1.5), Inches(4.5), [
    (ACCENT, [("星曆誤差沿軌主導", 15, INK, True)]),
    (MUTE, [("新鮮 TLE 沿軌 1,524 m 遠大於徑向 143 m", 13, INK, False)]),
    (ACCENT, [("MEME 改善以水平最大", 15, INK, True)]),
    (MUTE, [("水平約 300×、垂直約 30×；陳舊 72h 更達約 6,400×", 13, INK, False)]),
    (GOOD, [("對應價值④：軌道預測不確定性評估之延伸基礎", 13, GOOD, True)]),
])
_footer(s, 24)

# ══ 七、期末待辦與 roadmap ═════════════════════════════════════════════════
divider("七、期末待辦、roadmap 與建議", "ML 解鎖路徑與延伸研究")
_footer(prs.slides[-1], 25)

s = slide(); _chip(s, "期末 ML roadmap：深度模型的解鎖路徑")
bullets(s, Inches(0.7), Inches(1.5), Inches(12.0), [
    (ACCENT2, [("現況：ML 主力（融合）已勝出，深度序列模型受限於資料量", 15.5, INK, True)]),
    (MUTE, [("非能力問題——僅 597 個正 episode，稀疏正樣本下聚合特徵已足夠。", 13.5, INK, False)]),
    (ACCENT, [("自監督預訓練（最高潛力）", 15, INK, True)]),
    (MUTE, [("數千萬筆未標註 TLE 遮罩重建 → 小真值集微調，繞過正樣本太少", 13.5, INK, False)]),
    (ACCENT, [("真值擴增已就緒", 15, INK, True)]),
    (MUTE, [("Plan A（MEME）47,349 窗＋IDS 1,651 點火，正樣本層加厚", 13.5, INK, False)]),
    (ACCENT, [("域不變表徵為地基＋電推微結構為天然戰場", 15, INK, True)]),
    (MUTE, [("以成果⑤之 σ-SNR 表徵當輸入；422 段每軌微點火之序列訊號正是深度模型試驗場", 13.5, INK, False)]),
])
_footer(s, 26)

s = slide(); _chip(s, "建議事項")
cards_rows(s, [
    ("重點目標精密星曆", "對非合作＋高機動頻率目標爭取 MEME 級資料", "解鎖微型（電推 <0.1 m/s）機動偵測的唯一途徑", ACCENT),
    ("釋出 FORMOSAT GNSS（自有資產）", "FS-7 等自帶精密定軌，釋出即可訓練 L3", "ML 專屬有精密星曆之目標；σ 降至公尺級則本土衛星小機動亦可偵測（需遷移學習補正樣本）", GOOD),
    ("營運試行", "統一偵測摘要卡＋星系級週報 4–8 週試行", "收集誤報成本與判讀時效，校調操作點", ACCENT2),
    ("非合作星系常態監測", "以 Model 2＋NRLMSIS＋星系級三分析為主軸", "無需標籤、跨軌域通用；千帆基線已就緒", GOOD),
    ("延伸研究基礎", "碎片監測／接近行為／軌道預測不確定性", "以本可解釋 AI 架構為地基向外擴展", ACCENT),
], y0=Inches(1.3), ch=Inches(0.94), gap=Inches(1.05), tag_w=Inches(3.6))
_footer(s, 27)

# 送審重點：真值定義・驗收主指標・不適用場景（回應「一頁濃縮」）
s = slide(); _chip(s, "送審重點：真值定義・驗收主指標・不適用場景")
bullets(s, Inches(0.7), Inches(1.3), Inches(12.0), [
    (ACCENT2, [("真值定義：效能宣稱只採「嚴格真值」", 15, INK, True)]),
    (MUTE, [("嚴格真值＝MEME 精密星曆（公尺級）＋ILRS/IDS 點火日誌；代理真值（有無推進）僅描述覆蓋規模、不作效能宣稱；TLE 為特徵非真值", 12.5, INK, False)]),
    (ACCENT2, [("驗收主指標：唯一，避免敘事分散", 15, INK, True)]),
    (MUTE, [("＝unit 級 ROC-AUC ＋ FPR≤0.05 操作點下召回；satellite/episode/large-medium-small 皆為研究分析指標，不作驗收主張", 12.5, INK, False)]),
])
cards_rows(s, [
    ("small 機動（SNR<2）", "★核心限制", "貼近雜訊 σ 之物理下限；標「下限以下」，需精密星曆", ACCENT),
    ("極低更新／長缺口", "時效受限", "事件落取樣間隙；quality_flag 降權、不強判", ACCENT2),
    ("無嚴格真值目標", "改走物理路徑", "L3 無標籤可訓；Model 2＋NRLMSIS，不作量化召回", GOOD),
    ("域外量化召回外推", "僅域內成立", "L3 數字限 Starlink LEO；跨域僅作 FPR/一致性驗證", ACCENT),
], y0=Inches(3.55), ch=Inches(0.82), gap=Inches(0.92), tag_w=Inches(3.2))
bullets(s, Inches(0.7), Inches(6.55), Inches(12.0), [
    (GOOD, [("主流程：TLE→L1 廣掃(alarm)→L2 確認變點(score)→L3 融合(score+rank+alarm, FPR≤0.05)→路由守門(域外/再入)→摘要卡人工覆核", 12.5, INK, True)]),
])
_footer(s, 28)

# ══ 結語 ══════════════════════════════════════════════════════════════════
s = slide()
band = s.shapes.add_shape(1, 0, Inches(2.4), W, Inches(2.5)); _fill(band, PANEL)
_txt(s, Inches(0.9), Inches(2.65), Inches(11.6), Inches(0.9), [("一套可解釋的 AI 偵測架構，四個落地價值", 30, INK, True)])
_txt(s, Inches(0.9), Inches(3.7), Inches(11.6), Inches(0.7),
     [("僅憑公開軌道資料＋精密星曆真值，融合物理模型與人工智慧，支援低軌通訊衛星營運安全", 16, ACCENT, False)])
_txt(s, Inches(0.9), Inches(5.15), Inches(11.6), Inches(1.7),
     [[("・①營運安全　・②碰撞風險下的異常辨識　・③AI 太空監控示範　・④可解釋 AI 之延伸研究基礎", 14, INK, False)],
      [("・ML 為偵測主力：融合評分器 large 召回 0.973，勝過最佳統計通道 0.67——兌現「AI 抓到統計找不到的資訊」", 14, GOOD, True)],
      [("・核心資產：精密星曆＋機器學習機動偵測＋物理交叉驗證的可解釋能力", 14, INK, False)]], sp=1.4)
_footer(s, 29)

out = DOCS / "期中報告_正式簡報_20260716.pptx"
prs.save(str(out))
print(f"saved {out}  ({out.stat().st_size//1024} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
