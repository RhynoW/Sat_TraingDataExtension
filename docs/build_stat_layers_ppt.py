# -*- coding: utf-8 -*-
"""build_stat_layers_ppt.py — 統計分析基石：Layer 1 規則層與 Layer 2 統計層
產出：docs/統計分析基石_L1L2_20260720.pptx
結構：封面→大綱→定位→L1 方法/程序/成果→L2 方法(三法)/程序/成果→能力邊界→誠實聲明→結語
主題色沿用期中報告正式簡報（深海軍藍）。

所有數值來源（可重跑驗證）：
  data/benchmark/benchmark_v1_20260715.csv     三層總表（含 eval_basis）
  data/statistical_layer/metrics_20260714.csv  L2 八組 P/R/lead
  data/benchmark/injection_detect_rate.csv     合成注入偵測率 vs SNR
  data/benchmark/fusion_metrics_20260715.csv   融合評分器指標
  maneuver_strategies_july.py                  P1–P6 實作與預設參數
  statistical_detectors.py                     CUSUM/BOCPD/SSA/MAD 實作
  docs/paper1_tle_maneuver_detection_zh.md     L1 消融與 hold-out 驗證
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

DOCS = Path(__file__).parent
FONT = "Microsoft JhengHei"
BG      = RGBColor(0x0E, 0x1B, 0x2E)
PANEL   = RGBColor(0x16, 0x27, 0x40)
PANEL2  = RGBColor(0x1E, 0x33, 0x52)
ACCENT  = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT2 = RGBColor(0xFF, 0xD5, 0x4F)
INK     = RGBColor(0xEC, 0xF2, 0xF9)
MUTE    = RGBColor(0x9F, 0xB3, 0xC8)
GOOD    = RGBColor(0x66, 0xBB, 0x6A)
WARN    = RGBColor(0xEF, 0x9A, 0x9A)
W, H = Inches(13.333), Inches(7.5)


def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background(); shape.shadow.inherit = False


def _bg(slide):
    r = slide.shapes.add_shape(1, 0, 0, W, H); _fill(r, BG); r.shadow.inherit = False


def _txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp
        for (text, size, color, bold) in line:
            r = p.add_run(); r.text = text
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    return tb


def _chip(slide, text, color=ACCENT):
    c = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(0.14), Inches(0.5)); _fill(c, color)
    _txt(slide, Inches(0.88), Inches(0.41), Inches(11.8), Inches(0.6),
         [(text, 23, INK, True)], anchor=MSO_ANCHOR.MIDDLE)


def _footer(slide, n):
    _txt(slide, Inches(0.5), Inches(7.06), Inches(10), Inches(0.32),
         [("TASA-S-1150268　統計分析基石：Layer 1 規則層與 Layer 2 統計層", 9.5, MUTE, False)])
    _txt(slide, Inches(12.4), Inches(7.06), Inches(0.7), Inches(0.32),
         [(str(n), 9.5, MUTE, False)], align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, items, size=14.5, gap=1.22, h=Inches(4.9)):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, (mc, runs) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = gap; p.space_after = Pt(3)
        b = p.add_run(); b.text = "▍ "; b.font.name = FONT; b.font.size = Pt(size); b.font.color.rgb = mc
        for (text, sz, color, bold) in runs:
            r = p.add_run(); r.text = text
            r.font.name = FONT; r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    return tb


def cards_rows(slide, rows, y0=Inches(1.45), ch=Inches(1.02), gap=Inches(1.14), tag_w=Inches(2.7)):
    y = y0
    for tag, title, desc, col in rows:
        card = slide.shapes.add_shape(1, Inches(0.6), y, Inches(12.1), ch); _fill(card, PANEL)
        bar = slide.shapes.add_shape(1, Inches(0.6), y, Inches(0.16), ch); _fill(bar, col)
        _txt(slide, Inches(0.95), y + Inches(0.08), tag_w, ch - Inches(0.16),
             [(tag, 15.5, col, True)], anchor=MSO_ANCHOR.MIDDLE)
        _txt(slide, Inches(0.95) + tag_w, y + Inches(0.06), Inches(12.5) - tag_w - Inches(0.55), ch - Inches(0.12),
             [(title, 14.5, INK, True)] if not desc else
             [[(title, 14.5, INK, True)], [(desc, 12, MUTE, False)]], anchor=MSO_ANCHOR.MIDDLE, sp=1.08)
        y += gap


def table(slide, x, y, cols, rows, cw, size=11.5, rh=Inches(0.34), head=Inches(0.38),
          hi_rows=(), hi_col=GOOD):
    """簡易表格：cols=標題list，rows=list of list，cw=各欄寬(Inches)list。"""
    hx = x
    for c, w_ in zip(cols, cw):
        _txt(slide, hx, y, w_, head, [(c, size, ACCENT, True)], anchor=MSO_ANCHOR.MIDDLE)
        hx += w_
    ln = slide.shapes.add_shape(1, x, y + head, sum(cw), Inches(0.02)); _fill(ln, PANEL2)
    yy = y + head + Inches(0.06)
    for ri, r in enumerate(rows):
        if ri in hi_rows:
            band = slide.shapes.add_shape(1, x - Inches(0.06), yy, sum(cw) + Inches(0.12), rh)
            _fill(band, PANEL)
        cx = x
        for ci, (v, w_) in enumerate(zip(r, cw)):
            col = hi_col if ri in hi_rows and ci > 0 else (INK if ci == 0 else MUTE)
            bold = ri in hi_rows
            _txt(slide, cx, yy, w_, rh, [(str(v), size, col, bold)], anchor=MSO_ANCHOR.MIDDLE)
            cx += w_
        yy += rh


prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK); _bg(s); return s


def divider(title, sub):
    s = slide()
    band = s.shapes.add_shape(1, 0, Inches(2.7), W, Inches(2.0)); _fill(band, PANEL)
    bar = s.shapes.add_shape(1, Inches(0.9), Inches(2.9), Inches(0.18), Inches(1.6)); _fill(bar, ACCENT)
    _txt(s, Inches(1.3), Inches(3.0), Inches(11), Inches(1.0), [(title, 34, INK, True)])
    _txt(s, Inches(1.3), Inches(4.05), Inches(11), Inches(0.6), [(sub, 16, ACCENT, False)])
    return s


# ══ 1 封面 ══════════════════════════════════════════════════════════════════
s = slide()
band = s.shapes.add_shape(1, 0, Inches(2.3), W, Inches(2.7)); _fill(band, PANEL)
_txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.6),
     [("統計分析：本研究的基石", 22, ACCENT, True)])
_txt(s, Inches(0.9), Inches(2.55), Inches(11.6), Inches(1.3),
     [("Layer 1 規則層　與　Layer 2 統計層", 38, INK, True)])
_txt(s, Inches(0.9), Inches(3.55), Inches(11.6), Inches(0.9),
     [("方法・程序・成果", 30, INK, True)])
_txt(s, Inches(0.9), Inches(5.15), Inches(11.6), Inches(0.6),
     [("白箱規則 P1–P6　＋　CUSUM／BOCPD／SSA 變點偵測　—　可解釋、可稽核、可重跑", 15, ACCENT, False)])
_txt(s, Inches(0.9), Inches(6.55), Inches(11.6), Inches(0.5),
     [("計畫案號 TASA-S-1150268　｜　社團法人中華民國國防科技學術研究學會　｜　2026-07-20", 12, MUTE, False)])

# ══ 2 大綱 ══════════════════════════════════════════════════════════════════
s = slide(); _chip(s, "簡報大綱")
cards_rows(s, [
    ("一", "為何統計層是基石", "在三層架構中的角色定位：候選產生器與可稽核基線", ACCENT),
    ("二", "Layer 1 方法與程序", "白箱規則 P1–P6：偵測型／抑制型／調制型三種角色", ACCENT2),
    ("三", "Layer 1 成果", "消融實驗、54 天全量、獨立 MEME hold-out 驗證", GOOD),
    ("四", "Layer 2 方法與程序", "CUSUM／BOCPD／SSA 三種具名變點偵測法之原理", ACCENT2),
    ("五", "Layer 2 成果", "TLE 與 MEME 雙序列 × 四方法：八組 P/R/lead 實測", GOOD),
    ("六", "能力邊界與誠實聲明", "偵測率由訊噪比決定；三層指標不可橫向比較", WARN),
], y0=Inches(1.35), ch=Inches(0.78), gap=Inches(0.88), tag_w=Inches(1.0))
_footer(s, 2)

# ══ 一、定位 ═════════════════════════════════════════════════════════════════
divider("一、為何統計層是基石", "在三層架構中的角色定位")
_footer(prs.slides[-1], 3)

s = slide(); _chip(s, "三層分工：不是競爭，是各司其職")
cards_rows(s, [
    ("Layer 1　規則", "白箱規則 P1–P6 — 可稽核基線",
     "每個旗標都附「哪條規則、哪個物理量」；不需訓練資料、可直接部署於全庫 14,090 顆", ACCENT),
    ("Layer 2　統計", "CUSUM／BOCPD／SSA — 候選產生器",
     "不預設機動長什麼樣，逐點找變化點；設計為高召回，寧可多報由後段篩選", ACCENT2),
    ("Layer 3　ML", "融合評分器 — 判別器",
     "以 L2 的逐點分數為輸入特徵，綜合五通道給出可操作判別（ROC-AUC 0.982）", GOOD),
], y0=Inches(1.5), ch=Inches(1.25), gap=Inches(1.42), tag_w=Inches(2.6))
_txt(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(1.0),
     [[("關鍵：L2 的輸出「就是」L3 的輸入。", 15, ACCENT2, True),
       ("　融合評分器的五通道 ＝ CUSUM／BOCPD／SSA／3σ-MAD／NRLMSIS 阻力殘差，", 14, INK, False)],
      [("每通道取 ±24h 窗內的 max／mean／p90 共 15 維特徵——沒有統計層，就沒有機器學習層。", 14, INK, False)]], sp=1.3)
_footer(s, 4)

# ══ 二、L1 方法 ══════════════════════════════════════════════════════════════
divider("二、Layer 1 方法與程序", "白箱規則 P1–P6")
_footer(prs.slides[-1], 5)

s = slide(); _chip(s, "Layer 1：六條規則，三種角色")
cards_rows(s, [
    ("P2　偵測", "高度自適應 Δa 閾值",
     "拋物線左側曲線：vertex 700 km、floor 0.4 km；400 km 處門檻 2.0 km（低軌雜訊大→門檻高）", ACCENT),
    ("P4　偵測", "多窗口補充掃描",
     "滑動 3 筆內出現超閾值即補旗標——找回長觀測窗中被平均掉的單週機動", ACCENT),
    ("P6　偵測", "星座感知專屬閾值",
     "依傾角族群調整：53° 1.0×／SSO 1.2×／中傾角 1.1×；MEO/GEO/HEO 改用固定 0.05 km", ACCENT),
    ("P1　抑制", "單調衰減抑制",
     "5 筆窗內 ≥85% 為小幅負 Δa、無跳變（|Δa|<2 km）、B*>0 → 判為大氣阻力，非機動", WARN),
    ("P3　抑制", "B* 輔助抑制",
     "B* > max(中位+2σ, 5e-4) 且 Δa<0 且 |Δa|<1.5 km → 該下降可由高阻力解釋，移除誤報", WARN),
    ("P5　調制", "F10.7 太陽通量倍率",
     "拋物線：vertex 70 sfu、floor 1.0；200 sfu 時 1.6×——太陽活躍期阻力強，門檻同步放寬", ACCENT2),
], y0=Inches(1.3), ch=Inches(0.80), gap=Inches(0.90), tag_w=Inches(2.0))
_footer(s, 6)

s = slide(); _chip(s, "Layer 1 程序：偵測取聯集，再扣除抑制")
_txt(s, Inches(0.7), Inches(1.4), Inches(12.0), Inches(0.5),
     [("資料流（每顆衛星、每對相鄰 TLE）", 16, ACCENT2, True)])
cards_rows(s, [
    ("① 元素差分", "由 TLE 取 Δa／Δi／Δe／ΔΩ（J2 修正後殘差）",
     "僅用推論時可得的公開資料；MEME 不參與偵測，只在驗證階段當真值", ACCENT),
    ("② 閾值判定", "偵測型：flag = P2 ∪ P5 ∪ P6 ∪ P4 ∪ other(Δi/Δe/ΔΩ)",
     "任一指標超標即產生候選旗標——設計為寬進", ACCENT2),
    ("③ 物理抑制", "扣除：combined = detect AND NOT (P1 ∪ P3)",
     "把「可由大氣阻力解釋」的下降從嫌疑名單移除——這是誤報控制的主力", WARN),
    ("④ 輸出", "每旗標附理由字串（哪條規則、門檻多少、觸發幾筆）",
     "可稽核、可人工複核；這是「可解釋」的具體實作，非事後說明", GOOD),
], y0=Inches(2.0), ch=Inches(1.05), gap=Inches(1.19), tag_w=Inches(2.3))
_txt(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4),
     [("設計哲學：偵測寬進、物理嚴出——寧可先標記，再用物理原理逐條剔除。", 13, ACCENT2, True)])
_footer(s, 7)

# ══ 三、L1 成果 ══════════════════════════════════════════════════════════════
divider("三、Layer 1 成果", "消融實驗、全量部署、獨立驗證")
_footer(prs.slides[-1], 8)

s = slide(); _chip(s, "成果 ①　P1–P4 消融實驗：每條規則的獨立貢獻")
_txt(s, Inches(0.7), Inches(1.35), Inches(12.0), Inches(0.4),
     [("設定：2026-05-01～05-30（30 天）、14,019 顆 LEO 衛星", 13, MUTE, False)])
table(s, Inches(0.8), Inches(1.9),
      ["組態", "假陽性數", "精確率", "說明"],
      [["基準（無 P1–P4）", "68 顆", "94.8%", "僅固定閾值差分偵測"],
       ["＋P1 單調衰減抑制", "—", "—", "移除純阻力造成的誤報"],
       ["＋P2 高度自適應", "—", "—", "低軌高門檻、高軌低門檻"],
       ["＋P3 B* 輔助抑制", "—", "—", "高阻力個體的專屬校正"],
       ["完整 P1–P4", "29 顆", "97.5%", "另由 P4 找回 26 顆遺漏機動"]],
      [Inches(3.3), Inches(2.0), Inches(1.8), Inches(4.6)], hi_rows=(4,))
_txt(s, Inches(0.8), Inches(4.5), Inches(11.8), Inches(1.6),
     [[("關鍵發現：", 15, ACCENT2, True),
       ("假陽性由 68 顆降至 29 顆（減少 57%），精確率 94.8% → 97.5%。", 14, INK, False)],
      [("P4 多窗口補充額外找回 26 顆基準方法遺漏的真實機動衛星——", 14, INK, False),
       ("抑制與補充同時進行，不是單純調高門檻。", 14, ACCENT2, True)]], sp=1.35)
_txt(s, Inches(0.8), Inches(6.4), Inches(11.8), Inches(0.5),
     [("註：本頁為 30 天／P1–P4 設定，與次頁 54 天／P1–P6 設定不可直接比較。", 12, WARN, True)])
_footer(s, 9)

s = slide(); _chip(s, "成果 ②　54 天全量部署與獨立 hold-out 驗證")
_txt(s, Inches(0.7), Inches(1.32), Inches(12.0), Inches(0.4),
     [("設定：2026-05-01～06-23（54 天）、14,090 顆 LEO 衛星、啟用 P1–P6", 13, MUTE, False)])
# 左：全量指標
card = s.shapes.add_shape(1, Inches(0.6), Inches(1.85), Inches(5.9), Inches(4.3)); _fill(card, PANEL)
_txt(s, Inches(0.95), Inches(2.05), Inches(5.2), Inches(0.4), [("全量部署指標", 16, ACCENT, True)])
_txt(s, Inches(0.95), Inches(2.55), Inches(5.2), Inches(3.4),
     [[("Overall Recall", 13, MUTE, False), ("　26.9%", 20, INK, True)],
      [("誤報率 FAR", 13, MUTE, False), ("　5.4%", 20, GOOD, True)],
      [("Precision（整體）", 13, MUTE, False), ("　94.6%", 20, INK, True)],
      [("Precision@Top-1000", 13, MUTE, False), ("　98.2%", 20, GOOD, True)],
      [("真值規模", 13, MUTE, False), ("　10,186 顆（推進代理 GT）", 15, INK, False)]], sp=1.62)
# 右：hold-out
card = s.shapes.add_shape(1, Inches(6.8), Inches(1.85), Inches(5.9), Inches(4.3)); _fill(card, PANEL)
_txt(s, Inches(7.15), Inches(2.05), Inches(5.2), Inches(0.4),
     [("獨立 hold-out（MEME 精密星曆）", 16, ACCENT2, True)])
_txt(s, Inches(7.15), Inches(2.55), Inches(5.2), Inches(3.4),
     [[("真實事件數", 13, MUTE, False), ("　99 個", 20, INK, True)],
      [("事件級 Recall", 13, MUTE, False), ("　57.6%", 20, GOOD, True)],
      [("平均偵測前置時間", 13, MUTE, False), ("　24.4 小時", 20, GOOD, True)],
      [("跨時間段一致率", 13, MUTE, False), ("　89.7%", 20, INK, True)],
      [("意義", 13, MUTE, False), ("　完全獨立之第二真值源", 15, INK, False)]], sp=1.62)
_txt(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.9),
     [[("為何 Overall Recall 只有 26.9%？", 14, ACCENT2, True),
       ("因為代理真值是「該衛星有無推進能力」，", 13, INK, False)],
      [("而 54 天內有推進能力者未必實際機動——分母被高估。事件級 hold-out 的 57.6% 才是真正的偵測力。", 13, INK, False)]], sp=1.3)
_footer(s, 10)

# ══ 四、L2 方法 ══════════════════════════════════════════════════════════════
divider("四、Layer 2 方法與程序", "CUSUM／BOCPD／SSA 三種具名變點偵測")
_footer(prs.slides[-1], 11)

s = slide(); _chip(s, "Layer 2：三種具名方法 ＋ 一組對照基準")
cards_rows(s, [
    ("CUSUM", "累積和管制圖（Page, 1954）",
     "偵測序列均值的階躍變化。作用於 Δa 增量序列，穩健標準化後累積；k=0.5、h=5.0", ACCENT),
    ("BOCPD", "貝氏線上變點偵測（Adams & MacKay, 2007）",
     "Normal-Gamma 共軛、Student-t 預測；每時刻輸出「變點機率」；hazard λ=100", ACCENT2),
    ("SSA", "奇異譜分析（Singular Spectrum Analysis）",
     "滯後嵌入→SVD→取前 3 主成分重構「趨勢＋J2 振盪」，殘差尖峰即機動；L=min(24, n/3)", GOOD),
    ("3σ-MAD", "穩健 z 分數尖峰偵測（對照基準）",
     "z = |Δa − median| / (1.4826·MAD)，z>3 標為異常——最簡單，但實測召回最高", MUTE),
], y0=Inches(1.5), ch=Inches(1.12), gap=Inches(1.27), tag_w=Inches(2.0))
_txt(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6),
     [("全部以 numpy／scipy 自行實作，無外部相依；每個偵測器同時輸出「逐點分數」（供 L3 當特徵）與「事件索引」（供 L2 輸出）。",
       13, ACCENT2, True)])
_footer(s, 12)

s = slide(); _chip(s, "三法原理對照：各自擅長什麼")
cards_rows(s, [
    ("CUSUM　累積", "「小變化持續累積，總會超線」",
     "把每點偏差減去容差 k 後累加，超過 h 即報警並重置。對緩慢漂移敏感——單次大推力會讓 S 快速衝頂", ACCENT),
    ("BOCPD　機率", "「這段穩定期被打斷了嗎？」",
     "維護 run-length（已穩定多久）的後驗分布。用 P(run<3) 當分數，比標準 P(run=0) 對連續機動更穩健", ACCENT2),
    ("SSA　分解", "「先算出它『本來該長怎樣』，再看差多少」",
     "把序列拆成趨勢（大氣阻力）＋振盪（J2 攝動）＋殘差。機動落在殘差——與 Model 2 的物理殘差思路同源", GOOD),
], y0=Inches(1.5), ch=Inches(1.42), gap=Inches(1.60), tag_w=Inches(2.4))
_txt(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.7),
     [[("三法互補而非重複：", 14, ACCENT2, True),
       ("CUSUM 抓「累積」、BOCPD 抓「段落被打斷」、SSA 抓「偏離物理常態」。", 13.5, INK, False)],
      [("這正是它們能各自成為融合評分器獨立通道的原因——若三者高度相關，融合就沒有增益。", 13.5, INK, False)]], sp=1.3)
_footer(s, 13)

s = slide(); _chip(s, "Layer 2 程序：雙序列驗證設計")
cards_rows(s, [
    ("① 雙序列", "同時在 TLE sma 與 MEME mean-a 上跑同一組偵測器",
     "TLE＝推論時可得、可部署；MEME＝乾淨、已消 J2 振盪。比較兩者即可分離「方法限制」與「資料限制」", ACCENT),
    ("② 真值定義", "MEME transitions_full 中 da_severity ∈ {medium, large}（≥5 km）",
     "與主模型一致的門檻；共 1,758 個真值事件、283–284 顆衛星", ACCENT2),
    ("③ episode 合併", "相鄰機動時刻間隔 > 48 小時 → 視為不同 episode",
     "避免把一次連續機動的 per-8h 多個時刻重複計為多次事件", ACCENT2),
    ("④ 匹配與計分", "偵測事件 ↔ 真值 episode 以 ±24 小時容差匹配",
     "計算 precision／recall／lead-time（負值＝提前偵測）", GOOD),
], y0=Inches(1.5), ch=Inches(1.10), gap=Inches(1.25), tag_w=Inches(2.2))
_txt(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5),
     [("這個設計的價值：同一方法在兩種資料上的落差，直接量化了「公開 TLE 的資訊上限」。", 13.5, ACCENT2, True)])
_footer(s, 14)

# ══ 五、L2 成果 ══════════════════════════════════════════════════════════════
divider("五、Layer 2 成果", "八組實測：四方法 × 兩種輸入序列")
_footer(prs.slides[-1], 15)

s = slide(); _chip(s, "成果 ③　四方法 × 兩序列：完整實測結果")
_txt(s, Inches(0.6), Inches(1.28), Inches(12.0), Inches(0.4),
     [("真值 1,758 個 MEME 機動 episode；容差 ±24h；lead 負值＝提前偵測", 12.5, MUTE, False)])
table(s, Inches(0.7), Inches(1.78),
      ["方法", "輸入", "衛星數", "偵測數", "精確率", "召回率", "前置時間"],
      [["BOCPD", "MEME", "283", "355", "0.862", "0.148", "0.0 h"],
       ["CUSUM", "MEME", "283", "1,032", "0.968", "0.127", "−43.2 h"],
       ["3σ-MAD", "MEME", "283", "1,469", "0.953", "0.210", "−75.7 h"],
       ["SSA", "MEME", "283", "2,709", "0.952", "0.217", "−56.1 h"],
       ["BOCPD", "TLE", "284", "1,094", "0.480", "0.201", "−75.0 h"],
       ["CUSUM", "TLE", "284", "6,687", "0.496", "0.482", "−43.0 h"],
       ["3σ-MAD", "TLE", "284", "10,378", "0.526", "0.668", "−23.1 h"],
       ["SSA", "TLE", "284", "9,839", "0.493", "0.414", "−44.6 h"]],
      [Inches(1.6), Inches(1.2), Inches(1.3), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.7)],
      hi_rows=(7, 1), size=12.5, rh=Inches(0.40))
_txt(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(1.3),
     [[("TLE 上最佳召回：3σ-MAD 0.668", 14, GOOD, True),
       ("（期中報告所引之「統計層最佳 0.67」即為此值）", 13, INK, False)],
      [("MEME 上最佳精確率：CUSUM 0.968", 14, ACCENT2, True),
       ("——同一方法換乾淨資料，精確率由 0.496 躍升至 0.968。", 13, INK, False)]], sp=1.35)
_footer(s, 16)

s = slide(); _chip(s, "成果 ④　最重要的發現：精確率與召回率的翻轉")
cards_rows(s, [
    ("MEME 序列", "高精確率、低召回率　（0.86–0.97 ／ 0.13–0.22）",
     "資料乾淨 → 報出來的幾乎都對；但序列點數少、已被檔案平均平滑，小事件被抹平 → 漏抓多", GOOD),
    ("TLE 序列", "低精確率、高召回率　（0.48–0.53 ／ 0.20–0.67）",
     "雜訊大 → 一半是誤報；但取樣密、保留原始跳變 → 抓得到的事件反而多", ACCENT),
], y0=Inches(1.5), ch=Inches(1.30), gap=Inches(1.48), tag_w=Inches(2.3))
_txt(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.2),
     [[("這個翻轉為什麼重要？", 16, ACCENT2, True)],
      [("① 它證明「資料更好」不等於「偵測更多」——MEME 精度高 100 倍，召回率反而更低。", 14, INK, False)],
      [("② 它解釋了為何統計層必須設計為「候選產生器」：TLE 的高召回是可用的，", 14, INK, False),
       ("低精確率交給後段（物理抑制＋ML）處理。", 14, GOOD, True)],
      [("③ 它是三層架構的實證依據——單一層無法同時滿足高召回與高精確率。", 14, INK, False)]], sp=1.42)
_footer(s, 17)

# ══ 六、邊界 ═════════════════════════════════════════════════════════════════
divider("六、能力邊界與誠實聲明", "偵測率由訊噪比決定；指標不可橫向比較")
_footer(prs.slides[-1], 18)

s = slide(); _chip(s, "成果 ⑤　合成注入實驗：偵測率由「訊噪比」決定")
_txt(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(0.4),
     [("在三種雜訊帶的真實 TLE 序列上注入已知量級的人工機動，量測偵測率", 12.5, MUTE, False)])
table(s, Inches(0.75), Inches(1.8),
      ["雜訊帶", "σ (km)", "Δa=0.1", "Δa=0.2", "Δa=0.3", "Δa=0.5", "Δa=1.0", "Δa=2.0"],
      [["LEO 低軌 <450 km", "0.15", "2.5%", "3.5%", "8.5%", "25%", "92.5%", "100%"],
       ["LEO 中軌 450–700", "0.08", "4.5%", "17.5%", "41.5%", "84.5%", "100%", "100%"],
       ["LEO 高軌 >700 km", "0.05", "5.5%", "51%", "82%", "100%", "100%", "100%"]],
      [Inches(2.7), Inches(1.1), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3)],
      size=12.5, rh=Inches(0.42))
_txt(s, Inches(0.75), Inches(3.75), Inches(11.9), Inches(2.9),
     [[("同一個 Δa，在不同雜訊帶的偵測率天差地別。", 15, ACCENT2, True)],
      [("Δa = 0.2 km：低軌只有 3.5%，高軌卻有 51%——差 15 倍。", 14, INK, False)],
      [("但若改用訊噪比 SNR = Δa / σ 來看，三條曲線大致收攏為同一條：", 14, GOOD, True)],
      [("　　SNR < 2 → 幾乎測不到（三帶全距僅 0.02）　｜　SNR ≥ 7 → 100%（全距 0）", 14, INK, True)],
      [("　　轉折區 SNR 3–4.5 → 25%～51%（全距 0.26）：此區仍有殘餘差異，SNR 非唯一決定因素", 13, WARN, True)],
      [("結論：偵測能力主要由「相對於該星雜訊的量級」決定，而非絕對 Δa 大小。", 14, ACCENT2, True)],
      [("這正是後續 ML 層改用 σ 正規化（SNR）特徵的實證依據——統計層先發現了這個規律。", 13.5, INK, False)]], sp=1.38)
_footer(s, 19)

s = slide(); _chip(s, "統計層如何餵養機器學習層", GOOD)
cards_rows(s, [
    ("輸入", "四個統計通道的逐點分數 ＋ NRLMSIS 阻力殘差",
     "CUSUM／BOCPD／SSA／3σ-MAD 的 scores 陣列，與物理殘差並列為五通道", ACCENT),
    ("特徵工程", "每通道取 ±24 小時窗內的 max／mean／p90",
     "5 通道 × 3 統計量 ＝ 15 維特徵向量——這就是融合評分器的全部輸入", ACCENT2),
    ("輸出", "融合評分器 ROC-AUC 0.9823、AP 0.9583、FPR 0.0498",
     "large 召回 0.9728、medium 召回 0.9448——勝過任何單一統計通道", GOOD),
], y0=Inches(1.5), ch=Inches(1.22), gap=Inches(1.38), tag_w=Inches(2.2))
_txt(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.2),
     [[("但 small 召回 ＝ 0.000（目標 0.65）。", 15, WARN, True),
       ("這不是模型缺陷，而是統計層邊界的直接後果：", 14, INK, False)],
      [("small 級機動的 SNR 落在上頁「SNR < 2 幾乎測不到」的區間——", 14, INK, False),
       ("訊號本來就不在資料裡，任何下游模型都無法還原。", 14, WARN, True)]], sp=1.35)
_footer(s, 20)

s = slide(); _chip(s, "誠實聲明：三層指標不可橫向比較", WARN)
_txt(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5),
     [("三層各自使用不同的測試集與 Ground Truth。以下表格之 eval_basis 欄已明確標註。", 14, INK, False)])
table(s, Inches(0.7), Inches(1.9),
      ["層", "偵測器", "評估基準（Ground Truth）", "精確率", "召回率"],
      [["L1 規則", "P1–P6 整體", "推進代理 GT · 54 天 · 14,090 顆", "94.6", "26.9"],
       ["L1 規則", "P1–P6 @Top-1000", "推進代理 GT · Recall@N 排名", "98.2", "9.6"],
       ["L2 統計", "3σ-MAD (TLE)", "MEME episodes · 1,758 事件 · 284 顆", "52.6", "66.8"],
       ["L2 統計", "CUSUM (MEME)", "MEME episodes · 1,758 事件 · 283 顆", "96.8", "12.7"],
       ["L3 監督式", "LightGBM", "Plan B 自標籤測試集 · 2,104 顆", "99.5", "97.5"],
       ["L3 監督式", "LightGBM（MEME 外部）", "Plan A · MEME 真值 · 252 顆（外部）", "100.0", "39.7"]],
      [Inches(1.5), Inches(2.6), Inches(4.9), Inches(1.4), Inches(1.4)],
      size=12, rh=Inches(0.40), hi_rows=(5,), hi_col=WARN)
_txt(s, Inches(0.7), Inches(4.95), Inches(11.9), Inches(2.0),
     [[("三個必須說明的事實：", 15, ACCENT2, True)],
      [("① L1 的 26.9% 與 L3 的 97.5% 不是同一件事——分母、真值、時間窗全都不同。", 13.5, INK, False)],
      [("② L3 內部 97.5% 與外部 39.7% 的落差，就是監督式模型的域外泛化問題（已由域不變模型改善）。", 13.5, INK, False)],
      [("③ 本簡報所有數字皆可由 data/benchmark/benchmark_v1_20260715.csv 重跑複核。", 13.5, GOOD, True)]], sp=1.4)
_footer(s, 21)

# 同一擂台補充：可比較的版本（§13.2）
s = slide(); _chip(s, "同一擂台補充：把三層放到「完全相同」的條件比較（§13.2）", GOOD)
_txt(s, Inches(0.6), Inches(1.28), Inches(12.1), Inches(0.5),
     [("上頁強調「不可橫向比較」；為此另建一個可比較版本：284 顆 Starlink、同一 MEME episode unit、同一操作點 FPR≤0.05。", 13, INK, False)])
table(s, Inches(0.7), Inches(1.85),
      ["層／方法", "AUC", "精確率", "召回率", "large", "medium", "small"],
      [["L1 規則 P1–P6", "—", "0.321", "0.251", "0.333", "0.083", "0.000"],
       ["L2 CUSUM（最佳單通道）", "0.892", "0.509", "0.444", "0.437", "0.486", "0.000"],
       ["L2 五通道樸素 max", "0.818", "0.328", "0.209", "0.299", "0.022", "0.000"],
       ["L3 融合評分器 OOF", "0.982", "0.689", "0.950", "0.973", "0.950", "0.091"],
       ["naive 隨機（對照）", "—", "0.093", "0.044", "0.043", "0.043", "0.073"]],
      [Inches(3.5), Inches(1.2), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4)],
      size=12.5, rh=Inches(0.42), hi_rows=(3,))
_txt(s, Inches(0.7), Inches(4.55), Inches(11.9), Inches(2.3),
     [[("L2 四通道同一擂台 AUC：CUSUM 0.892 ＞ 3σ-MAD 0.847 ＞ SSA 0.801 ＞ BOCPD 0.713", 14, ACCENT2, True)],
      [("——四法各有盲區、皆無法單獨涵蓋，正是 L3 融合的動機（融合後 AUC 0.982）。", 13.5, INK, False)],
      [("關鍵對照：naive 隨機僅召回 0.044 ", 14, GOOD, True),
       ("——證明此擂台具鑑別力（若擂台無效，隨機應接近滿分）；同標準下 L3 ＞ L2 ＞ L1 ＞ naive。", 13.5, INK, False)],
      [("誠實界定：本擂台僅 Starlink，測不到 L1 域先驗與跨域優勢；L1 門檻為全庫調校非為此集最佳化。", 12.5, WARN, True)]], sp=1.36)
_footer(s, 22)

# ══ 結語 ═════════════════════════════════════════════════════════════════════
s = slide()
band = s.shapes.add_shape(1, 0, Inches(2.4), W, Inches(2.3)); _fill(band, PANEL)
_txt(s, Inches(0.9), Inches(2.65), Inches(11.5), Inches(1.0),
     [("統計層是基石，不是配角", 32, INK, True)])
_txt(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.7),
     [("沒有可稽核的規則基線與高召回的候選產生器，機器學習層沒有輸入、也沒有對照", 16, ACCENT, False)])
_txt(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.8),
     [[("・Layer 1 提供可解釋基線：14,090 顆全量部署，Precision@1000 98.2%、誤報率 5.4%", 14, INK, False)],
      [("・Layer 2 提供 ML 的全部輸入：五通道 × 三統計量 ＝ 15 維特徵，融合後 ROC-AUC 0.982", 14, INK, False)],
      [("・Layer 2 先發現「偵測力由訊噪比決定」，才有後續 ML 的 σ 正規化域不變模型", 14, GOOD, True)],
      [("・能力邊界已量化：SNR<2 測不到——這是資料的物理限制，不是演算法的失敗", 14, ACCENT2, True)]], sp=1.42)
_txt(s, Inches(0.9), Inches(6.95), Inches(11.6), Inches(0.4),
     [("TASA-S-1150268　｜　2026-07-20", 11, MUTE, False)])

out = DOCS / "統計分析基石_L1L2_20260720.pptx"
prs.save(str(out))
print(f"saved {out}  ({out.stat().st_size//1024} KB, {len(prs.slides._sldIdLst)} slides)")
