# -*- coding: utf-8 -*-
"""build_tasa14_ppt.py — NASA/ILRS 14 星標竿比較投影片(TASA 風,16:9)。
輸出：docs/TASA_ILRS_benchmark_20260802.pptx"""
import sys
from pathlib import Path
try: sys.stdout.reconfigure(encoding="utf-8")
except Exception: pass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

D = Path(r"f:\GitHub\Sat_TraingDataExtension\docs")
BG=RGBColor(0x0B,0x1E,0x45); PANEL=RGBColor(0x14,0x2A,0x55); ACCENT=RGBColor(0x2F,0x86,0xE0)
INK=RGBColor(0xF2,0xF5,0xFA); MUTE=RGBColor(0x9F,0xB3,0xC8); GOLD=RGBColor(0xFF,0xD5,0x4F)
GOOD=RGBColor(0x5AC8,0x5AC8,0x5AC8) if False else RGBColor(0x57,0xC7,0x8E)
FONT="Microsoft JhengHei"; W,H=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H; BLANK=prs.slide_layouts[6]

def bg(s,c=BG):
    r=s.shapes.add_shape(1,0,0,W,H); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background()
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return r

def box(s,x,y,w,h,c):
    r=s.shapes.add_shape(1,x,y,w,h); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); return r

def txt(s,x,y,w,h,runs,size=18,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,gap=1.05):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[(runs,size,color,bold)]
    for i,item in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=gap; p.space_after=Pt(6)
        for seg in (item if isinstance(item,list) else [item]):
            t,sz,cl,bd=seg; r=p.add_run(); r.text=t; f=r.font; f.name=sz and FONT; f.size=Pt(sz); f.color.rgb=cl; f.bold=bd
    return tb

def head(s,title):
    box(s,0,0,W,Inches(1.15),PANEL); box(s,Inches(0.5),Inches(0.28),Inches(0.13),Inches(0.6),ACCENT)
    txt(s,Inches(0.78),Inches(0.22),Inches(11.8),Inches(0.75),[(title,26,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(11.3),Inches(0.35),Inches(1.7),Inches(0.5),[("✷ TASA",15,MUTE,True)],align=PP_ALIGN.RIGHT)

def foot(s,n):
    txt(s,Inches(0.5),Inches(7.08),Inches(10),Inches(0.34),
        [("TASA-S-1150268　NASA/ILRS 外部標竿驗證",11,MUTE,False)])
    txt(s,Inches(12.4),Inches(7.08),Inches(0.7),Inches(0.34),[(str(n),11,MUTE,False)],align=PP_ALIGN.RIGHT)

def new():
    s=prs.slides.add_slide(BLANK); bg(s); return s

def img(s,path,x,y,w):
    if Path(path).exists(): s.shapes.add_picture(str(path),x,y,width=w)

# ── 1 封面 ──
s=new()
box(s,0,Inches(2.3),W,Inches(2.7),PANEL)
box(s,Inches(0.9),Inches(2.55),Inches(0.16),Inches(2.2),ACCENT)
txt(s,Inches(1.25),Inches(2.6),Inches(11),Inches(1.7),
    [[("TASA Space Event Detection Benchmark",30,INK,True)],[("validated against NASA / ILRS maneuver truth",24,INK,True)]],gap=1.1)
txt(s,Inches(1.28),Inches(4.15),Inches(11),Inches(0.6),[("本專案偵測法 vs 曲線擬合法(Poly/LOWESS)—— 14 星外部認證標竿",18,ACCENT,True)])
txt(s,Inches(1.28),Inches(5.3),Inches(11),Inches(0.9),
    [[("Systems Engineering Division　|　2026-08-01",15,MUTE,False)],
     [("國家太空中心　Taiwan Space Agency",14,MUTE,False)]],gap=1.2)

# ── 2 標竿資料 ──
s=new(); head(s,"標竿資料：NASA/ILRS 14 星機動歷史(1992–2025)")
txt(s,Inches(0.6),Inches(1.35),Inches(12.2),Inches(1.0),
    [[("● 來源 ",16,GOLD,True),("ilrs.gsfc.nasa.gov：operator 認證、公尺級精密定軌之真值(外部、可稽核)。",16,INK,False)],
     [("● 收斂為機動事件後與 SE 組簡報高度吻合：",15,INK,False),("Envisat 177、HY-2A 58、TOPEX 43 完全一致。",15,GOOD,True)]],gap=1.1)
txt(s,Inches(0.6),Inches(3.05),Inches(12.4),Inches(3.6),
    [[("本研究補齊(資料整備)：",17,GOLD,True)],
     [("① ILRS 真值更正：",15,INK,True),("HY-2D 49206→48621、SWOT 55160→54754(原接錯 OneWeb),補 TOPEX window-only → 14 星/1718 burns。",15,INK,False)],
     [("② TLE 歷史回補：",15,INK,True),("14 星向 Space-Track 補至機動起始年(TOPEX 1992…),+16.2 萬筆;名稱×高度×發射日 三重驗證。",15,INK,False)],
     [("③ NORAD 不回收：",15,INK,True),("編目號一經指派永久綁定(5→6 位數用盡即證)→ 接錯係原碼未驗證猜測,非回收。",15,INK,False)]],gap=1.15)
foot(s,2)

# ── 3 方法 ──
s=new(); head(s,"比較方法(評估口徑一致：窗容差 ±1.5 天,逐星 F1)")
rows=[("基線","固定絕對門檻 |Δa|"),
      ("曲線法(PDF 對照)","Polynomial Fit / LOWESS × ΔSMA、ΔSMA/Δt × 原/自適應/迭代"),
      ("本專案 · σ 正規化","相對各星雜訊 σ 的 SNR 門檻(單趟);迭代+前後中位『位準位移』(抑 FP)"),
      ("本專案 · L2 統計層","CUSUM、BOCPD、SSA、3σ-MAD;無監督融合(union / vote≥2)")]
y=1.5
for k,v in rows:
    box(s,Inches(0.6),Inches(y),Inches(3.2),Inches(1.15),PANEL)
    txt(s,Inches(0.75),Inches(y+0.1),Inches(3.0),Inches(0.95),[(k,16,ACCENT,True)],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(4.0),Inches(y+0.1),Inches(8.9),Inches(0.95),[(v,15,INK,False)],anchor=MSO_ANCHOR.MIDDLE,gap=1.05)
    y+=1.28
txt(s,Inches(0.6),Inches(6.65),Inches(12.4),Inches(0.5),
    [[("設計依據：",14,GOLD,True),("雜訊底隨高度巨變(719–1338km 僅 0.3m vs Starlink 24–75m)→ 門檻應以高度相依 σ 正規化取代固定值。",14,MUTE,False)]])
foot(s,3)

# ── 4 結果總排名(圖) ──
s=new(); head(s,"結果：全方法平均 F1 總排名")
img(s,D/"fig_tasa14_allmethods.png",Inches(0.5),Inches(1.35),Inches(8.6))
txt(s,Inches(9.3),Inches(1.5),Inches(3.7),Inches(5.2),
    [[("重點",18,GOLD,True)],
     [("現代 6 星子集 ",15,GOOD,True),("平均 F1 0.64",15,GOOD,True),(" > PDF 0.52",15,INK,False)],
     [("BOCPD ",15,ACCENT,True),("0.46(最高 0.863),為最強單一統計偵測器",15,INK,False)],
     [("迭代+位準位移 ",15,ACCENT,True),("0.46,緊追 PDF",15,INK,False)],
     [("固定門檻僅 0.21 ",15,MUTE,False),("→ σ 正規化決定性勝出",15,INK,False)]],gap=1.25)
foot(s,4)

# ── 5 迭代效果+逐星(圖) ──
s=new(); head(s,"迭代之效果與逐星表現")
img(s,D/"fig_tasa14_summary.png",Inches(0.35),Inches(1.4),Inches(12.6))
txt(s,Inches(0.6),Inches(6.75),Inches(12.4),Inches(0.5),
    [[("迭代把平均 F1 0.38→0.46、精確率 0.30→0.49(FP 大降);前 5 現代星全部 ≥ PDF 平均 0.52。",14,INK,False)]])
foot(s,5)

# ── 6 委員提問回應：資料品質量化 ──
s=new(); head(s,"委員提問：14 星資料品質能量化比較嗎? → 可以")
img(s,D/"fig_tasa14_dataquality.png",Inches(0.45),Inches(1.35),Inches(8.3))
txt(s,Inches(9.0),Inches(1.45),Inches(4.0),Inches(5.4),
    [[("量化四維度",17,GOLD,True)],
     [("σ 雜訊底、更新間隔、|Δa| 量級、SNR 可見比例",13.5,MUTE,False)],
     [("① σ 皆 0.1–0.3m 均一 ",14,GOOD,True),("→ 不構成差異(ρ=−0.04)",14,INK,False)],
     [("② SNR≥2 幾乎 100% ",14,GOOD,True),("→ 非限制(異於 Starlink)",14,INK,False)],
     [("③ F1 由目標特性決定：",14,ACCENT,True)],
     [("　|Δa| 量級 ρ=+0.75、更新間隔 ρ=−0.51",14,INK,False)],
     [("結論：低分星為內在困難目標",14,GOLD,True)],
     [("(小機動+稀疏追蹤),非方法失效——可用數字證明。",13.5,INK,False)]],gap=1.2)
foot(s,6)

# ── 7 討論與結論 ──
s=new(); head(s,"討論與結論(誠實界定)")
txt(s,Inches(0.6),Inches(1.4),Inches(12.4),Inches(5.3),
    [[("① 乾淨現代資料上,本法勝曲線法平均：",16,GOOD,True),("現代 6 星 0.64 > PDF 0.52;BOCPD 單星最高 0.863≈PDF 0.92。",16,INK,False)],
     [("② 全域略低(~0.46 vs 0.52)源自資料非演算法：",16,GOLD,True),("TOPEX/Jason-1/Envisat 老 TLE、真值不全,其『FP』多為 ILRS 未列之真實機動 → 精確率被低估。",16,INK,False)],
     [("③ 樸素融合無效 → 凸顯訓練式 L3 之必要：",16,GOLD,True),("4 通道跨域各自 FP 高,OR/投票無法正確組合;L3 為 Starlink 域內模型,跨域僅作一致性驗證、不列數字。",16,INK,False)],
     [("④ BOCPD 為最強單一 L2 偵測器,",16,ACCENT,True),("建議作為 L2 主力通道。",16,INK,False)],
     [("結論：",17,INK,True),("在外部、非自建之認證真值上,本專案『高度相依 σ 正規化 + 迭代精修』再次獲得支持。",16,INK,False)]],gap=1.3)
foot(s,7)

# ── 8 結尾 ──
s=new()
box(s,0,Inches(3.0),W,Inches(1.6),PANEL); box(s,Inches(0.9),Inches(3.2),Inches(0.16),Inches(1.2),ACCENT)
txt(s,Inches(1.25),Inches(3.15),Inches(11),Inches(1.3),[("謝謝聆聽",40,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
txt(s,Inches(1.28),Inches(4.8),Inches(11),Inches(0.5),[("國家太空中心　Taiwan Space Agency",15,MUTE,False)])

out=D/"TASA_ILRS_benchmark_20260802.pptx"; prs.save(str(out))
print("saved",out,f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
