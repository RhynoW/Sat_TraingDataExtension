"""
export_behavior_pptx.py
========================
將 analyze_maneuver_behavior.py 架構說明輸出為 PowerPoint 簡報。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from lxml import etree
import copy

OUT_PATH = r"output/analyze_maneuver_behavior_py_20260525架構說明.pptx"

# ── 主題色彩 ─────────────────────────────────────────────────────────────────
C_NAVY    = RGBColor(0x1F, 0x49, 0x7D)   # 深藍（標題）
C_BLUE    = RGBColor(0x2E, 0x74, 0xB5)   # 中藍（副標題、重點）
C_LBLUE   = RGBColor(0xBD, 0xD7, 0xEE)   # 淺藍（背景帶）
C_ORANGE  = RGBColor(0xED, 0x7D, 0x31)   # 橙色（強調）
C_GREEN   = RGBColor(0x70, 0xAD, 0x47)   # 綠色（✅）
C_RED     = RGBColor(0xFF, 0x00, 0x00)   # 紅色（警告）
C_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)   # 淺灰（表格交替）
C_DGRAY   = RGBColor(0x59, 0x59, 0x59)   # 深灰（內文）
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK   = RGBColor(0x00, 0x00, 0x00)

# 投影片尺寸：16:9 寬螢幕
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── 輔助函式 ─────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # 空白版面
    return prs.slides.add_slide(layout)


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(14), bold=False, color=C_BLACK,
                align=PP_ALIGN.LEFT, wrap=True, font_name="微軟正黑體"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_label(slide, text, left, top, width, height,
              bg=C_BLUE, fg=C_WHITE, font_size=Pt(13), bold=True):
    """圓角矩形標籤"""
    shape = slide.shapes.add_shape(5, left, top, width, height)  # 5=圓角矩形
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = fg
    run.font.name = "微軟正黑體"
    return shape


def slide_header(slide, title: str, subtitle: str = "", slide_num: int = 0):
    """統一的投影片頂部標題帶"""
    # 深藍色頂部色帶
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), fill_color=C_NAVY)
    # 標題文字
    add_textbox(slide, title,
                Inches(0.4), Inches(0.1), Inches(11.5), Inches(0.7),
                font_size=Pt(24), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.72), Inches(10), Inches(0.35),
                    font_size=Pt(13), color=C_LBLUE, align=PP_ALIGN.LEFT)
    if slide_num:
        add_textbox(slide, str(slide_num),
                    Inches(12.7), Inches(0.15), Inches(0.5), Inches(0.4),
                    font_size=Pt(14), color=RGBColor(0xBD,0xD7,0xEE),
                    align=PP_ALIGN.RIGHT)


def add_table(slide, headers, rows,
              left, top, width, height,
              hdr_color=C_NAVY, font_size=Pt(11)):
    """手工繪製表格（支援自訂色彩）"""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    col_w = width // n_cols
    for i in range(n_cols):
        tbl.columns[i].width = col_w

    # 標題列
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True
        run.font.color.rgb = C_WHITE
        run.font.size = font_size
        run.font.name = "微軟正黑體"

    # 資料列
    for ri, row in enumerate(rows):
        bg = C_GRAY if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                run = p.runs[0]
                run.font.size = font_size
                run.font.name = "微軟正黑體"
                run.font.color.rgb = C_DGRAY
    return tbl


def bullet_box(slide, items, left, top, width, height,
               font_size=Pt(13), color=C_DGRAY, bullet="▸ "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.name = "微軟正黑體"
    return txBox


def code_box(slide, text, left, top, width, height, font_size=Pt(10)):
    bg = add_rect(slide, left, top, width, height,
                  fill_color=RGBColor(0x1E,0x1E,0x1E))
    bg.line.color.rgb = C_BLUE
    bg.line.width = Pt(1)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.05),
        width - Inches(0.2), height - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
    return txBox


# ── 各張投影片 ────────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = blank_slide(prs)
    # 全版深藍背景
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_NAVY)
    # 橙色裝飾條
    add_rect(slide, 0, Inches(3.5), Inches(0.12), Inches(2.5), fill_color=C_ORANGE)
    add_rect(slide, Inches(0.22), Inches(3.5), Inches(0.04), Inches(2.5),
             fill_color=RGBColor(0xED,0x7D,0x31))
    # 主標題
    add_textbox(slide, "analyze_maneuver_behavior.py",
                Inches(0.5), Inches(1.8), Inches(12), Inches(1.0),
                font_size=Pt(36), bold=True, color=C_WHITE,
                align=PP_ALIGN.LEFT)
    # 副標題
    add_textbox(slide, "LEO 衛星機動行為分類器　架構說明",
                Inches(0.5), Inches(2.8), Inches(10), Inches(0.7),
                font_size=Pt(22), color=C_LBLUE, align=PP_ALIGN.LEFT)
    # 分隔線
    add_rect(slide, Inches(0.5), Inches(3.65), Inches(8), Inches(0.03),
             fill_color=C_ORANGE)
    # 元資訊
    info = [
        "資料期間：2026-05-01 ～ 2026-05-22（21 天）",
        "分析樣本：1,323 顆 LEO 衛星（maneuver_detected = True）",
        "版本日期：2026-05-25",
    ]
    for i, t in enumerate(info):
        add_textbox(slide, t,
                    Inches(0.5), Inches(3.9 + i*0.45), Inches(10), Inches(0.4),
                    font_size=Pt(14), color=RGBColor(0xBD,0xD7,0xEE))
    # 右下角標誌文字
    add_textbox(slide, "LEO Maneuver Annotation Pipeline",
                Inches(8.5), Inches(6.8), Inches(4.5), Inches(0.5),
                font_size=Pt(11), color=RGBColor(0x80,0x90,0xA8),
                align=PP_ALIGN.RIGHT)


def slide_pipeline(prs):
    slide = blank_slide(prs)
    slide_header(slide, "1. 整體架構：五階段資料管線", "Data Pipeline Overview", 1)

    # 管線方塊
    stages = [
        ("① 讀取\nCSV", C_BLUE),
        ("② 讀取\n標註", C_BLUE),
        ("③ DuckDB\nTLE 查詢", C_ORANGE),
        ("④ TLE 差分\n計算", RGBColor(0x70,0xAD,0x47)),
        ("⑤ 行為\n分類", RGBColor(0xBF,0x38,0x8C)),
        ("⑥ 輸出\nCSV/MD", RGBColor(0x59,0x59,0x59)),
    ]
    box_w, box_h = Inches(1.7), Inches(1.1)
    top = Inches(1.5)
    gap = Inches(0.25)
    start_x = Inches(0.4)
    for i, (label, color) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        add_label(slide, label, x, top, box_w, box_h, bg=color, font_size=Pt(13))
        # 箭頭
        if i < len(stages) - 1:
            add_textbox(slide, "→",
                        x + box_w, top + Inches(0.3),
                        gap, Inches(0.5),
                        font_size=Pt(18), color=C_NAVY, align=PP_ALIGN.CENTER)

    # 各階段說明
    details = [
        ("① 讀 validation_full.csv\n→ 篩 maneuver_detected=True\n→ 1,323 顆 NORAD", Inches(0.4)),
        ("② 讀 annotations_leo_full.csv\n→ norad_id 索引\n→ 取 propulsion_class", Inches(2.35)),
        ("③ 一次性批次查詢\n→ 62,005 筆 TLE\n→ 4 個軌道根數欄位", Inches(4.30)),
        ("④ detect_transitions()\n→ 逐對相鄰 TLE\n→ Δa, Δi, Δe, ΔRAAN", Inches(6.25)),
        ("⑤ classify_behavior()\n→ 10 種行為規則\n→ 決策樹分類", Inches(8.20)),
        ("⑥ CSV 15 欄位\nMarkdown 報告\n6 章節", Inches(10.15)),
    ]
    for text, x in details:
        txBox = slide.shapes.add_textbox(x, Inches(2.85), Inches(1.85), Inches(1.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for line in text.split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = line
            run.font.size = Pt(10)
            run.font.name = "微軟正黑體"
            run.font.color.rgb = C_DGRAY

    # 輸入/輸出標示
    add_label(slide, "輸入", Inches(0.4), Inches(4.5), Inches(1.0), Inches(0.4),
              bg=C_LBLUE, fg=C_NAVY, font_size=Pt(11), bold=False)
    add_textbox(slide,
        "validation_full.csv  ·  annotations_leo_full.csv  ·  space_db.duckdb（DuckDB）",
        Inches(1.5), Inches(4.5), Inches(9), Inches(0.4),
        font_size=Pt(12), color=C_DGRAY)

    add_label(slide, "輸出", Inches(0.4), Inches(5.1), Inches(1.0), Inches(0.4),
              bg=C_ORANGE, fg=C_WHITE, font_size=Pt(11), bold=False)
    add_textbox(slide,
        "maneuver_behavior_classified.csv（1,323 列 × 15 欄）  ·  maneuver_behavior_report.md",
        Inches(1.5), Inches(5.1), Inches(10), Inches(0.4),
        font_size=Pt(12), color=C_DGRAY)

    # 關鍵數字
    kpis = [
        ("1,323", "機動衛星"),
        ("62,005", "TLE 筆數"),
        ("10", "行為類別"),
        ("21 天", "觀測窗口"),
    ]
    for i, (val, lbl) in enumerate(kpis):
        x = Inches(0.4 + i * 3.1)
        add_rect(slide, x, Inches(5.8), Inches(2.8), Inches(1.35),
                 fill_color=C_LBLUE)
        add_textbox(slide, val, x, Inches(5.85), Inches(2.8), Inches(0.75),
                    font_size=Pt(28), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, x, Inches(6.55), Inches(2.8), Inches(0.4),
                    font_size=Pt(12), color=C_BLUE, align=PP_ALIGN.CENTER)


def slide_physics(prs):
    slide = blank_slide(prs)
    slide_header(slide, "2. 物理量計算：detect_transitions()", "Per-TLE-pair Orbital Difference", 2)

    # 左：四個差分說明
    add_label(slide, "四個軌道差分", Inches(0.4), Inches(1.3), Inches(3.5), Inches(0.4),
              bg=C_BLUE, font_size=Pt(12))
    add_table(slide,
        ["變數", "公式", "物理意義"],
        [
            ["da",        "sma[i] − sma[i-1]",          "半長軸差（+ 升軌 / − 降軌）"],
            ["di",        "inc[i] − inc[i-1]",           "傾角差"],
            ["de",        "ecc[i] − ecc[i-1]",           "離心率差"],
            ["draan_res", "ΔRAAN_raw − J2進動量×Δt",      "J2修正後 RAAN 殘差"],
        ],
        left=Inches(0.4), top=Inches(1.8), width=Inches(6.5), height=Inches(1.9),
        font_size=Pt(11))

    # 左：閾值表
    add_label(slide, "旗標閾值（任一超過 → flagged=True）",
              Inches(0.4), Inches(3.9), Inches(4.5), Inches(0.4),
              bg=C_ORANGE, font_size=Pt(12))
    add_table(slide,
        ["參數", "閾值"],
        [
            ["THR_DA_SM", "> 1.0 km"],
            ["THR_DI",    "> 0.02°"],
            ["THR_DE",    "> 0.001"],
            ["THR_DRAAN", "> 0.1°"],
        ],
        left=Inches(0.4), top=Inches(4.4), width=Inches(3.5), height=Inches(1.6),
        font_size=Pt(12))

    # 右：J2 公式
    add_label(slide, "J2 RAAN 進動修正公式",
              Inches(7.2), Inches(1.3), Inches(5.7), Inches(0.4),
              bg=C_NAVY, font_size=Pt(12))
    code_box(slide,
"# j2_raan_rate(sma, ecc, inc_deg)\n"
"n  = sqrt(MU / sma³)      # 平均角速度 rad/s\n"
"p  = sma × (1 − ecc²)     # 半通徑 km\n\n"
"dΩ/dt = −(3/2) × J2\n"
"       × (RE/p)²\n"
"       × n × cos(inc)\n"
"               # 單位：deg/s\n\n"
"draan_res = RAAN_raw_diff\n"
"          − dΩ/dt × Δt_s",
        Inches(7.2), Inches(1.85), Inches(5.7), Inches(3.3))

    # 右下：間隔說明
    add_rect(slide, Inches(7.2), Inches(5.35), Inches(5.7), Inches(1.5),
             fill_color=C_LBLUE)
    add_textbox(slide, "⚠  跳過條件",
                Inches(7.3), Inches(5.4), Inches(5.5), Inches(0.4),
                font_size=Pt(13), bold=True, color=C_NAVY)
    add_textbox(slide,
        "dt ≤ 0  或  dt > 7 天（604,800 秒）\n→ 避免長期缺測造成差分誤差放大",
        Inches(7.3), Inches(5.78), Inches(5.5), Inches(0.9),
        font_size=Pt(12), color=C_DGRAY)


def slide_features(prs):
    slide = blank_slide(prs)
    slide_header(slide, "3. 衍生特徵計算：classify_behavior() 前半段",
                 "Derived Feature Engineering", 3)

    blocks = [
        ("觸發維度（布林）", C_BLUE,
         "has_da    = flagged 轉移中 |da| > 1.0 km 的存在\n"
         "has_di    = flagged 轉移中 |di| > 0.02° 的存在\n"
         "has_de    = flagged 轉移中 |de| > 0.001 的存在\n"
         "has_draan = flagged 轉移中 |draan_res| > 0.1°",
         Inches(0.3), Inches(1.3), Inches(6.0), Inches(2.3)),

        ("幅度指標", C_ORANGE,
         "max_da  = max(|da|) 於所有 flagged 轉移  (km)\n"
         "net_da  = sum(da) 跨所有轉移             (km)\n"
         "          → 21 天軌道累積漂移量",
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(1.8)),

        ("頻率指標", RGBColor(0x70,0xAD,0x47),
         "flag_rate       = n_flagged / n_transitions\n"
         "burn_freq_per_day = n_flagged / 21",
         Inches(6.7), Inches(3.4), Inches(6.3), Inches(1.5)),
    ]

    for label, color, code, x, y, w, h in blocks:
        add_label(slide, label, x, y, w * 0.55, Inches(0.38),
                  bg=color, font_size=Pt(12))
        code_box(slide, code, x, y + Inches(0.45), w, h - Inches(0.5))

    # 單調性指標（重點，獨立大框）
    add_label(slide, "單調性指標（核心）— 識別大氣阻力衰減 vs 主動機動",
              Inches(0.3), Inches(3.7), Inches(6.0), Inches(0.38),
              bg=RGBColor(0xBF,0x38,0x8C), font_size=Pt(12))
    code_box(slide,
"neg_streak = 連續 da < −0.3 km 的最長序列長度\n"
"total_drop = abs(net_da)  if net_da < 0  else 0\n\n"
"monotone_decay =  (neg_streak  ≥ 5)\n"
"             AND  (total_drop  > 5 km)\n"
"             AND  (net_da      < −3 km)",
        Inches(0.3), Inches(4.15), Inches(6.0), Inches(2.2))

    # 右下說明框
    add_rect(slide, Inches(6.7), Inches(5.1), Inches(6.3), Inches(2.0),
             fill_color=C_LBLUE)
    add_textbox(slide, "物理意義",
                Inches(6.8), Inches(5.15), Inches(6.0), Inches(0.4),
                font_size=Pt(13), bold=True, color=C_NAVY)
    add_textbox(slide,
        "大氣阻力衰減：da 持續為負，無回升（單調性強）\n\n"
        "主動機動：     da 有正有負、呈現非單調特性\n\n"
        "→ 三個條件同時成立才判定為 monotone_decay，\n"
        "   避免短暫阻力補償燒被誤判。",
        Inches(6.8), Inches(5.55), Inches(6.1), Inches(1.4),
        font_size=Pt(12), color=C_DGRAY)


def slide_tree(prs):
    slide = blank_slide(prs)
    slide_header(slide, "4. 規則決策樹：classify_behavior()", "Rule-based Decision Tree", 4)

    add_textbox(slide,
        "最頂層分歧：推進系統標註是否存在　（monotone_decay 在有推進分支中最優先判斷）",
        Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.4),
        font_size=Pt(13), color=C_DGRAY)

    code_box(slide,
"propulsion_class ∈ {Electric_EP, Chemical, Micro/ColdGas, Hybrid/Other}?\n"
"│\n"
"├─ 否（無推進標註）\n"
"│   ├─ monotone_decay = True  →  DragDecay          大氣阻力衰減（FP）\n"
"│   ├─ has_draan only         →  UnknownRAANAnomaly  RAAN 異常待查\n"
"│   └─ 其他                   →  UnknownFP           未知 FP\n"
"│\n"
"└─ 是（有推進標註）\n"
"    ├─ monotone_decay = True       →  DragMakeup       大氣阻力補償  ★最優先\n"
"    ├─ net_da > +5 km              →  OrbitRaising      主動升軌\n"
"    ├─ net_da < −5 km              →  OrbitLowering     主動降軌/離軌\n"
"    ├─ has_da AND has_di           →  ComplexManeuver   複合機動\n"
"    ├─ has_di only                 →  InclinationChange 傾角調整\n"
"    ├─ has_draan only              →  RAANMaintenance   RAAN 維持\n"
"    ├─ flag_rate > 25%             →  Stationkeeping    高頻小幅站位保持\n"
"    ├─ n_flagged ≤ 3 & max_da > 5  →  PhasingManeuver   少次大幅相位調整\n"
"    ├─ flag_rate ≤ 15%             →  DragMakeup        低頻小幅阻力補償\n"
"    └─ 其他                        →  Stationkeeping",
        Inches(0.4), Inches(1.7), Inches(12.5), Inches(5.1))


def slide_categories(prs):
    slide = blank_slide(prs)
    slide_header(slide, "5. 十種行為類別與分類結果",
                 "10 Behavior Categories — Results Summary", 5)

    add_table(slide,
        ["類別代號", "中文名稱", "顆數", "佔比", "關鍵判斷條件", "代表案例"],
        [
            ["DragMakeup",        "大氣阻力補償",     "758", "57.3%",
             "monotone_decay=T 或 flag_rate≤15%",  "Starlink（主流）"],
            ["OrbitLowering",     "軌道降低/離軌",    "252", "19.0%",
             "net_da < −5 km",                     "Starlink 離軌程序"],
            ["PhasingManeuver",   "相位調整",         "91",  " 6.9%",
             "n_flagged≤3 且 max_da>5 km",          "Starlink 新部署"],
            ["OrbitRaising",      "軌道抬升",         "81",  " 6.1%",
             "net_da > +5 km",                     "KUIPER-00262 +64 km"],
            ["Stationkeeping",    "站位保持",         "69",  " 5.2%",
             "flag_rate > 25%",                    "SAR/EO 精密星座"],
            ["UnknownFP",         "未知 FP",          "39",  " 2.9%",
             "無推進標註，非單調",                   "OBJECT 系列"],
            ["DragDecay",         "阻力衰減（無推進）","21",  " 1.6%",
             "monotone_decay=T 且無推進",            "RHOK-SAT 268 km"],
            ["InclinationChange",  "傾角調整",         " 6",  " 0.5%",
             "has_di=T，has_da=F",                  "Starlink 面調整"],
            ["ComplexManeuver",   "複合機動",         " 4",  " 0.3%",
             "has_da=T 且 has_di=T",                "KUIPER-00262"],
            ["UnknownRAANAnomaly","RAAN 異常",        " 2",  " 0.2%",
             "has_draan=T，無 da/di",               "IXPE / ORS 5"],
        ],
        left=Inches(0.3), top=Inches(1.25), width=Inches(12.7), height=Inches(5.9),
        font_size=Pt(10.5))


def slide_propulsion_cross(prs):
    slide = blank_slide(prs)
    slide_header(slide, "6. 推進類別 × 行為分佈交叉表",
                 "Propulsion Class × Behavior Cross-tabulation", 6)

    add_table(slide,
        ["推進類別", "DragMakeup", "OrbitLowering", "PhasingManeuver",
         "OrbitRaising", "Stationkeeping", "UnknownFP", "DragDecay",
         "InclinationChange", "ComplexManeuver", "UnknownRAANAnomaly"],
        [
            ["Electric_EP",  "742", "249", "90", "79", "69", "0", "0", "6", "3", "0"],
            ["Chemical",     " 13", "  1", " 0", " 2", " 0", "0", "0", "0", "1", "0"],
            ["Micro/ColdGas","  2", "  2", " 0", " 0", " 0", "0", "0", "0", "0", "0"],
            ["Hybrid/Other", "  1", "  0", " 1", " 0", " 0", "0", "0", "0", "0", "0"],
            ["未標註",        "  0", "  0", " 0", " 0", " 0","39","21", "0", "0", "2"],
        ],
        left=Inches(0.3), top=Inches(1.3), width=Inches(12.7), height=Inches(3.2),
        font_size=Pt(11))

    # 重點觀察
    add_label(slide, "關鍵觀察", Inches(0.3), Inches(4.7), Inches(2.0), Inches(0.38),
              bg=C_NAVY, font_size=Pt(12))
    obs = [
        "Electric_EP 幾乎包辦全部有推進機動（1,238/1,261 TP），反映 LEO 星座電推主流地位",
        "未標註衛星全部落入 UnknownFP / DragDecay / UnknownRAANAnomaly，無跨界情形",
        "Chemical 的 DragMakeup=13 顆：CSS 空間站三模組 + SZ-22 + YAOGAN 系列的軌道維持燃燒",
        "Micro/ColdGas 的 OrbitLowering=2：ASTROCAST-0103 推進劑耗盡後的終端降軌",
    ]
    bullet_box(slide, obs, Inches(0.3), Inches(5.15), Inches(12.7), Inches(2.0),
               font_size=Pt(12))


def slide_limits(prs):
    slide = blank_slide(prs)
    slide_header(slide, "7. 已知限制與潛在問題", "Known Limitations", 7)

    limits = [
        ("DragMakeup 定義過寬", C_RED,
         "monotone_decay=True + 有推進 → DragMakeup。但 JILIN-01 GAOFEN 3J（net_da=−117 km）"
         "是真正的終端降軌，被誤分為 DragMakeup。\n"
         "→ 改進：加入 total_drop > 50 km 門檻，區分「終端降軌」vs「阻力補償」"),
        ("OrbitLowering 與 Starlink 離軌混淆", C_ORANGE,
         "Starlink 離軌程序（1–3 次大幅燒）與短期相位調整後未補回的衛星，"
         "特徵相似（net_da < −5 km），目前無法從 TLE 差分層次區分。"),
        ("PhasingManeuver 的 n_flagged ≤ 3 是硬編碼", C_ORANGE,
         "Starlink 新部署（1–2 次燒）與 CSS 單次大幅軌維燒，條件相同但語意不同，"
         "需要加入軌道高度或衛星質量資訊才能區分。"),
        ("monotone_decay 未考慮 B* 值", C_BLUE,
         "演算法改進記憶（project_algorithm_improvements.md）建議同時參考 B*。"
         "目前僅用 da 序列，對高阻力小衛星（B* 大）判斷準確，"
         "對中等高度標準衛星可能有漏判。"),
    ]

    row_h = Inches(1.4)
    for i, (title, color, desc) in enumerate(limits):
        y = Inches(1.3) + i * row_h
        add_rect(slide, Inches(0.3), y, Inches(12.7), row_h - Inches(0.05),
                 fill_color=C_GRAY)
        # 色條
        add_rect(slide, Inches(0.3), y, Inches(0.12), row_h - Inches(0.05),
                 fill_color=color)
        add_textbox(slide, title,
                    Inches(0.55), y + Inches(0.05), Inches(12.2), Inches(0.38),
                    font_size=Pt(13), bold=True, color=color)
        add_textbox(slide, desc,
                    Inches(0.55), y + Inches(0.42), Inches(12.2), Inches(0.9),
                    font_size=Pt(11.5), color=C_DGRAY)

    # Recall 原因
    add_label(slide, "Recall 偏低（12.4%）原因", Inches(0.3), Inches(7.0),
              Inches(4.0), Inches(0.38), bg=C_NAVY, font_size=Pt(12))
    reasons = "觀測窗口僅 21 天  ·  閾值 1 km 過濾微小機動  ·  TLE epoch 間機動不可見  ·  升降互抵累積 da ≈ 0"
    add_textbox(slide, reasons,
                Inches(4.5), Inches(7.05), Inches(8.5), Inches(0.38),
                font_size=Pt(11.5), color=C_DGRAY)


def slide_output(prs):
    slide = blank_slide(prs)
    slide_header(slide, "8. 輸出格式", "Output Specification", 8)

    # 左：CSV
    add_label(slide, "maneuver_behavior_classified.csv",
              Inches(0.3), Inches(1.25), Inches(5.8), Inches(0.4),
              bg=C_BLUE, font_size=Pt(12))
    add_table(slide,
        ["欄位", "說明"],
        [
            ["norad_id / sat_name",   "衛星識別"],
            ["propulsion_class",      "推進類別（標註來源）"],
            ["behavior",              "行為類別（10 種）"],
            ["net_da_km",             "21 天累積 Δa（km）"],
            ["max_da_km",             "單次最大 |Δa|（km）"],
            ["n_flagged",             "旗標轉移次數"],
            ["n_transitions",         "總轉移次數"],
            ["flag_rate",             "旗標率 = n_flagged / n_trans"],
            ["burn_freq_per_day",     "日均機動次數"],
            ["monotone_decay",        "是否判定大氣阻力單調衰減"],
            ["triggers",              "觸發維度（da/di/de/draan）"],
            ["mean_sma_km",           "平均 SMA（km）"],
            ["sub_tags",              "細分描述標籤"],
        ],
        left=Inches(0.3), top=Inches(1.75), width=Inches(5.8), height=Inches(5.4),
        font_size=Pt(10.5))

    # 右：Markdown
    add_label(slide, "maneuver_behavior_report.md",
              Inches(6.5), Inches(1.25), Inches(6.5), Inches(0.4),
              bg=C_ORANGE, font_size=Pt(12))
    add_table(slide,
        ["§", "章節內容"],
        [
            ["1", "行為類別總覽（10 種顆數與佔比）"],
            ["2", "推進類別 × 行為分佈交叉表"],
            ["3", "各類別詳細說明 + 代表衛星（前 15 顆）"],
            ["4", "特殊案例（da 最大 Top5、旗標最多 Top5、高頻 Top10）"],
            ["5", "殘餘 FP 完整列表（62 顆）"],
            ["6", "訓練資料建議 + Recall 偏低原因分析"],
        ],
        left=Inches(6.5), top=Inches(1.75), width=Inches(6.5), height=Inches(2.7),
        font_size=Pt(11))

    # 訓練建議表
    add_label(slide, "訓練正例使用建議",
              Inches(6.5), Inches(4.6), Inches(6.5), Inches(0.4),
              bg=C_NAVY, font_size=Pt(12))
    add_table(slide,
        ["行為類別", "推薦", "說明"],
        [
            ["OrbitRaising / Lowering / Phasing", "✅ 高信心", "特徵明確，適合直接納入"],
            ["ComplexManeuver / Stationkeeping",   "✅ 中信心", "可用，需注意標籤精度"],
            ["DragMakeup",                          "⚠ 低信心", "難以與自然衰減區分"],
            ["DragDecay / UnknownFP",               "❌ 排除",  "FP，不納入訓練正例"],
        ],
        left=Inches(6.5), top=Inches(5.1), width=Inches(6.5), height=Inches(2.1),
        font_size=Pt(11))


def slide_callchain(prs):
    slide = blank_slide(prs)
    slide_header(slide, "9. 完整呼叫鏈", "Full Call Chain", 9)

    code_box(slide,
"main()\n"
" ├─ pd.read_csv(validation_full.csv)\n"
" │   └─ 過濾 maneuver_detected=True  →  1,323 顆 NORAD\n"
" ├─ pd.read_csv(annotations_leo_full.csv)\n"
" │   └─ set_index('norad_id')  →  propulsion_class 快速查找\n"
" ├─ duckdb.connect(space_db.duckdb, read_only=True)\n"
" │   └─ SELECT sma_km, eccentricity, inclination_deg, raan_deg\n"
" │      FROM tle_table\n"
" │      WHERE norad_id IN (1323 個 NORAD)\n"
" │        AND date_tag BETWEEN '2026-05-01' AND '2026-05-22'\n"
" │      ORDER BY norad_id, date_tag\n"
" │      → 62,005 筆 TLE\n"
" └─ for norad in norads:                   # 逐顆處理\n"
"     ├─ detect_transitions(sat_tle)\n"
"     │   └─ for i in 1..len(tle):          # 逐對相鄰 TLE\n"
"     │       ├─ da  = sma[i] − sma[i-1]\n"
"     │       ├─ di  = inc[i] − inc[i-1]\n"
"     │       ├─ de  = ecc[i] − ecc[i-1]\n"
"     │       └─ draan_res = angle_diff() − j2_raan_rate() × dt\n"
"     ├─ classify_behavior(transitions, {propulsion_class})\n"
"     │   ├─ 計算 has_da/di/de/draan, net_da, max_da\n"
"     │   ├─ 計算 neg_streak  →  monotone_decay\n"
"     │   ├─ 計算 flag_rate, burn_freq_per_day\n"
"     │   └─ 決策樹（10 條規則）  →  behavior 字串\n"
"     └─ records.append(行為分類結果 dict)\n"
" ├─ DataFrame(records).to_csv('maneuver_behavior_classified.csv')\n"
" └─ generate_report(classified)\n"
"     └─ Path(OUT_MD).write_text(report, encoding='utf-8')",
        Inches(0.3), Inches(1.25), Inches(12.7), Inches(6.1))


# ── 主程式 ───────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_cover(prs)
    slide_pipeline(prs)
    slide_physics(prs)
    slide_features(prs)
    slide_tree(prs)
    slide_categories(prs)
    slide_propulsion_cross(prs)
    slide_limits(prs)
    slide_output(prs)
    slide_callchain(prs)
    prs.save(OUT_PATH)
    print(f"[DONE] {OUT_PATH}  ({len(prs.slides)} 張投影片)")


if __name__ == "__main__":
    main()
